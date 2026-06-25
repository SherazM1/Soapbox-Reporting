from pptx import Presentation
from pathlib import Path

path = Path('templates/Audit_Template_New.pptx')
prs = Presentation(path)
for idx, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        try:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                txt = shape.text.strip().replace('\n', ' | ')
                if txt:
                    texts.append((shape.shape_type, txt, shape.left.inches, shape.top.inches, shape.width.inches, shape.height.inches))
        except Exception:
            pass
    if any('Search & Discoverability Benchmarking' in t[1] for t in texts) or any('Current Visibility Structure' in t[1] for t in texts) or any('Competitive Search Benchmark' in t[1] for t in texts):
        print('SLIDE', idx)
        for item in texts:
            print(item)
        print('---')
