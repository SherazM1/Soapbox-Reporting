from __future__ import annotations

import io
import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation

from app.audit_helpers.image_guides import get_image_guide_page
from audit_export import build_audit_export_plan
from audit_powerpoint_new import (
    build_slide4_pdp_benchmark_payload,
    generate_new_audit_powerpoint_from_template,
)
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "templates" / "Audit_Template_New.pptx"


def _walk_shapes(shapes):
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from _walk_shapes(shape.shapes)


def _make_images(directory: Path, prefix: str, count: int, size: tuple[int, int]) -> list[dict]:
    images = []
    for index in range(count):
        path = directory / f"{prefix}-{index + 1}.png"
        Image.new(
            "RGB",
            size,
            color=((index * 31) % 255, (index * 53) % 255, (index * 79) % 255),
        ).save(path)
        images.append(
            {
                "index": index,
                "url": path.resolve().as_uri(),
                "is_hero": index == 0,
                "width": 2200,
                "height": 2200,
                "dimensions": "2200 x 2200",
                "dimensions_text": "2200 x 2200",
            }
        )
    return images


def _record(brand: str, images: list[dict], record_id: str) -> dict:
    return {
        "record_id": record_id,
        "brand": brand,
        "product_title": f"{brand} Nut Butter",
        "item_id": record_id,
        "source_url": f"https://example.test/{record_id}",
        "category": "Food/Pantry/Peanut butter & spreads",
        "subcategory": "Nut Butters & Spreads",
        "images": images,
        "image_count": len(images),
        "ingest_metadata": {"content_score": "95"},
        "reviews_summary": {},
        "description_body": "Creamy peanut butter spread with protein for pantry snacks.",
        "description_bullets": ["Great for breakfast, snacks, and recipes"],
        "key_features": ["Protein", "Fresh roasted peanut taste"],
        "review_count": 245,
        "average_rating": 4.7,
        "sold_by_walmart": True,
        "shipped_by_walmart": True,
        "enhanced_brand_content": "present",
    }


def _slide4_lines(slide):
    return [
        shape
        for shape in _walk_shapes(slide.shapes)
        if "LINE" in str(getattr(shape, "shape_type", "")).upper()
    ]


def _slide4_vertical_divider_count(slide) -> int:
    return sum(
        1
        for shape in _slide4_lines(slide)
        if int(getattr(shape, "width", 0) or 0) <= int(Inches(0.03))
        and int(getattr(shape, "height", 0) or 0) >= int(Inches(1.0))
    )


def _image_analysis(*, image_count: int, formats: list[str], signals: list[list[str]]) -> dict:
    return {
        "status": "complete",
        "image_count": image_count,
        "analyzed_image_count": len(formats),
        "images": [
            {
                "status": "analyzed",
                "position": index + 1,
                "probable_format": fmt,
                "white_background_ratio": 0.86 if index == 0 else 0.15,
                "detected_signals": signals[index] if index < len(signals) else [],
                "ocr_tokens": ["ingredient", "protein", "spread"] if index == 1 else [],
            }
            for index, fmt in enumerate(formats)
        ],
        "stack_signals": {"duplicate_image_count": 0},
    }


class Slide4IntegrationTest(unittest.TestCase):
    def test_slide4_uses_full_carousels_and_image_guide_recommendations(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            temp_dir = Path(tmp)
            client_record = _record(
                "Client Brand", _make_images(temp_dir, "client", 10, (800, 800)), "client-1"
            )
            competitor_1 = _record(
                "Competitor Alpha", _make_images(temp_dir, "comp1", 7, (800, 600)), "comp-1"
            )
            competitor_2 = _record(
                "Competitor Beta", _make_images(temp_dir, "comp2", 6, (600, 800)), "comp-2"
            )
            primary_entry = {
                "entry_id": "entry-client-1",
                "record_id": "client-1",
                "product_title": client_record["product_title"],
                "item_id": client_record["item_id"],
                "cached_record": client_record,
                "selected_primary_image": {
                    "record_id": "client-1",
                    "image_index": 0,
                    "url": client_record["images"][0]["url"],
                },
                "selected_primary_images": [
                    {
                        "record_id": "client-1",
                        "image_index": 0,
                        "url": client_record["images"][0]["url"],
                    }
                ],
                "include_in_export": True,
            }
            plan = build_audit_export_plan(
                audit_record={
                    "client_name": "Client Company",
                    "retailer": "Walmart",
                    "audit_date": "2026-06-23",
                    "status": "generated_mvp",
                },
                primary_entries=[primary_entry],
                competitor_assignments=[],
                competitor_records=[competitor_1, competitor_2],
            )

            pdp_slide = plan["product_slide_pairs"][0]["pdp_slide"]
            self.assertEqual(len(pdp_slide["ordered_images"]), 10)
            self.assertEqual(pdp_slide["ordered_images"][0]["index"], 0)
            self.assertEqual(pdp_slide["ordered_images"][-1]["index"], 9)
            self.assertEqual(pdp_slide["brand"], "Client Brand")
            self.assertEqual(
                pdp_slide["category"],
                "Food/Pantry/Peanut butter & spreads",
            )
            self.assertEqual(pdp_slide["product_type"], "Nut Butters & Spreads")
            self.assertIn("style_guide_match", pdp_slide)
            self.assertEqual(pdp_slide["ingest_metadata"], {"content_score": "95"})
            self.assertEqual(
                pdp_slide["ordered_images"][0],
                {
                    "index": 0,
                    "url": client_record["images"][0]["url"],
                    "is_hero": True,
                    "width": 2200,
                    "height": 2200,
                    "dimensions": "2200 x 2200",
                    "dimensions_text": "2200 x 2200",
                },
            )

            guide_page = get_image_guide_page(
                "Food/Pantry/Peanut butter & spreads", "Nut Butters & Spreads"
            )
            self.assertIsNotNone(guide_page)
            self.assertEqual(
                guide_page["page_key"],
                "savory_dips_spreads_nut_butters_jarred_single",
            )
            self.assertEqual(
                guide_page["required_slots"],
                [
                    "silo_front",
                    "graphic_ingredients",
                    "silo_alt_in_pack",
                    "lifestyle_in_use",
                    "graphic_guarantee",
                    "feature_graphic",
                    "dimensions",
                    "graphic_nutrition",
                ],
            )

            payload = build_slide4_pdp_benchmark_payload(
                plan, competitor_records=[competitor_1, competitor_2]
            )
            self.assertEqual(payload["layout_mode"], "3-column")
            self.assertEqual(payload["debug"]["layout_debug"]["divider_cleanup"], "preserve_three_column_chrome")
            self.assertEqual(payload["debug"]["layout_debug"]["render_slots"], [0, 1, 2])
            self.assertEqual(
                [column["label"] for column in payload["columns"]],
                ["Client Company", "Competitor Alpha", "Competitor Beta"],
            )
            self.assertEqual(
                [len(column["ordered_images"]) for column in payload["columns"]],
                [10, 7, 6],
            )
            self.assertEqual(
                payload["columns"][0]["image_guide_match"]["required_slots"],
                guide_page["required_slots"],
            )

            deck_bytes = generate_new_audit_powerpoint_from_template(
                export_plan=plan,
                template_path=str(TEMPLATE),
                include_slide_9=True,
                competitor_records=[competitor_1, competitor_2],
            )
            presentation = Presentation(io.BytesIO(deck_bytes))
            deck_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in _walk_shapes(slide.shapes)
                if getattr(shape, "has_text_frame", False)
            )
            self.assertNotIn("{{", deck_text)
            self.assertTrue(
                any(
                    any(
                        "Walmart Cash Program Visibility" in (shape.text or "")
                        for shape in _walk_shapes(slide.shapes)
                        if getattr(shape, "has_text_frame", False)
                    )
                    for slide in presentation.slides
                )
            )
            slide4 = next(
                slide
                for slide in presentation.slides
                if any(
                    "PDP Content Benchmarking" in (shape.text or "")
                    for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False)
                )
            )
            self.assertGreaterEqual(_slide4_vertical_divider_count(slide4), 2)
            all_text = "\n".join(
                shape.text
                for shape in _walk_shapes(slide4.shapes)
                if getattr(shape, "has_text_frame", False)
            )
            for expected in ("Client Company", "Competitor Alpha", "Competitor Beta"):
                self.assertIn(expected, all_text)
            for sample in ("Honest", "CeraVe", "Jergens"):
                self.assertNotIn(sample, all_text)
            self.assertNotIn("{{", all_text)
            self.assertIn("Peanut butter positioning reinforces pantry relevance", all_text)
            self.assertIn("Protein cues strengthen snack and pantry relevance", all_text)
            self.assertIn("Clear pack and nutrition detail", all_text)

            pictures = [shape for shape in slide4.shapes if hasattr(shape, "image")]
            self.assertEqual(len(pictures), 23)
            intro_bottom = max(
                shape.top + shape.height
                for shape in slide4.shapes
                if getattr(shape, "has_text_frame", False)
                and "Strong Walmart PDP patterns" in (shape.text or "")
            )
            bullet_shapes = sorted(
                (
                    shape
                    for shape in slide4.shapes
                    if getattr(shape, "has_text_frame", False)
                    and any(
                        phrase in (shape.text or "").lower()
                        for phrase in ("positioning", "carousel supports", "pack and nutrition")
                    )
                ),
                key=lambda shape: shape.left,
            )
            self.assertEqual(len(bullet_shapes), 3)
            bullet_centers = [
                (
                    shape.left + shape.width / 2,
                    shape.top,
                )
                for shape in bullet_shapes
            ]
            for picture in pictures:
                self.assertGreaterEqual(picture.left, 0)
                self.assertGreaterEqual(picture.top, intro_bottom)
                self.assertLessEqual(picture.left + picture.width, presentation.slide_width)
                self.assertLessEqual(picture.top + picture.height, presentation.slide_height)
                picture_center = picture.left + picture.width / 2
                column_index = min(
                    range(3),
                    key=lambda index: abs(picture_center - bullet_centers[index][0]),
                )
                self.assertLessEqual(
                    picture.top + picture.height,
                    bullet_centers[column_index][1],
                )
                expected_ratio = (1.0, 4 / 3, 3 / 4)[column_index]
                self.assertAlmostEqual(
                    picture.width / picture.height,
                    expected_ratio,
                    places=2,
                )

            without_slide_9 = generate_new_audit_powerpoint_from_template(
                export_plan=plan,
                template_path=str(TEMPLATE),
                include_slide_9=False,
                competitor_records=[competitor_1, competitor_2],
            )
            presentation_without_slide_9 = Presentation(io.BytesIO(without_slide_9))
            self.assertFalse(
                any(
                    any(
                        "Walmart Cash Program Visibility" in (shape.text or "")
                        for shape in _walk_shapes(slide.shapes)
                        if getattr(shape, "has_text_frame", False)
                    )
                    for slide in presentation_without_slide_9.slides
                )
            )

    def test_missing_second_competitor_is_not_duplicated(self) -> None:
        plan = {
            "audit_metadata": {"client_name": "Client Company"},
            "product_slide_pairs": [{"pdp_slide": {"brand": "Client Brand", "ordered_images": []}}],
        }
        payload = build_slide4_pdp_benchmark_payload(
            plan,
            competitor_records=[{"brand": "Only Competitor", "images": []}],
        )
        self.assertEqual(payload["columns"][1]["label"], "Only Competitor")
        self.assertFalse(payload["columns"][2]["active"])
        self.assertIn("Competitor 2", [item["label"] for item in payload["hidden_columns"]])
        self.assertEqual(payload["columns"][2]["ordered_images"], [])
        self.assertEqual(payload["columns"][2]["bullets"], [])

    @unittest.skipUnless(TEMPLATE.exists(), "New strategic template is unavailable")
    def test_one_competitor_clears_competitor_2_in_generated_deck(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            temp_dir = Path(tmp)
            client_record = _record(
                "Nutella", _make_images(temp_dir, "nutella", 8, (800, 800)), "client-1"
            )
            client_record.update(
                {
                    "product_title": "Nutella Hazelnut Spread with Cocoa, Breakfast Favorite",
                    "description_body": "Hazelnut spread with cocoa for breakfast, snack, and recipe occasions.",
                    "description_bullets": ["Great on toast", "Family snacking pantry staple"],
                    "key_features": ["Hazelnut", "Cocoa", "Breakfast"],
                    "review_count": 85,
                }
            )
            competitor_1 = _record(
                "Jif", _make_images(temp_dir, "jif", 7, (800, 800)), "comp-1"
            )
            competitor_1.update(
                {
                    "product_title": "Jif Creamy Peanut Butter, 16 oz Jar",
                    "description_body": "Fresh roasted peanut taste with protein for snacks and pantry use.",
                    "description_bullets": ["Protein per serving", "Creamy spreadability"],
                    "key_features": ["Peanut butter", "Protein", "Pantry"],
                    "review_count": 1200,
                }
            )
            primary_entry = {
                "entry_id": "entry-client-1",
                "record_id": "client-1",
                "product_title": client_record["product_title"],
                "item_id": client_record["item_id"],
                "cached_record": client_record,
                "selected_primary_image": {"record_id": "client-1", "image_index": 0, "url": client_record["images"][0]["url"]},
                "selected_primary_images": [{"record_id": "client-1", "image_index": 0, "url": client_record["images"][0]["url"]}],
                "include_in_export": True,
            }
            plan = build_audit_export_plan(
                audit_record={"client_name": "Nutella", "client_company_name": "Nutella"},
                primary_entries=[primary_entry],
                competitor_assignments=[],
                competitor_records=[competitor_1],
            )
            payload = build_slide4_pdp_benchmark_payload(plan, competitor_records=[competitor_1])
            self.assertEqual(payload["layout_mode"], "2-column")
            layout_debug = payload["debug"]["layout_debug"]
            self.assertEqual(layout_debug["active_column_count"], 2)
            self.assertEqual(layout_debug["render_slots"], [0, 2])
            self.assertEqual(layout_debug["divider_cleanup"], "suppress_middle_column_label_rule_and_vertical_3_column_dividers")
            client_bullets = payload["columns"][0]["bullets"]
            competitor_bullets = payload["columns"][1]["bullets"]
            self.assertEqual(len(client_bullets), 6)
            self.assertEqual(len(competitor_bullets), 6)
            self.assertFalse(set(client_bullets) & set(competitor_bullets))
            self.assertEqual(len(client_bullets + competitor_bullets), len(set(client_bullets + competitor_bullets)))
            self.assertTrue(any("hazelnut" in bullet.lower() or "cocoa" in bullet.lower() for bullet in client_bullets))
            self.assertTrue(any("breakfast" in bullet.lower() or "snack" in bullet.lower() for bullet in client_bullets))
            self.assertTrue(any("peanut" in bullet.lower() or "protein" in bullet.lower() for bullet in competitor_bullets))
            self.assertTrue(any("pantry" in bullet.lower() or "nutrition" in bullet.lower() for bullet in competitor_bullets))
            forbidden = " ".join(client_bullets + competitor_bullets).lower()
            for term in ("sales", "rank", "share of search", "best-in-class", "secondary benchmark", "benchmark cue"):
                self.assertNotIn(term, forbidden)

            deck_bytes = generate_new_audit_powerpoint_from_template(
                export_plan=plan,
                template_path=str(TEMPLATE),
                include_slide_9=False,
                competitor_records=[competitor_1],
            )
            presentation = Presentation(io.BytesIO(deck_bytes))
            slide4 = next(
                slide
                for slide in presentation.slides
                if any(
                    "PDP Content Benchmarking" in (shape.text or "")
                    for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False)
                )
            )
            all_text = "\n".join(
                shape.text
                for shape in _walk_shapes(slide4.shapes)
                if getattr(shape, "has_text_frame", False)
            )
            self.assertIn("Nutella", all_text)
            self.assertIn("Jif", all_text)
            self.assertNotIn("Competitor 2", all_text)
            self.assertNotIn("CeraVe", all_text)
            self.assertNotIn("Jergens", all_text)
            self.assertIn("hazelnut", all_text.lower())
            self.assertIn("peanut", all_text.lower())
            self.assertIn("benchmark", all_text.lower())
            self.assertEqual(_slide4_vertical_divider_count(slide4), 0)
            bullet_shapes = sorted(
                [
                    shape
                    for shape in _walk_shapes(slide4.shapes)
                    if getattr(shape, "has_text_frame", False)
                    and any(
                        phrase in (shape.text or "").lower()
                        for phrase in ("hazelnut", "peanut", "protein", "breakfast")
                    )
                ],
                key=lambda shape: shape.left,
            )
            self.assertEqual(len(bullet_shapes), 2)
            rendered_font_sizes = {
                run.font.size
                for shape in bullet_shapes
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
                if paragraph.text.strip() and run.text.strip()
            }
            rendered_line_spacing = {
                paragraph.line_spacing
                for shape in bullet_shapes
                for paragraph in shape.text_frame.paragraphs
                if paragraph.text.strip()
            }
            rendered_space_after = {
                paragraph.space_after
                for shape in bullet_shapes
                for paragraph in shape.text_frame.paragraphs
                if paragraph.text.strip()
            }
            self.assertEqual(len(rendered_font_sizes), 1)
            self.assertEqual(rendered_font_sizes, {Pt(11)})
            self.assertEqual(rendered_line_spacing, {0.9})
            self.assertEqual(len(rendered_space_after), 1)
            self.assertEqual(payload["debug"]["render_targets"]["target_bullet_count"], 6)
            self.assertEqual(
                set(payload["debug"]["render_targets"]["final_bullet_counts"].values()),
                {6},
            )
            centers = [shape.left + shape.width / 2 for shape in bullet_shapes]
            self.assertAlmostEqual(
                sum(centers) / 2,
                presentation.slide_width / 2,
                delta=150000,
            )
            self.assertGreater(bullet_shapes[0].left, 1_000_000)
            self.assertLess(
                bullet_shapes[1].left + bullet_shapes[1].width,
                presentation.slide_width - 1_000_000,
            )

    def test_slide4_balances_bullet_families_and_uses_guide_debug(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            temp_dir = Path(tmp)
            client_record = _record(
                "Client Brand", _make_images(temp_dir, "client-family", 6, (800, 800)), "client-1"
            )
            competitor_1 = _record(
                "Competitor Alpha", _make_images(temp_dir, "comp-family", 6, (800, 800)), "comp-1"
            )
            plan = {
                "audit_metadata": {"client_name": "Client Brand"},
                "product_slide_pairs": [{"pdp_slide": client_record}],
            }
            payload = build_slide4_pdp_benchmark_payload(plan, competitor_records=[competitor_1])

            for column in payload["columns"][:2]:
                self.assertEqual(len(column["bullets"]), 6)
                families = [debug.get("bullet_family") for debug in column["bullet_debug"]]
                self.assertGreaterEqual(len(set(families)), 3)
                self.assertLessEqual(max(families.count(family) for family in set(families)), 2)
                for debug in column["bullet_debug"]:
                    self.assertIn(debug.get("evidence_family_source"), {
                        "title/product_identification",
                        "title_formula",
                        "pdp_detail_compliance",
                        "style_title_and_pdp_detail",
                        "pdp_storytelling",
                        "image_support",
                        "image_guide_support",
                        "review_trust",
                        "strategic_cue_engine",
                    })
                    self.assertTrue(debug.get("product_type_context"))
                    self.assertTrue(debug.get("guide_context"))
                    self.assertTrue(debug.get("selected_because"))
            guide_debug = payload["columns"][0]["bullet_debug"][0]["evidence_context"]["identity"]
            self.assertIn("style_guides", guide_debug["style_guide_path"])
            self.assertIn("image_guides", guide_debug["image_guide_path"])

    def test_slide4_blocks_wrong_category_language_for_beauty(self) -> None:
        client_record = {
            "brand": "Glow",
            "product_title": "Glow Hydrating Facial Cleanser for Sensitive Skin",
            "item_id": "beauty-client",
            "category": "Beauty/Skin Care",
            "product_type": "Facial Cleansers",
            "images": [],
            "image_count": 0,
            "description_body": "Gentle cleanser for daily use with hydrating skin benefits.",
            "key_features": ["Hydrating", "Sensitive skin", "Fragrance free"],
            "review_count": 42,
        }
        competitor_record = {
            "brand": "Cera",
            "product_title": "Cera Foaming Facial Cleanser",
            "item_id": "beauty-comp",
            "category": "Beauty/Skin Care",
            "product_type": "Facial Cleansers",
            "images": [],
            "image_count": 0,
            "description_body": "Daily face wash for normal skin.",
            "key_features": ["Cleanser", "Daily use"],
            "review_count": 800,
        }
        payload = build_slide4_pdp_benchmark_payload(
            {
                "audit_metadata": {"client_name": "Glow"},
                "product_slide_pairs": [{"pdp_slide": client_record}],
            },
            competitor_records=[competitor_record],
        )
        joined = " ".join(
            bullet
            for column in payload["columns"][:2]
            for bullet in column["bullets"]
        ).lower()
        self.assertIn("facial cleanser", joined)
        self.assertTrue("hydrating" in joined or "sensitive" in joined or "skin-benefit" in joined)
        for forbidden in ("nutrition", "protein", "breakfast", "snack", "recipe", "pantry"):
            self.assertNotIn(forbidden, joined)
        self.assertEqual(payload["columns"][0]["bullet_debug"][0]["evidence_context"]["identity"]["category_key"], "beauty")

    def test_slide4_uses_evidence_based_bullets_when_image_analysis_exists(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            temp_dir = Path(tmp)
            client_record = _record(
                "Client Brand", _make_images(temp_dir, "client-evidence", 6, (800, 800)), "client-1"
            )
            client_record["image_analysis"] = _image_analysis(
                image_count=6,
                formats=["product_silo", "text_heavy_graphic", "nutrition_or_ingredients"],
                signals=[[], ["feature_or_benefit_claim"], ["nutrition_or_ingredients"]],
            )
            competitor_1 = _record(
                "Competitor Alpha", _make_images(temp_dir, "comp-evidence", 6, (800, 800)), "comp-1"
            )
            competitor_1["image_analysis"] = _image_analysis(
                image_count=6,
                formats=["product_silo", "lifestyle_or_scene", "mixed_product_graphic"],
                signals=[[], ["recipe_or_serving"], ["feature_or_benefit_claim"]],
            )
            primary_entry = {
                "entry_id": "entry-client-1",
                "record_id": "client-1",
                "product_title": client_record["product_title"],
                "item_id": client_record["item_id"],
                "cached_record": client_record,
                "selected_primary_image": {
                    "record_id": "client-1",
                    "image_index": 0,
                    "url": client_record["images"][0]["url"],
                },
                "selected_primary_images": [
                    {
                        "record_id": "client-1",
                        "image_index": 0,
                        "url": client_record["images"][0]["url"],
                    }
                ],
                "include_in_export": True,
            }
            plan = build_audit_export_plan(
                audit_record={
                    "client_name": "Client Company",
                    "retailer": "Walmart",
                    "audit_date": "2026-06-23",
                    "status": "generated_mvp",
                },
                primary_entries=[primary_entry],
                competitor_assignments=[],
                competitor_records=[competitor_1],
            )
            self.assertEqual(plan["slide4_findings"]["client"]["analyzed_pdp_count"], 1)
            payload = build_slide4_pdp_benchmark_payload(
                plan, competitor_records=[competitor_1]
            )
            self.assertTrue(any("peanut" in bullet.lower() or "nutrition" in bullet.lower() for bullet in payload["columns"][0]["bullets"]))
            self.assertTrue(any("carousel" in bullet.lower() or "visual" in bullet.lower() for bullet in payload["columns"][1]["bullets"]))
            self.assertTrue(any("benefit" in bullet.lower() or "recipe" in bullet.lower() or "education" in bullet.lower() for bullet in payload["columns"][1]["bullets"]))

            deck_bytes = generate_new_audit_powerpoint_from_template(
                export_plan=plan,
                template_path=str(TEMPLATE),
                include_slide_9=False,
                competitor_records=[competitor_1],
            )
            presentation = Presentation(io.BytesIO(deck_bytes))
            slide4 = next(
                slide
                for slide in presentation.slides
                if any(
                    "PDP Content Benchmarking" in (shape.text or "")
                    for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False)
                )
            )
            all_text = "\n".join(
                shape.text
                for shape in _walk_shapes(slide4.shapes)
                if getattr(shape, "has_text_frame", False)
            )
            self.assertIn("Client Company", all_text)
            self.assertIn("Competitor Alpha", all_text)
            self.assertIn("Peanut butter positioning reinforces pantry relevance", all_text)
            self.assertNotIn("Carousel: 6 ordered images", all_text)
            self.assertGreaterEqual(len([shape for shape in slide4.shapes if hasattr(shape, "image")]), 12)


if __name__ == "__main__":
    unittest.main()
