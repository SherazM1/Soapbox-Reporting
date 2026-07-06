import io
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from app.audit_helpers.slide2_summary import (
    COMPETITIVE_BENCHMARK_LABELS,
    CONSUMER_DEMAND_LABELS,
    RATING_SCALES,
    WALMART_OPPORTUNITY_LABELS,
    _validate_rating,
    build_slide2_summary_payload,
    resolve_slide2_phrases,
)
from audit_export import build_audit_export_plan
from audit_powerpoint import generate_audit_powerpoint_from_template, resolve_audit_template_path
from audit_powerpoint_new import generate_new_audit_powerpoint_from_template


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "templates" / "Audit_Template_New.pptx"


def _walk_shapes(shapes):
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from _walk_shapes(shape.shapes)


def _record(
    *,
    category: str = "Food/Pantry/Jams, Jellies & Preserves",
    product_type: str = "Jams, Jellies & Preserves",
    title: str = "Strawberry Preserves",
    rating: float | None = 4.6,
    rating_count: int | None = 120,
    gap_count: int = 0,
) -> dict:
    reviews = {}
    if rating is not None:
        reviews["average_rating"] = rating
    if rating_count is not None:
        reviews["ratings_count"] = rating_count
    return {
        "record_id": f"{product_type}-{rating}-{rating_count}",
        "brand": "Test Brand",
        "product_title": title,
        "category": category,
        "product_type": product_type,
        "subcategory": product_type,
        "images": [],
        "image_count": 0,
        "reviews_summary": reviews,
        "_slide2_gap_count": gap_count,
    }


def _primary_entry(record: dict, *, generated_outputs: dict | None = None) -> dict:
    return {
        "entry_id": record["record_id"],
        "record_id": record["record_id"],
        "product_title": record["product_title"],
        "item_id": "123",
        "cached_record": record,
        "include_in_export": True,
        "generated_outputs": generated_outputs or {},
    }


class Slide2SummaryTest(unittest.TestCase):
    def test_rating_labels_stay_in_section_specific_scales(self) -> None:
        payload = build_slide2_summary_payload(
            primary_records=[_record(gap_count=4)],
            competitor_records=[_record()],
            slide4_findings={
                "client": {"opportunities": [{"signal": "missing_lifestyle_storytelling"}]},
                "competitor_1": {"strengths": [{"signal": "usage_or_recipe_storytelling"} for _ in range(4)]},
                "competitor_2": {"strengths": []},
            },
            audit_metadata={"client_name": "test"},
        )
        sections = payload["sections"]
        self.assertIn(sections["consumer_demand"]["rating"], CONSUMER_DEMAND_LABELS)
        self.assertIn(sections["walmart_opportunity"]["rating"], WALMART_OPPORTUNITY_LABELS)
        self.assertIn(sections["competitive_benchmark"]["rating"], COMPETITIVE_BENCHMARK_LABELS)
        self.assertNotIn(sections["consumer_demand"]["rating"], WALMART_OPPORTUNITY_LABELS)
        self.assertNotIn(sections["walmart_opportunity"]["rating"], CONSUMER_DEMAND_LABELS)
        self.assertNotIn(sections["competitive_benchmark"]["rating"], WALMART_OPPORTUNITY_LABELS)
        self.assertEqual(
            sections["consumer_demand"]["allowed_ratings"],
            list(RATING_SCALES["consumer_demand"]["allowed"]),
        )

    def test_invalid_rating_labels_fall_back_to_section_defaults(self) -> None:
        warnings: list[str] = []
        self.assertEqual(_validate_rating("consumer_demand", "Significant", warnings), "Emerging")
        self.assertEqual(_validate_rating("walmart_opportunity", "Strong", warnings), "Meaningful")
        self.assertEqual(_validate_rating("competitive_benchmark", "Meaningful", warnings), "Evolving")
        self.assertEqual(len(warnings), 3)
        self.assertTrue(all("replaced with default" in warning for warning in warnings))

    def test_phrase_resolver_supported_categories_and_generic_fallback(self) -> None:
        self.assertEqual(
            resolve_slide2_phrases([_record(category="Baby/Bath", product_type="Baby Wash")])["category_phrase"],
            "baby care",
        )
        self.assertEqual(
            resolve_slide2_phrases([_record(category="Food", product_type="Jams, Jellies & Preserves")])["category_phrase"],
            "jams and fruit spreads",
        )
        self.assertEqual(
            resolve_slide2_phrases([_record(category="Food", product_type="Nut Butters & Spreads")])["category_phrase"],
            "nut butter and spreads",
        )
        self.assertEqual(
            resolve_slide2_phrases([_record(category="", product_type="", title="Generic Product")])["category_phrase"],
            "the category",
        )

    def test_consumer_demand_rating_from_reviews(self) -> None:
        strong = build_slide2_summary_payload([_record(rating=4.7, rating_count=200)], audit_metadata={"client_name": "test"})
        emerging = build_slide2_summary_payload([_record(rating=3.9, rating_count=25)], audit_metadata={"client_name": "test"})
        limited = build_slide2_summary_payload([_record(rating=None, rating_count=None)], audit_metadata={"client_name": "test"})
        self.assertEqual(strong["sections"]["consumer_demand"]["rating"], "Strong")
        self.assertEqual(emerging["sections"]["consumer_demand"]["rating"], "Emerging")
        self.assertEqual(limited["sections"]["consumer_demand"]["rating"], "Limited")

    def test_walmart_opportunity_rating_from_content_and_image_gaps(self) -> None:
        selective = build_slide2_summary_payload([_record(gap_count=0)], audit_metadata={"client_name": "test"})
        meaningful = build_slide2_summary_payload([_record(gap_count=1)], audit_metadata={"client_name": "test"})
        significant = build_slide2_summary_payload([_record(gap_count=4)], audit_metadata={"client_name": "test"})
        self.assertEqual(selective["sections"]["walmart_opportunity"]["rating"], "Selective")
        self.assertEqual(meaningful["sections"]["walmart_opportunity"]["rating"], "Meaningful")
        self.assertEqual(significant["sections"]["walmart_opportunity"]["rating"], "Significant")

    def test_competitive_benchmark_rating_from_client_vs_competitor_evidence(self) -> None:
        limited = build_slide2_summary_payload([_record()], competitor_records=[], audit_metadata={"client_name": "test"})
        evolving = build_slide2_summary_payload([_record()], competitor_records=[_record(), _record()], audit_metadata={"client_name": "test"})
        competitive = build_slide2_summary_payload(
            [_record()],
            competitor_records=[_record()],
            slide4_findings={
                "client": {"opportunities": [{"signal": "a"}, {"signal": "b"}]},
                "competitor_1": {"strengths": [{"signal": "a"} for _ in range(4)]},
            },
            audit_metadata={"client_name": "test"},
        )
        self.assertEqual(limited["sections"]["competitive_benchmark"]["rating"], "Evolving")
        self.assertEqual(evolving["sections"]["competitive_benchmark"]["rating"], "Evolving")
        self.assertEqual(competitive["sections"]["competitive_benchmark"]["rating"], "Competitive")

    def test_bullet_count_is_four_per_section(self) -> None:
        payload = build_slide2_summary_payload([_record(gap_count=4)], competitor_records=[_record()], audit_metadata={"client_name": "test"})
        for section in payload["sections"].values():
            self.assertEqual(len(section["bullets"]), 4)
            self.assertEqual(len(section["bullets"]), len(section["bullet_debug"]))
            self.assertTrue(all(len(bullet) <= 64 for bullet in section["bullets"]))

    def test_cue_refinement_preserves_section_ownership_and_limits_swaps(self) -> None:
        payload = build_slide2_summary_payload(
            [_record(gap_count=4)],
            competitor_records=[_record(), _record()],
            slide4_findings={
                "client": {"opportunities": [{"signal": "missing_lifestyle_storytelling"} for _ in range(3)]},
                "competitor_1": {"strengths": [{"signal": "usage_or_recipe_storytelling"} for _ in range(4)]},
            },
            audit_metadata={"client_name": "test"},
        )
        for section_key, section in payload["sections"].items():
            source_tags = [item.get("source_tag") for item in section["bullet_debug"]]
            self.assertLessEqual(source_tags.count("cue_refined_swap"), 2, section_key)
            self.assertGreaterEqual(
                sum(1 for tag in source_tags if tag in {"section_bank_original", "section_bank_fallback"}),
                2,
                section_key,
            )
            self.assertEqual(section["cue_refinement_debug"]["mode"], "controlled_swap")
            self.assertLessEqual(section["cue_refinement_debug"]["swap_count"], 2)

        consumer_text = " ".join(payload["sections"]["consumer_demand"]["bullets"]).lower()
        opportunity_text = " ".join(payload["sections"]["walmart_opportunity"]["bullets"]).lower()
        benchmark_text = " ".join(payload["sections"]["competitive_benchmark"]["bullets"]).lower()
        all_bullets_text = " ".join(
            bullet
            for section in payload["sections"].values()
            for bullet in section["bullets"]
        ).lower()
        self.assertTrue(any(term in consumer_text for term in ("trust", "review", "confidence", "fit")))
        self.assertTrue(any(term in opportunity_text for term in ("walmart", "pdp", "shelf", "guidance", "conversion")))
        self.assertTrue(any(term in benchmark_text for term in ("competitor", "benchmark", "category leaders")))
        self.assertNotIn("more ownable", all_bullets_text)
        self.assertNotIn("value communication", all_bullets_text)

    def test_slide2_debug_explains_ratings_sources_and_final_validation(self) -> None:
        payload = build_slide2_summary_payload(
            [_record(gap_count=4)],
            competitor_records=[_record()],
            audit_metadata={"client_name": "test"},
        )
        for section_key, section in payload["sections"].items():
            self.assertTrue(section["rating_reason"])
            self.assertEqual(section["rating_inputs"]["section_key"], section_key)
            self.assertEqual(section["rating_inputs"]["validated_rating"], section["rating"])
            self.assertEqual(section["rating_signals"], section["signals"])
            self.assertTrue(all(item.get("source_tag") for item in section["bullet_debug"]))
            validation = section["final_validation"]
            self.assertEqual(validation["final_bullet_count"], 4)
            self.assertEqual(validation["dedupe_result"], "passed")
            self.assertEqual(validation["fit_result"], "passed")
            self.assertTrue(validation["final_rating_valid"])
            self.assertTrue(validation["required_count_met"])
            for swap in section["cue_refinement_debug"]["accepted_swaps"]:
                self.assertIn("original_bullet_text", swap)
                self.assertIn("replacement_bullet_text", swap)
                self.assertIn("reason", swap)

    def test_section_fallback_refill_preserves_four_bullets_after_dedupe(self) -> None:
        payload = build_slide2_summary_payload(
            [_record(category="", product_type="", title="Generic Product", rating=None, rating_count=None)],
            competitor_records=[],
            audit_metadata={"client_name": "test"},
        )
        for section in payload["sections"].values():
            self.assertEqual(len(section["bullets"]), 4)
            self.assertTrue(section["final_validation"]["required_count_met"])
        source_tags = [
            item.get("source_tag")
            for section in payload["sections"].values()
            for item in section["bullet_debug"]
        ]
        self.assertIn("section_bank_fallback", source_tags)

    def test_final_guardrail_blocks_wrong_category_and_robotic_consumer_language(self) -> None:
        payload = build_slide2_summary_payload(
            [
                _record(
                    category="Beauty/Skin Care",
                    product_type="Facial Cleansers",
                    title="Hydrating Facial Cleanser",
                    rating=None,
                    rating_count=None,
                )
            ],
            competitor_records=[],
            audit_metadata={"client_name": "test"},
        )
        sections = payload["sections"]
        self.assertEqual(
            {key: len(section["bullets"]) for key, section in sections.items()},
            {
                "consumer_demand": 4,
                "walmart_opportunity": 4,
                "competitive_benchmark": 4,
            },
        )
        consumer_text = " ".join(sections["consumer_demand"]["bullets"]).lower()
        all_text = " ".join(
            bullet
            for section in sections.values()
            for bullet in section["bullets"]
        ).lower()
        self.assertTrue(any(term in consumer_text for term in ("trust", "review", "relevant", "confidence", "fit")))
        for forbidden in ("pantry routine", "breakfast", "snack", "recipe", "signals support", "cues support", "support confidence"):
            self.assertNotIn(forbidden, all_text)
        for section in sections.values():
            self.assertTrue(section["final_validation"]["required_count_met"])

    def test_bullet_debug_is_traceable_and_powerpoint_bullets_remain_text_only(self) -> None:
        payload = build_slide2_summary_payload(
            [_record(gap_count=4)],
            competitor_records=[_record(), _record()],
            slide4_findings={
                "client": {
                    "analyzed_pdp_count": 10,
                    "opportunities": [
                        {
                            "signal": "missing_usage_or_recipe_storytelling",
                            "supporting_pdps": 6,
                            "analyzed_pdps": 10,
                        }
                    ],
                },
                "competitor_1": {
                    "strengths": [
                        {
                            "signal": "usage_or_recipe_storytelling",
                            "supporting_pdps": 2,
                            "analyzed_pdps": 2,
                        }
                    ]
                },
            },
            audit_metadata={"client_name": "test"},
        )
        opportunity = payload["sections"]["walmart_opportunity"]
        self.assertTrue(all(isinstance(bullet, str) for bullet in opportunity["bullets"]))
        traced = next(item for item in opportunity["cue_translation_debug"])
        self.assertIn(traced["classification"], {"opportunity", "pressure", "context", "strength"})
        self.assertIn("coverage_ratio", traced)
        self.assertIn("Translated from aggregated strategic cue evidence", traced["reason"])
        self.assertIn("strategic_cues", payload["debug"])
        self.assertIn("candidate_cues", payload["debug"]["strategic_cues"])

    def test_missing_data_returns_defaults_bullets_and_warnings(self) -> None:
        payload = build_slide2_summary_payload([], competitor_records=[], audit_metadata={"client_name": "test"})
        self.assertEqual(payload["sections"]["consumer_demand"]["rating"], "Emerging")
        self.assertEqual(payload["sections"]["walmart_opportunity"]["rating"], "Meaningful")
        self.assertEqual(payload["sections"]["competitive_benchmark"]["rating"], "Evolving")
        for section in payload["sections"].values():
            self.assertEqual(len(section["bullets"]), 4)
            self.assertTrue(section["debug_warnings"])
        self.assertTrue(payload["debug"]["warnings"])

    def test_new_strategic_powerpoint_populates_slide2_in_place(self) -> None:
        record = _record(product_type="Nut Butters & Spreads", title="Protein Peanut Butter")
        entry = _primary_entry(
            record,
            generated_outputs={
                "image_recommendations": ["Add usage imagery", "Add benefit graphic"],
                "description_recommendations": ["Add clearer shopper guidance"],
                "key_features_recommendations": [],
                "recommended_title": "Better Title",
            },
        )
        competitor = _record(product_type="Nut Butters & Spreads", title="Competitor Peanut Butter")
        plan = build_audit_export_plan(
            audit_record={
                "client_name": "test",
                "retailer": "Walmart",
                "audit_date": "2026-06-23",
                "status": "generated_mvp",
            },
            primary_entries=[entry],
            competitor_assignments=[],
            competitor_records=[competitor],
        )
        deck_bytes = generate_new_audit_powerpoint_from_template(
            export_plan=plan,
            template_path=str(TEMPLATE),
            include_slide_9=False,
            competitor_records=[competitor],
        )
        presentation = Presentation(io.BytesIO(deck_bytes))
        slide2 = next(
            slide
            for slide in presentation.slides
            if any(
                "Walmart eCommerce Opportunity" in (shape.text or "")
                for shape in _walk_shapes(slide.shapes)
                if getattr(shape, "has_text_frame", False)
            )
        )
        all_text = "\n".join(
            shape.text
            for shape in _walk_shapes(slide2.shapes)
            if getattr(shape, "has_text_frame", False)
        )
        self.assertIn("test", all_text)
        self.assertNotIn("Honest", all_text)
        self.assertNotIn("The Honest Company", all_text)
        self.assertIn("nut butter and spreads", all_text)
        self.assertIn("Consumer Demand", all_text)
        self.assertIn("Walmart Opportunity", all_text)
        self.assertIn("Competitive Benchmark", all_text)
        self.assertIn(plan["slide2_summary"]["sections"]["consumer_demand"]["rating"], all_text)
        self.assertIn(plan["slide2_summary"]["sections"]["walmart_opportunity"]["rating"], all_text)
        self.assertIn(plan["slide2_summary"]["sections"]["competitive_benchmark"]["rating"], all_text)
        render_fit = plan["slide2_summary"]["debug"]["render_fit"]
        self.assertEqual(
            set(render_fit),
            {"consumer_demand", "walmart_opportunity", "competitive_benchmark"},
        )
        self.assertEqual(
            len({item["font_fallback"] for item in render_fit.values()}),
            1,
        )

    def test_old_current_powerpoint_export_still_generates(self) -> None:
        record = _record()
        plan = build_audit_export_plan(
            audit_record={
                "client_name": "test",
                "retailer": "Walmart",
                "audit_date": "2026-06-23",
                "status": "generated_mvp",
            },
            primary_entries=[_primary_entry(record)],
            competitor_assignments=[],
            competitor_records=[],
        )
        deck_bytes = generate_audit_powerpoint_from_template(
            export_plan=plan,
            template_path=resolve_audit_template_path(),
        )
        self.assertGreater(len(deck_bytes), 1000)


if __name__ == "__main__":
    unittest.main()
