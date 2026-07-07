from __future__ import annotations

import base64
import contextlib
import io
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

from app.audit_helpers.slide5_brand_shop import (
    DIMENSION_PRIORITY,
    build_slide5_brand_shop,
)
from audit_powerpoint_new import _decode_data_image, generate_new_audit_powerpoint_from_template


ROOT = Path(__file__).resolve().parents[1]
NEW_TEMPLATE = ROOT / "templates" / "Audit_Template_New.pptx"

CONSTANT_TEXT = (
    "Brand Shop Content Benchmarking",
    "Leading Walmart brands are evolving Brand Shops beyond static product collections into immersive, conversion-focused shopping destinations that drive discovery, education, and cross-category engagement.",
    "Current Structure",
    "Competitive Benchmark",
    "Competitors are leveraging Walmart Brand Shops as full-funnel merchandising environments, not just product collections.",
)
CLIENT_SAMPLE_BULLETS = (
    "Strong product-led presentation",
    "Cohesive brand identity and visual consistency",
    "Clean, premium lifestyle imagery",
)
COMPETITOR_SAMPLE_BULLETS = (
    "Lifestyle-led merchandising journeys",
    "Advanced category segmentation",
    "Educational and promotional storytelling",
)


def _image_data_url(color: tuple[int, int, int], size: tuple[int, int] = (900, 1400)) -> str:
    buffer = io.BytesIO()
    Image.new("RGB", size, color).save(buffer, format="JPEG", quality=85)
    return "data:image/jpeg;base64," + base64.b64encode(buffer.getvalue()).decode("ascii")


def _rich_capture(role: str, row: int, brand: str, color: tuple[int, int, int]) -> dict:
    return {
        "sourceRow": row,
        "role": role,
        "inputBrandName": brand,
        "extractedBrandName": brand,
        "extractionStatus": "success",
        "screenshotDataUrl": _image_data_url(color),
        "modules": [
            {"type": "HeroPov", "heading": "Hydration for everyday care"},
            {"type": "PovCard", "heading": "Body care solutions"},
            {"type": "PovCard", "heading": "Daily moisture support"},
            {"type": "ItemCarousel"},
            {"type": "ItemCarousel"},
            {"type": "VideoPlayer", "title": "The hydration story"},
        ],
        "categoryNavigation": ["Skin Care", "Face Care", "Hair Care", "Body Care"],
        "videoPresent": True,
        "videoTitles": ["The hydration story"],
        "productCount": 12,
        "editorialHeadings": [
            "Hydration for everyday care",
            "Body care solutions",
            "Daily moisture support",
        ],
        "promotionalCopy": [
            "Explore hydration benefits for everyday use",
            "Discover product solutions across body care",
        ],
        "destinationLinks": [
            {"label": "Skin Care", "url": "/skin"},
            {"label": "Face Care", "url": "/face"},
            {"label": "Hair Care", "url": "/hair"},
            {"label": "Body Care", "url": "/body"},
        ],
    }


def _weak_capture(role: str, row: int) -> dict:
    return {
        "sourceRow": row,
        "role": role,
        "extractionStatus": "partial",
        "screenshotDataUrl": _image_data_url((180, 180, 180)),
        "modules": [{"type": "ItemCarousel"}],
        "productCount": 3,
    }


def _nested_capture(role: str, row: int, brand: str, color: tuple[int, int, int]) -> dict:
    data = _rich_capture(role, row, brand, color)
    return {
        "sourceRow": row,
        "role": role,
        "originalRole": role,
        "status": "success",
        "data": data,
    }


def _slide5(prs: Presentation):
    return next(
        slide
        for slide in prs.slides
        if any(
            getattr(shape, "has_text_frame", False)
            and shape.text.strip() == "Brand Shop Content Benchmarking"
            for shape in slide.shapes
        )
    )


def _slide_texts(slide) -> list[str]:
    return [
        shape.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    ]


class Slide5BrandShopTests(unittest.TestCase):
    def test_standard_mode_remains_default(self) -> None:
        payload = build_slide5_brand_shop(
            [_rich_capture("Client", 7, "Good Good", (20, 80, 140))],
            [_rich_capture("Competitor", 8, "Palmer's", (140, 60, 20))],
        )
        self.assertEqual(payload["mode"], "standard")
        self.assertTrue(payload["client_has_brand_shop"])
        self.assertIsNotNone(payload["client"])
        self.assertEqual(len(payload["client"]["bullet_debug"]), 7)
        self.assertTrue(
            all(item["template_id"].startswith("strategic_cue_") for item in payload["client"]["bullet_debug"])
        )

    def test_no_brand_shop_mode_uses_only_competitor_evidence(self) -> None:
        payload = build_slide5_brand_shop(
            [_rich_capture("Client", 7, "Uploaded Client", (20, 80, 140))],
            [_rich_capture("Competitor", 8, "Palmer's", (140, 60, 20))],
            client_has_brand_shop=False,
        )
        self.assertEqual(payload["mode"], "no_brand_shop")
        self.assertFalse(payload["client_has_brand_shop"])
        self.assertIsNone(payload["client"])
        self.assertTrue(any("user selection was honored" in warning for warning in payload["warnings"]))
        bullets = payload["competitor"]["bullet_debug"]
        self.assertEqual(len(bullets), 7)
        self.assertEqual(
            [item["type"] for item in bullets],
            [
                "competitor_strength",
                "competitor_strength",
                "competitor_strength",
                "competitor_strength",
                "competitor_strength",
                "importance",
                "client_opportunity",
            ],
        )
        strength_dimensions = [
            item["dimension"] for item in bullets if item["type"] == "competitor_strength"
        ]
        self.assertEqual(len(strength_dimensions), len(set(strength_dimensions)))
        self.assertEqual(len({item["text"].lower() for item in bullets}), 7)

    def test_role_mapping_and_earliest_valid_source_row(self) -> None:
        clients = [
            _rich_capture("Client", 11, "Later Client", (20, 80, 140)),
            _rich_capture("Competitor", 1, "Wrong Role", (10, 10, 10)),
            _rich_capture("Client", 7, "Selected Client", (30, 90, 150)),
        ]
        competitors = [
            _rich_capture("Competitor", 12, "Later Competitor", (140, 60, 20)),
            _rich_capture("Client", 2, "Wrong Role", (10, 10, 10)),
            _rich_capture("Competitor", 8, "Selected Competitor", (150, 70, 30)),
        ]
        payload = build_slide5_brand_shop(clients, competitors)
        self.assertEqual(payload["client"]["source_row"], 2)
        self.assertEqual(payload["client"]["brand_name"], "Wrong Role")
        self.assertEqual(payload["competitor"]["source_row"], 1)
        self.assertEqual(payload["competitor"]["brand_name"], "Wrong Role")
        self.assertEqual(payload["debug"]["unused_client_source_rows"], [7, 11])
        self.assertEqual(payload["debug"]["unused_competitor_source_rows"], [8, 12])
        self.assertTrue(any("source row 2" in warning for warning in payload["warnings"]))

    def test_nested_schema2_screenshot_and_module_fields_are_supported(self) -> None:
        payload = build_slide5_brand_shop(
            [_nested_capture("Client", 7, "Nested Client", (20, 80, 140))],
            [_nested_capture("Competitor", 8, "Nested Competitor", (140, 60, 20))],
        )
        self.assertEqual(payload["client"]["brand_name"], "Nested Client")
        self.assertEqual(payload["competitor"]["brand_name"], "Nested Competitor")
        self.assertTrue(payload["client"]["screenshot"].startswith("data:image/jpeg;base64,"))
        self.assertTrue(payload["competitor"]["screenshot"].startswith("data:image/jpeg;base64,"))
        self.assertEqual(len(payload["client"]["bullets"]), 7)
        self.assertEqual(len(payload["competitor"]["bullets"]), 7)

    def test_full_page_screenshot_is_preferred_over_viewport_capture(self) -> None:
        viewport = _image_data_url((20, 80, 140), size=(500, 500))
        full_page = _image_data_url((30, 90, 150), size=(900, 1600))
        client = _rich_capture("Client", 7, "Full Page Client", (20, 80, 140))
        client["screenshotDataUrl"] = viewport
        client["fullPageScreenshotDataUrl"] = full_page
        payload = build_slide5_brand_shop(
            [client],
            [_rich_capture("Competitor", 8, "Competitor", (140, 60, 20))],
        )
        self.assertEqual(payload["client"]["screenshot"], full_page)
        self.assertEqual(payload["client"]["screenshot_source"], "fullPageScreenshotDataUrl")

    def test_current_and_benchmark_role_aliases_are_mapped_without_brand_inference(self) -> None:
        payload = build_slide5_brand_shop(
            [_rich_capture("Current", 7, "Current Alias", (20, 80, 140))],
            [_rich_capture("Benchmark", 8, "Benchmark Alias", (140, 60, 20))],
        )
        self.assertEqual(payload["client"]["brand_name"], "Current Alias")
        self.assertEqual(payload["competitor"]["brand_name"], "Benchmark Alias")

    def test_all_dimensions_and_exact_bullet_composition(self) -> None:
        payload = build_slide5_brand_shop(
            [_rich_capture("Client", 7, "Good Good", (20, 80, 140))],
            [_rich_capture("Competitor", 8, "Palmer's", (140, 60, 20))],
        )
        client = payload["client"]
        competitor = payload["competitor"]
        self.assertEqual(tuple(client["dimension_scores"]), DIMENSION_PRIORITY)
        self.assertEqual(len(client["bullets"]), 7)
        self.assertEqual(len(competitor["bullets"]), 7)
        client_types = [item["type"] for item in client["bullet_debug"]]
        self.assertGreaterEqual(client_types.count("strength"), 1)
        self.assertIn("strategic_cues", client)
        self.assertEqual(
            len({item["dimension"] for item in client["bullet_debug"]}),
            7,
        )
        self.assertEqual(
            len({item["dimension"] for item in competitor["bullet_debug"]}),
            7,
        )
        self.assertTrue(all(item["type"] for item in competitor["bullet_debug"]))
        self.assertFalse(set(client["bullets"]) & set(competitor["bullets"]))
        all_bullet_text = " ".join(client["bullets"] + competitor["bullets"]).lower()
        self.assertNotIn("uses hero modules", all_bullet_text)
        self.assertNotIn("product modules organize", all_bullet_text)
        self.assertNotIn("merchandising advantage", all_bullet_text)
        self.assertIn("assortment", all_bullet_text)
        self.assertIn("storytelling", all_bullet_text)
        self.assertTrue(
            any(
                "cross_side_mirror_reworded" in item.get("signals", [])
                for item in competitor["bullet_debug"]
            )
        )

    def test_category_specific_language_and_traceability(self) -> None:
        payload = build_slide5_brand_shop(
            [_rich_capture("Client", 7, "Good Good", (20, 80, 140))],
            [_rich_capture("Competitor", 8, "Palmer's", (140, 60, 20))],
        )
        all_bullets = [
            *payload["client"]["bullet_debug"],
            *payload["competitor"]["bullet_debug"],
        ]
        category_bullets = [
            item for item in all_bullets if item["dimension"] in {"category_grouping", "cross_category_navigation"}
        ]
        self.assertTrue(
            any("skin care" in item["text"].lower() and "navigation" in item["text"].lower() for item in category_bullets)
        )
        for item in all_bullets:
            self.assertTrue(item["template_id"])
            self.assertTrue(item["reason"])
            self.assertIn(item["score"], {"Strong", "Present", "Limited", "Missing", "strength", "opportunity", "pressure", "context"})
            self.assertLessEqual(len(item["text"]), 78)
            self.assertLessEqual(len(item["text"].split()), 11)

    def test_weak_evidence_uses_restrained_language(self) -> None:
        payload = build_slide5_brand_shop(
            [_weak_capture("Client", 2)],
            [_weak_capture("Competitor", 3)],
        )
        self.assertTrue(payload["client"]["warnings"])
        self.assertTrue(payload["competitor"]["warnings"])
        self.assertEqual(len(payload["client"]["bullets"]), 7)
        self.assertFalse(
            any("Strong branded" in bullet for bullet in payload["client"]["bullets"])
        )

    def test_screenshot_data_decoding(self) -> None:
        valid = _image_data_url((20, 80, 140))
        self.assertIsNotNone(_decode_data_image(valid))
        self.assertIsNone(_decode_data_image("data:image/png;base64,not-valid"))
        self.assertIsNone(_decode_data_image("https://example.com/image.png"))

    @unittest.skipUnless(NEW_TEMPLATE.exists(), "New strategic template is unavailable")
    def test_powerpoint_populates_only_existing_slide5_targets(self) -> None:
        source = Presentation(str(NEW_TEMPLATE))
        source_slide = _slide5(source)
        source_shape_count = len(source_slide.shapes)
        source_pictures = sorted(
            [shape for shape in source_slide.shapes if hasattr(shape, "image")],
            key=lambda shape: shape.left,
        )
        source_bounds = [
            (shape.left, shape.top, shape.width, shape.height)
            for shape in source_pictures
        ]
        source_picture_blobs = [shape.image.blob for shape in source_pictures]
        source_bullet_boxes = sorted(
            [
                shape
                for shape in source_slide.shapes
                if getattr(shape, "has_text_frame", False)
                and len(shape.text_frame.paragraphs) >= 6
                and shape.top > 2_000_000
            ],
            key=lambda shape: shape.left,
        )
        source_bullet_bounds = [
            (shape.left, shape.top, shape.width, shape.height)
            for shape in source_bullet_boxes
        ]
        payload = build_slide5_brand_shop(
            [_rich_capture("Client", 7, "Good Good", (20, 80, 140))],
            [_rich_capture("Competitor", 8, "Palmer's", (140, 60, 20))],
        )
        deck_bytes = generate_new_audit_powerpoint_from_template(
            export_plan={
                "audit_metadata": {"client_name": "Test"},
                "slide5_brand_shop": payload,
            },
            template_path=str(NEW_TEMPLATE),
            include_slide_9=True,
        )
        generated = Presentation(io.BytesIO(deck_bytes))
        slide = _slide5(generated)
        self.assertEqual(len(slide.shapes), source_shape_count)
        texts = _slide_texts(slide)
        for constant in CONSTANT_TEXT:
            self.assertIn(constant, texts)
        pictures = sorted(
            [shape for shape in slide.shapes if hasattr(shape, "image")],
            key=lambda shape: shape.left,
        )
        self.assertEqual(len(pictures), 2)
        for picture, target in zip(pictures, source_bounds):
            self.assertGreaterEqual(picture.left, target[0])
            self.assertGreaterEqual(picture.top, target[1])
            self.assertLessEqual(picture.left + picture.width, target[0] + target[2])
            self.assertLessEqual(picture.top + picture.height, target[1] + target[3])
        self.assertNotEqual(pictures[0].image.blob, source_picture_blobs[0])
        self.assertNotEqual(pictures[1].image.blob, source_picture_blobs[1])
        bullet_boxes = sorted(
            [
                shape
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
                and len(shape.text_frame.paragraphs) >= 6
                and shape.top > 2_000_000
            ],
            key=lambda shape: shape.left,
        )
        self.assertEqual(
            [(shape.left, shape.top, shape.width, shape.height) for shape in bullet_boxes],
            source_bullet_bounds,
        )
        left_text = [paragraph.text.strip() for paragraph in bullet_boxes[0].text_frame.paragraphs if paragraph.text.strip()]
        right_text = [paragraph.text.strip() for paragraph in bullet_boxes[1].text_frame.paragraphs if paragraph.text.strip()]
        self.assertEqual(left_text, payload["client"]["bullets"])
        self.assertEqual(right_text, payload["competitor"]["bullets"])
        rendered_font_sizes = {
            run.font.size
            for box in bullet_boxes
            for paragraph in box.text_frame.paragraphs
            for run in paragraph.runs
            if paragraph.text.strip() and run.text.strip()
        }
        rendered_line_spacing = {
            paragraph.line_spacing
            for box in bullet_boxes
            for paragraph in box.text_frame.paragraphs
            if paragraph.text.strip()
        }
        rendered_space_after = {
            paragraph.space_after
            for box in bullet_boxes
            for paragraph in box.text_frame.paragraphs
            if paragraph.text.strip()
        }
        self.assertEqual(len(rendered_font_sizes), 1)
        self.assertEqual(rendered_font_sizes, {Pt(11)})
        self.assertEqual(rendered_line_spacing, {0.86})
        self.assertEqual(len(rendered_space_after), 1)
        render_fit = payload["debug"]["render_fit"]
        self.assertEqual({item["target_bullet_count"] for item in render_fit.values()}, {7})
        self.assertEqual({item["rendered_bullet_count"] for item in render_fit.values()}, {7})
        self.assertEqual({item["font_size_selected"] for item in render_fit.values()}, {11})
        self.assertEqual({item["shared_fallback_font_size_used"] for item in render_fit.values()}, {False})
        self.assertTrue(all(item["visible_count_expectation_met"] for item in render_fit.values()))
        for sample in (*CLIENT_SAMPLE_BULLETS, *COMPETITOR_SAMPLE_BULLETS):
            self.assertNotIn(sample, left_text + right_text)
        text_box_count = sum(
            1
            for shape in slide.shapes
            if str(shape.shape_type).endswith("TEXT_BOX (17)")
        )
        source_text_box_count = sum(
            1
            for shape in source_slide.shapes
            if str(shape.shape_type).endswith("TEXT_BOX (17)")
        )
        self.assertEqual(text_box_count, source_text_box_count)

    @unittest.skipUnless(NEW_TEMPLATE.exists(), "New strategic template is unavailable")
    def test_powerpoint_inserts_nested_schema2_brand_shop_screenshots(self) -> None:
        source = Presentation(str(NEW_TEMPLATE))
        source_slide = _slide5(source)
        source_picture_blobs = [
            shape.image.blob
            for shape in sorted(
                [shape for shape in source_slide.shapes if hasattr(shape, "image")],
                key=lambda shape: shape.left,
            )
        ]
        payload = build_slide5_brand_shop(
            [_nested_capture("Client", 7, "Nested Client", (20, 80, 140))],
            [_nested_capture("Competitor", 8, "Nested Competitor", (140, 60, 20))],
        )
        generated = Presentation(
            io.BytesIO(
                generate_new_audit_powerpoint_from_template(
                    export_plan={
                        "audit_metadata": {"client_name": "Test"},
                        "slide5_brand_shop": payload,
                    },
                    template_path=str(NEW_TEMPLATE),
                    include_slide_9=True,
                )
            )
        )
        pictures = sorted(
            [shape for shape in _slide5(generated).shapes if hasattr(shape, "image")],
            key=lambda shape: shape.left,
        )
        self.assertEqual(len(pictures), 2)
        self.assertNotEqual(pictures[0].image.blob, source_picture_blobs[0])
        self.assertNotEqual(pictures[1].image.blob, source_picture_blobs[1])

    @unittest.skipUnless(NEW_TEMPLATE.exists(), "New strategic template is unavailable")
    def test_bad_screenshot_data_warns_and_preserves_template_picture(self) -> None:
        source = Presentation(str(NEW_TEMPLATE))
        source_slide = _slide5(source)
        source_picture_blobs = [
            shape.image.blob
            for shape in sorted(
                [shape for shape in source_slide.shapes if hasattr(shape, "image")],
                key=lambda shape: shape.left,
            )
        ]
        bad_client = _rich_capture("Client", 7, "Bad Client", (20, 80, 140))
        bad_client["screenshotDataUrl"] = "data:image/png;base64,not-valid"
        payload = build_slide5_brand_shop(
            [bad_client],
            [_rich_capture("Competitor", 8, "Good Competitor", (140, 60, 20))],
        )
        output = io.StringIO()
        with contextlib.redirect_stdout(output):
            generated = Presentation(
                io.BytesIO(
                    generate_new_audit_powerpoint_from_template(
                        export_plan={
                            "audit_metadata": {"client_name": "Test"},
                            "slide5_brand_shop": payload,
                        },
                        template_path=str(NEW_TEMPLATE),
                        include_slide_9=True,
                    )
                )
            )
        self.assertIn("screenshot was invalid", output.getvalue())
        pictures = sorted(
            [shape for shape in _slide5(generated).shapes if hasattr(shape, "image")],
            key=lambda shape: shape.left,
        )
        self.assertEqual(pictures[0].image.blob, source_picture_blobs[0])
        self.assertNotEqual(pictures[1].image.blob, source_picture_blobs[1])

    @unittest.skipUnless(NEW_TEMPLATE.exists(), "New strategic template is unavailable")
    def test_missing_side_preserves_that_template_side(self) -> None:
        source = Presentation(str(NEW_TEMPLATE))
        source_slide = _slide5(source)
        source_pictures = sorted(
            [shape for shape in source_slide.shapes if hasattr(shape, "image")],
            key=lambda shape: shape.left,
        )
        source_bullet_boxes = sorted(
            [
                shape
                for shape in source_slide.shapes
                if getattr(shape, "has_text_frame", False)
                and len(shape.text_frame.paragraphs) >= 6
                and shape.top > 2_000_000
            ],
            key=lambda shape: shape.left,
        )
        payload = build_slide5_brand_shop(
            [_rich_capture("Client", 7, "Good Good", (20, 80, 140))],
            [],
        )
        deck_bytes = generate_new_audit_powerpoint_from_template(
            export_plan={
                "audit_metadata": {"client_name": "Test"},
                "slide5_brand_shop": payload,
            },
            template_path=str(NEW_TEMPLATE),
            include_slide_9=True,
        )
        generated = Presentation(io.BytesIO(deck_bytes))
        slide = _slide5(generated)
        pictures = sorted(
            [shape for shape in slide.shapes if hasattr(shape, "image")],
            key=lambda shape: shape.left,
        )
        bullet_boxes = sorted(
            [
                shape
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
                and len(shape.text_frame.paragraphs) >= 6
                and shape.top > 2_000_000
            ],
            key=lambda shape: shape.left,
        )
        self.assertEqual(pictures[1].image.blob, source_pictures[1].image.blob)
        self.assertEqual(bullet_boxes[1].text, source_bullet_boxes[1].text)
        self.assertNotEqual(pictures[0].image.blob, source_pictures[0].image.blob)

    @unittest.skipUnless(NEW_TEMPLATE.exists(), "New strategic template is unavailable")
    def test_no_brand_shop_layout_uses_fixed_existing_shapes(self) -> None:
        payload = build_slide5_brand_shop(
            [_rich_capture("Client", 7, "Uploaded Client", (20, 80, 140))],
            [_rich_capture("Competitor", 8, "Palmer's", (140, 60, 20))],
            client_has_brand_shop=False,
        )
        deck_bytes = generate_new_audit_powerpoint_from_template(
            export_plan={
                "audit_metadata": {
                    "client_name": "Test",
                    "client_has_brand_shop": False,
                },
                "slide5_brand_shop": payload,
            },
            template_path=str(NEW_TEMPLATE),
            include_slide_9=True,
        )
        generated = Presentation(io.BytesIO(deck_bytes))
        slide = _slide5(generated)
        texts = _slide_texts(slide)
        for constant in (
            CONSTANT_TEXT[0],
            CONSTANT_TEXT[1],
            CONSTANT_TEXT[3],
            CONSTANT_TEXT[4],
        ):
            self.assertIn(constant, texts)
        self.assertNotIn("Current Structure", texts)
        self.assertFalse(any(sample in "\n".join(texts) for sample in CLIENT_SAMPLE_BULLETS))
        self.assertEqual(len(slide.shapes), 7)
        header = next(
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and shape.text == "Competitive Benchmark"
        )
        self.assertEqual(header.left, Inches(3.45))
        self.assertEqual(header.width, Inches(3.2))
        divider = next(
            shape
            for shape in slide.shapes
            if "LINE" in str(shape.shape_type).upper()
        )
        self.assertEqual(divider.left, Inches(1.4))
        self.assertEqual(divider.width, Inches(11.2))
        picture = next(shape for shape in slide.shapes if hasattr(shape, "image"))
        target = (Inches(1.4), Inches(2.88), Inches(7.4), Inches(3.77))
        self.assertGreaterEqual(picture.left, target[0])
        self.assertGreaterEqual(picture.top, target[1])
        self.assertLessEqual(picture.left + picture.width, target[0] + target[2])
        self.assertLessEqual(picture.top + picture.height, target[1] + target[3])
        bullet_box = next(
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and len(shape.text_frame.paragraphs) >= 6
        )
        self.assertEqual(bullet_box.left, 8659907)
        bullet_text = [
            paragraph.text.strip()
            for paragraph in bullet_box.text_frame.paragraphs
            if paragraph.text.strip()
        ]
        self.assertEqual(bullet_text, payload["competitor"]["bullets"])
        self.assertEqual(
            {paragraph.line_spacing for paragraph in bullet_box.text_frame.paragraphs if paragraph.text.strip()},
            {0.86},
        )
        text_box_count = sum(
            1 for shape in slide.shapes if str(shape.shape_type).endswith("TEXT_BOX (17)")
        )
        self.assertEqual(text_box_count, 1)

    @unittest.skipUnless(NEW_TEMPLATE.exists(), "New strategic template is unavailable")
    def test_no_brand_shop_missing_competitor_leaves_slide5_unchanged(self) -> None:
        payload = build_slide5_brand_shop(
            [],
            [],
            client_has_brand_shop=False,
        )
        source = Presentation(str(NEW_TEMPLATE))
        source_slide = _slide5(source)
        generated = Presentation(
            io.BytesIO(
                generate_new_audit_powerpoint_from_template(
                    export_plan={
                        "audit_metadata": {
                            "client_name": "Test",
                            "client_has_brand_shop": False,
                        },
                        "slide5_brand_shop": payload,
                    },
                    template_path=str(NEW_TEMPLATE),
                    include_slide_9=True,
                )
            )
        )
        self.assertEqual(source_slide.element.xml, _slide5(generated).element.xml)

    @unittest.skipUnless(NEW_TEMPLATE.exists(), "New strategic template is unavailable")
    def test_tactical_slide_is_removed_from_new_strategic_exports(self) -> None:
        generated = Presentation(
            io.BytesIO(
                generate_new_audit_powerpoint_from_template(
                    export_plan={"audit_metadata": {"client_name": "Test"}},
                    template_path=str(NEW_TEMPLATE),
                    include_slide_9=True,
                )
            )
        )
        all_text = "\n".join(
            shape.text
            for slide in generated.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        self.assertNotIn("If they only have bandwidth for 5 things:", all_text)
        ordered_titles = []
        for slide in generated.slides:
            slide_text = "\n".join(
                shape.text
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            )
            for title in (
                "Walmart eCommerce Opportunity",
                "Search & Discoverability Benchmarking",
                "PDP Content Benchmarking",
                "Brand Shop Content Benchmarking",
                "Digital Shelf Ownership",
            ):
                if title in slide_text:
                    ordered_titles.append(title)
                    break
        self.assertEqual(
            ordered_titles,
            [
                "Walmart eCommerce Opportunity",
                "Search & Discoverability Benchmarking",
                "PDP Content Benchmarking",
                "Brand Shop Content Benchmarking",
                "Digital Shelf Ownership",
            ],
        )

    @unittest.skipUnless(NEW_TEMPLATE.exists(), "New strategic template is unavailable")
    def test_slide5_population_does_not_change_slides_2_3_4_or_6(self) -> None:
        base_plan = {"audit_metadata": {"client_name": "Test"}}
        payload = build_slide5_brand_shop(
            [_rich_capture("Client", 7, "Good Good", (20, 80, 140))],
            [_rich_capture("Competitor", 8, "Palmer's", (140, 60, 20))],
        )
        without_slide5 = Presentation(
            io.BytesIO(
                generate_new_audit_powerpoint_from_template(
                    export_plan=base_plan,
                    template_path=str(NEW_TEMPLATE),
                    include_slide_9=True,
                )
            )
        )
        with_slide5 = Presentation(
            io.BytesIO(
                generate_new_audit_powerpoint_from_template(
                    export_plan={**base_plan, "slide5_brand_shop": payload},
                    template_path=str(NEW_TEMPLATE),
                    include_slide_9=True,
                )
            )
        )
        for slide_number in (2, 3, 4, 6):
            self.assertEqual(
                without_slide5.slides[slide_number - 1].element.xml,
                with_slide5.slides[slide_number - 1].element.xml,
            )

if __name__ == "__main__":
    unittest.main()
