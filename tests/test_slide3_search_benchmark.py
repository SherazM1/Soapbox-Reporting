from __future__ import annotations

import base64
import io
import re
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation

from app.audit_helpers.slide3_search_benchmark import build_slide3_search_benchmark
from app.audit_helpers.slide6_visibility import build_slide6_visibility
from audit_powerpoint_new import generate_new_audit_powerpoint_from_template

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "templates" / "Audit_Template_New.pptx"

CONSTANT_TEXT = (
    "Search & Discoverability Benchmarking",
    "Current Visibility Structure",
    "Competitive Search Benchmark",
    "Competitive Walmart search environments increasingly reward brands that combine strong visual merchandising, clear benefit communication, and trust-based positioning directly within the digital shelf experience.",
)


def _image_data_url(color: tuple[int, int, int]) -> str:
    buffer = io.BytesIO()
    Image.new("RGB", (900, 1200), color).save(buffer, format="JPEG", quality=85)
    return "data:image/jpeg;base64," + base64.b64encode(buffer.getvalue()).decode("ascii")


def _slide3(prs: Presentation):
    return next(
        slide
        for slide in prs.slides
        if any(
            getattr(shape, "has_text_frame", False)
            and "Search & Discoverability Benchmarking" in (shape.text or "")
            for shape in slide.shapes
        )
    )


def _slide_texts(slide) -> list[str]:
    return [shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)]


class Slide3SearchBenchmarkTests(unittest.TestCase):
    def test_build_slide3_search_benchmark_selects_earliest_valid_evidence(self) -> None:
        payload = build_slide3_search_benchmark(
            {
                "current": [
                    {
                        "role": "Current",
                        "sourceRow": 7,
                        "searchTerm": "nut butter spread",
                        "screenshotDataUrl": _image_data_url((10, 20, 30)),
                        "orderedMainResultProducts": [
                            {"position": 1, "title": "Client Brand Nut Butter", "brand": "Client Brand", "sponsored": False},
                            {"position": 2, "title": "Natural Peanut Butter", "brand": "Natural Co", "sponsored": True},
                        ],
                    },
                    {
                        "role": "Current",
                        "sourceRow": 5,
                        "searchTerm": "nut butter spread",
                        "screenshotDataUrl": _image_data_url((20, 30, 40)),
                        "orderedMainResultProducts": [
                            {"position": 1, "title": "Client Brand Nut Butter", "brand": "Client Brand", "sponsored": False},
                            {"position": 2, "title": "Natural Peanut Butter", "brand": "Natural Co", "sponsored": True},
                        ],
                    },
                ],
                "benchmark": [
                    {
                        "role": "Benchmark",
                        "sourceRow": 6,
                        "searchTerm": "low sugar spreads",
                        "screenshotDataUrl": _image_data_url((30, 40, 50)),
                        "orderedMainResultProducts": [
                            {"position": 1, "title": "Low Sugar Fruit Spread", "brand": "Brand A", "sponsored": False},
                            {"position": 2, "title": "Organic Jam", "brand": "Brand B", "sponsored": True},
                        ],
                    }
                ],
            },
            client_name="Client Brand",
        )
        self.assertEqual(payload["current"]["source_row"], 5)
        self.assertEqual(payload["current"]["search_term"], "nut butter spread")
        self.assertEqual(payload["current"]["category_phrase"], "nut butter spreads")
        self.assertEqual(payload["benchmark"]["source_row"], 6)
        self.assertEqual(payload["benchmark"]["search_term"], "low sugar spreads")
        self.assertEqual(payload["benchmark"]["category_phrase"], "low sugar spreads")
        self.assertEqual(len(payload["current"]["bullets"]), 4)
        self.assertEqual(len(payload["benchmark"]["bullets"]), 4)
        self.assertTrue(all(len(bullet.split()) <= 9 for bullet in payload["current"]["bullets"]))
        self.assertTrue(all(len(bullet.split()) <= 9 for bullet in payload["benchmark"]["bullets"]))
        self.assertTrue(any("bullets summarize the broader valid evidence set" in warning.lower() for warning in payload["warnings"]))

    def test_search_term_fallbacks_to_url_query_and_label(self) -> None:
        payload = build_slide3_search_benchmark(
            {
                "current": [
                    {
                        "role": "Current",
                        "sourceRow": 1,
                        "URL": "https://example.test/search?q=strawberry%20jam",
                        "label": "strawberry jam",
                    }
                ],
                "benchmark": [
                    {
                        "role": "Benchmark",
                        "sourceRow": 2,
                        "URL": "https://example.test/search?query=low%20sugar%20fruit%20spread",
                    }
                ],
            },
            client_name="Client Brand",
        )
        self.assertEqual(payload["current"]["search_term"], "strawberry jam")
        self.assertEqual(payload["current"]["category_phrase"], "strawberry jam")
        self.assertEqual(payload["benchmark"]["search_term"], "low sugar fruit spread")
        self.assertEqual(payload["benchmark"]["category_phrase"], "low sugar fruit spreads")

    def test_schema2_nested_search_fields_are_supported(self) -> None:
        payload = build_slide3_search_benchmark(
            {
                "current": [
                    {
                        "role": "Current",
                        "sourceRow": 3,
                        "status": "success",
                        "data": {
                            "searchTerm": "peanut butter spread",
                            "screenshotDataUrl": _image_data_url((40, 50, 60)),
                            "products": [
                                {"position": 2, "title": "Client Brand Peanut Butter Spread", "brand": "Client Brand", "reviewCount": 120, "badges": ["Best seller"]},
                                {"position": 3, "title": "Great Value Peanut Butter", "brand": "Great Value", "reviewCount": 450},
                            ],
                        },
                    }
                ],
                "benchmark": [
                    {
                        "role": "Benchmark",
                        "sourceRow": 4,
                        "data": {
                            "url": "https://example.test/search?q=nut+butter",
                            "screenshotDataUrl": _image_data_url((70, 80, 90)),
                            "capturedProducts": [
                                {"position": 1, "title": "Almond Nut Butter", "brand": "Brand A", "reviewCount": 90},
                                {"position": 2, "title": "Cashew Nut Butter", "brand": "Brand B", "reviewCount": 25},
                            ],
                        },
                    }
                ],
            },
            client_name="Client Brand",
        )
        self.assertEqual(payload["current"]["search_term"], "peanut butter spread")
        self.assertEqual(payload["current"]["category_phrase"], "peanut butter spreads")
        self.assertEqual(payload["benchmark"]["search_term"], "nut butter")
        self.assertEqual(payload["benchmark"]["category_phrase"], "nut butters")
        self.assertEqual(payload["current"]["product_count"], 2)
        self.assertEqual(payload["benchmark"]["product_count"], 2)
        self.assertEqual(payload["current"]["client_products"][0]["position"], 2)
        self.assertIn("Great Value", payload["current"]["top_brands"])
        self.assertIn("Best seller", payload["current"]["badges"])
        self.assertIn(120, payload["current"]["review_counts"])
        self.assertEqual(len({item["dimension"] for item in payload["current"]["bullet_debug"]}), 4)
        self.assertEqual(len({item["dimension"] for item in payload["benchmark"]["bullet_debug"]}), 4)

    def test_slide3_builds_sides_separately_and_suppresses_overlap(self) -> None:
        payload = build_slide3_search_benchmark(
            {
                "current": [
                    {
                        "role": "Current",
                        "sourceRow": 1,
                        "searchTerm": "facial cleanser",
                        "screenshotDataUrl": _image_data_url((12, 24, 36)),
                        "orderedMainResultProducts": [
                            {"position": 5, "title": "Glow Hydrating Facial Cleanser Sensitive Skin", "brand": "Glow", "reviewCount": 42},
                            {"position": 8, "title": "Foaming Face Cleanser", "brand": "Brand B", "reviewCount": 12},
                        ],
                    }
                ],
                "benchmark": [
                    {
                        "role": "Benchmark",
                        "sourceRow": 2,
                        "searchTerm": "facial cleanser",
                        "screenshotDataUrl": _image_data_url((36, 24, 12)),
                        "orderedMainResultProducts": [
                            {"position": 1, "title": "CeraVe Hydrating Facial Cleanser", "brand": "CeraVe", "reviewCount": 1200, "sponsored": True},
                            {"position": 2, "title": "Cetaphil Gentle Skin Cleanser", "brand": "Cetaphil", "reviewCount": 900, "sponsored": True},
                            {"position": 3, "title": "Neutrogena Hydro Boost Cleanser", "brand": "Neutrogena", "reviewCount": 450},
                        ],
                    }
                ],
            },
            client_name="Glow",
        )
        self.assertEqual(len(payload["current"]["bullets"]), 4)
        self.assertEqual(len(payload["benchmark"]["bullets"]), 4)
        self.assertNotEqual(payload["current"]["bullets"], payload["benchmark"]["bullets"])
        self.assertEqual(
            payload["debug"]["side_build_order"],
            [
                "current_candidates_built",
                "benchmark_candidates_built",
                "cross_side_overlap_suppression_applied",
            ],
        )
        self.assertTrue(payload["current"]["debug"]["side_built_separately"])
        self.assertTrue(payload["benchmark"]["debug"]["side_built_separately"])
        overlap = payload["debug"]["overlap_suppression"]["rejected_overlapping_bullets"]
        self.assertTrue(any(item.get("side") == "current" and item.get("replacement") for item in overlap))
        self.assertTrue(any("benefit-led coverage" in bullet.lower() for bullet in payload["current"]["bullets"]))
        self.assertTrue(any("sponsored" in bullet.lower() and "pressure" in bullet.lower() for bullet in payload["benchmark"]["bullets"]))

    def test_slide3_outputs_search_native_phrasing_without_malformed_templates(self) -> None:
        payload = build_slide3_search_benchmark(
            {
                "current": [
                    {
                        "role": "Current",
                        "sourceRow": 1,
                        "searchTerm": "peanut butter",
                        "screenshotDataUrl": _image_data_url((44, 55, 66)),
                        "orderedMainResultProducts": [
                            {"position": 2, "title": "Client Brand Peanut Butter", "brand": "Client Brand", "reviewCount": 220, "badges": ["Best seller"]},
                            {"position": 4, "title": "Natural Peanut Butter", "brand": "Brand B", "reviewCount": 125},
                        ],
                    }
                ],
                "benchmark": [
                    {
                        "role": "Benchmark",
                        "sourceRow": 2,
                        "searchTerm": "natural peanut butter",
                        "screenshotDataUrl": _image_data_url((66, 55, 44)),
                        "orderedMainResultProducts": [
                            {"position": 1, "title": "Jif Natural Peanut Butter", "brand": "Jif", "reviewCount": 1500},
                            {"position": 2, "title": "Skippy Natural Peanut Butter", "brand": "Skippy", "reviewCount": 700},
                        ],
                    }
                ],
            },
            client_name="Client Brand",
        )
        all_bullets = payload["current"]["bullets"] + payload["benchmark"]["bullets"]
        forbidden = " ".join(all_bullets).lower()
        for phrase in (
            "enhanced to improve",
            "review and confidence signals",
            "enhanced active",
            "generic discovery",
            "benchmark cue",
            "benchmark shelf breadth is still concentrated",
            "higher review threshold reaches",
            "visible review threshold",
            "visible reviews average near",
            "median review count near",
            "review counts near",
            "query fit is uneven",
            "query fit is clearest",
            "query fit varies",
        ):
            self.assertNotIn(phrase, forbidden)
        self.assertFalse(any(re.search(r"\d{3,}", bullet) for bullet in all_bullets))
        self.assertTrue(
            all(
                any(term in bullet.lower() for term in ("query", "queries", "shelf", "review", "brand", "position", "badge", "coverage", "sponsored", "trust"))
                for bullet in all_bullets
            )
        )
        self.assertTrue(all(len(bullet.split()) <= 9 for bullet in all_bullets))
        for row in payload["current"]["bullet_debug"] + payload["benchmark"]["bullet_debug"]:
            self.assertIn("bullet_family", row)
            self.assertIn("evidence_summary", row)
            self.assertTrue(row["reason"])

    def test_slide3_suppresses_deck_name_leakage_and_awkward_benchmark_phrasing(self) -> None:
        payload = build_slide3_search_benchmark(
            {
                "current": [
                    {
                        "role": "Current",
                        "sourceRow": 1,
                        "searchTerm": "facial cleanser",
                        "screenshotDataUrl": _image_data_url((12, 34, 56)),
                        "orderedMainResultProducts": [
                            {"position": 3, "title": "Hydrating Face Cleanser", "brand": "Glow", "reviewCount": 42},
                        ],
                    }
                ],
                "benchmark": [
                    {
                        "role": "Benchmark",
                        "sourceRow": 2,
                        "searchTerm": "facial cleanser",
                        "screenshotDataUrl": _image_data_url((56, 34, 12)),
                        "orderedMainResultProducts": [
                            {"position": 1, "title": "Hydrating Facial Cleanser", "brand": "CeraVe", "reviewCount": 15055},
                        ],
                    }
                ],
            },
            client_name="test_deck.pptx",
        )
        all_text = " ".join(payload["current"]["bullets"] + payload["benchmark"]["bullets"]).lower()
        self.assertNotIn("test_deck", all_text)
        self.assertNotIn(".pptx", all_text)
        self.assertNotIn("presence is narrow", all_text)
        self.assertNotIn("benchmark shelf breadth is still concentrated", all_text)
        self.assertNotIn("higher review threshold reaches", all_text)
        self.assertNotIn("15055", all_text)
        self.assertTrue(any("client brand" in bullet.lower() for bullet in payload["current"]["bullets"]))
        self.assertTrue(any("shelf trust" in bullet.lower() for bullet in payload["benchmark"]["bullets"]))

    def test_slide3_reduces_side_label_starters_and_uses_representative_evidence(self) -> None:
        payload = build_slide3_search_benchmark(
            {
                "current": [
                    {
                        "role": "Current",
                        "sourceRow": 1,
                        "searchTerm": "face cleanser",
                        "screenshotDataUrl": _image_data_url((11, 22, 33)),
                        "orderedMainResultProducts": [
                            {"position": 8, "title": "Glow Face Cleanser", "brand": "Glow", "reviewCount": 28},
                        ],
                    },
                    {
                        "role": "Current",
                        "sourceRow": 2,
                        "searchTerm": "sensitive skin cleanser",
                        "screenshotDataUrl": _image_data_url((33, 22, 11)),
                        "orderedMainResultProducts": [
                            {"position": 5, "title": "Gentle Skin Cleanser", "brand": "Brand B", "reviewCount": 34},
                            {"position": 7, "title": "Hydrating Face Wash", "brand": "Brand C", "reviewCount": 42},
                        ],
                    },
                ],
                "benchmark": [
                    {
                        "role": "Benchmark",
                        "sourceRow": 3,
                        "searchTerm": "face cleanser",
                        "screenshotDataUrl": _image_data_url((44, 55, 66)),
                        "orderedMainResultProducts": [
                            {"position": 1, "title": "CeraVe Hydrating Cleanser", "brand": "CeraVe", "reviewCount": 10917, "sponsored": True},
                        ],
                    },
                    {
                        "role": "Benchmark",
                        "sourceRow": 4,
                        "searchTerm": "sensitive skin cleanser",
                        "screenshotDataUrl": _image_data_url((66, 55, 44)),
                        "orderedMainResultProducts": [
                            {"position": 1, "title": "Cetaphil Gentle Cleanser", "brand": "Cetaphil", "reviewCount": 8200, "sponsored": True},
                            {"position": 2, "title": "Neutrogena Skin Cleanser", "brand": "Neutrogena", "reviewCount": 4100},
                        ],
                    },
                ],
            },
            client_name="test7/7.pptx",
        )
        all_bullets = payload["current"]["bullets"] + payload["benchmark"]["bullets"]
        all_text = " ".join(all_bullets).lower()
        self.assertEqual(len(payload["current"]["bullets"]), 4)
        self.assertEqual(len(payload["benchmark"]["bullets"]), 4)
        self.assertNotIn("test7", all_text)
        self.assertNotIn(".pptx", all_text)
        self.assertFalse(any(re.search(r"\d{3,}", bullet) for bullet in all_bullets))
        self.assertLessEqual(
            sum(1 for bullet in all_bullets if bullet.startswith(("Current", "Benchmark"))),
            1,
        )
        self.assertTrue(any("brand variety" in bullet.lower() or "comparison set" in bullet.lower() for bullet in all_bullets))
        self.assertTrue(any("shelf trust" in bullet.lower() for bullet in payload["benchmark"]["bullets"]))
        representative = payload["debug"]["representative_evidence"]
        self.assertEqual(representative["current"]["valid_capture_count"], 2)
        self.assertEqual(representative["benchmark"]["valid_capture_count"], 2)
        self.assertEqual(payload["current"]["product_count"], 3)
        self.assertEqual(payload["benchmark"]["product_count"], 3)

    def test_slide3_uses_slide6_shared_search_framework_for_theme_selection(self) -> None:
        slide6_payload = build_slide6_visibility(
            [
                {
                    "category": "Beauty/Skin Care",
                    "product_type": "Facial Cleansers",
                    "product_title": "Hydrating Fragrance Free Face Cleanser for Sensitive Skin",
                    "description": "Sensitive skin cleanser with hydrating fragrance free positioning.",
                }
            ],
            [
                {
                    "category": "Beauty/Skin Care",
                    "product_type": "Facial Cleansers",
                    "product_title": "CeraVe Hydrating Facial Cleanser",
                    "description": "Gentle facial cleanser for sensitive skin.",
                }
            ],
            audit_metadata={"client_company_name": "Glow"},
        )
        payload = build_slide3_search_benchmark(
            {
                "current": [
                    {
                        "role": "Current",
                        "sourceRow": 1,
                        "searchTerm": "face cleanser",
                        "screenshotDataUrl": _image_data_url((10, 20, 30)),
                        "orderedMainResultProducts": [
                            {"position": 6, "title": "Glow Face Cleanser Sensitive Skin", "brand": "Glow", "reviewCount": 28},
                        ],
                    }
                ],
                "benchmark": [
                    {
                        "role": "Benchmark",
                        "sourceRow": 2,
                        "searchTerm": "face cleanser",
                        "screenshotDataUrl": _image_data_url((30, 20, 10)),
                        "orderedMainResultProducts": [
                            {"position": 1, "title": "CeraVe Hydrating Facial Cleanser", "brand": "CeraVe", "reviewCount": 1200, "sponsored": True},
                            {"position": 2, "title": "Cetaphil Gentle Skin Cleanser", "brand": "Cetaphil", "reviewCount": 900, "sponsored": True},
                        ],
                    }
                ],
            },
            client_name="Glow",
            slide6_visibility=slide6_payload,
        )
        self.assertEqual(len(payload["current"]["bullets"]), 4)
        self.assertEqual(len(payload["benchmark"]["bullets"]), 4)
        all_text = " ".join(payload["current"]["bullets"] + payload["benchmark"]["bullets"]).lower()
        self.assertIn("face cleanser", all_text)
        self.assertNotIn("product_type_anchor", all_text)
        self.assertNotIn("candidate_bucket", all_text)
        self.assertNotIn("benchmark shelf breadth is limited", all_text)
        self.assertNotIn("shelf breadth is concentrated", all_text)
        self.assertTrue(payload["current"]["debug"]["shared_search_framework_used"])
        self.assertTrue(payload["benchmark"]["debug"]["shared_search_framework_used"])
        self.assertEqual(
            payload["debug"]["shared_search_framework"]["source"],
            "slide6_shared_search_framework",
        )
        current_framework_dims = {
            row["dimension"]
            for row in payload["current"]["bullet_debug"]
            if row["evidence_summary"]["framework_source"]
        }
        self.assertIn("shared_search_query_alignment", current_framework_dims)
        self.assertIn("shared_search_breadth", current_framework_dims)
        current_lead = payload["current"]["bullet_debug"][0]
        benchmark_lead = payload["benchmark"]["bullet_debug"][0]
        self.assertNotEqual(current_lead["bullet_family"], benchmark_lead["bullet_family"])
        self.assertEqual(current_lead["bullet_family"], "side_specific")
        self.assertEqual(benchmark_lead["bullet_family"], "shelf_breadth")
        self.assertNotIn("query fit", payload["current"]["bullets"][0].lower())
        self.assertNotIn("query fit", payload["benchmark"]["bullets"][0].lower())

    def test_slide3_shared_framework_keeps_outlier_query_from_dominating(self) -> None:
        slide6_payload = build_slide6_visibility(
            [
                {
                    "category": "Beauty/Skin Care",
                    "product_type": "Facial Cleansers",
                    "product_title": "Hydrating Face Cleanser",
                }
            ],
            [],
            audit_metadata={"client_company_name": "Glow"},
        )
        payload = build_slide3_search_benchmark(
            {
                "current": [
                    {
                        "role": "Current",
                        "sourceRow": 1,
                        "searchTerm": "moisturizer",
                        "screenshotDataUrl": _image_data_url((1, 2, 3)),
                        "orderedMainResultProducts": [
                            {"position": 9, "title": "Unrelated Moisturizer", "brand": "Other", "reviewCount": 12},
                        ],
                    },
                    {
                        "role": "Current",
                        "sourceRow": 2,
                        "searchTerm": "face cleanser",
                        "screenshotDataUrl": _image_data_url((4, 5, 6)),
                        "orderedMainResultProducts": [
                            {"position": 5, "title": "Glow Hydrating Face Cleanser", "brand": "Glow", "reviewCount": 34},
                        ],
                    },
                    {
                        "role": "Current",
                        "sourceRow": 3,
                        "searchTerm": "face cleanser",
                        "screenshotDataUrl": _image_data_url((7, 8, 9)),
                        "orderedMainResultProducts": [
                            {"position": 7, "title": "Glow Sensitive Face Cleanser", "brand": "Glow", "reviewCount": 42},
                        ],
                    },
                ],
                "benchmark": [],
            },
            client_name="Glow",
            slide6_visibility=slide6_payload,
        )
        self.assertEqual(payload["current"]["search_term"], "face cleanser")
        all_text = " ".join(payload["current"]["bullets"]).lower()
        self.assertIn("face cleanser", all_text)
        self.assertNotIn("moisturizer queries carry", all_text)
        representative = payload["debug"]["representative_evidence"]["current"]
        self.assertEqual(representative["valid_capture_count"], 3)
        self.assertIn("face cleanser", representative["framework_terms_considered"])

    @unittest.skipUnless(TEMPLATE.exists(), "New strategic template is unavailable")
    def test_slide3_generation_populates_screenshots_and_text_without_changing_constants(self) -> None:
        payload = build_slide3_search_benchmark(
            {
                "current": [
                    {
                        "role": "Current",
                        "sourceRow": 3,
                        "searchTerm": "nut butter spread",
                        "screenshotDataUrl": _image_data_url((100, 20, 20)),
                        "orderedMainResultProducts": [
                            {"position": 1, "title": "Client Brand Nut Butter", "brand": "Client Brand", "reviewCount": 120, "sponsored": False},
                        ],
                    }
                ],
                "benchmark": [
                    {
                        "role": "Benchmark",
                        "sourceRow": 4,
                        "searchTerm": "low sugar spreads",
                        "screenshotDataUrl": _image_data_url((20, 20, 100)),
                        "orderedMainResultProducts": [
                            {"position": 1, "title": "Low Sugar Fruit Spread", "brand": "Brand A", "reviewCount": 90, "sponsored": True},
                        ],
                    }
                ],
            },
            client_name="Client Brand",
        )
        source = Presentation(str(TEMPLATE))
        source_slide = _slide3(source)
        source_shape_count = len(source_slide.shapes)
        source_picture_blobs = [
            shape.image.blob
            for shape in sorted(
                [shape for shape in source_slide.shapes if hasattr(shape, "image")],
                key=lambda shape: shape.left,
            )
        ]
        deck_bytes = generate_new_audit_powerpoint_from_template(
            export_plan={
                "audit_metadata": {"client_name": "Client Brand", "client_company_name": "Client Brand"},
                "slide3_search_benchmark": payload,
            },
            template_path=str(TEMPLATE),
            include_slide_9=True,
        )
        slide3 = _slide3(Presentation(io.BytesIO(deck_bytes)))
        texts = _slide_texts(slide3)
        all_text = "\n".join(texts)
        self.assertEqual(len(slide3.shapes), source_shape_count)
        self.assertTrue(any("“nut butter spread”" in text for text in texts))
        self.assertTrue(any("“low sugar spreads”" in text for text in texts))
        self.assertTrue(any("nut butter spreads" in text for text in texts))
        self.assertTrue(any("low sugar spreads" in text for text in texts))
        self.assertNotIn("baby bath products", all_text)
        self.assertNotIn("clean baby care", all_text)
        self.assertNotIn("baby bath and clean baby care categories", all_text)
        for constant in CONSTANT_TEXT:
            self.assertTrue(any(constant in text for text in texts))
        pictures = sorted([shape for shape in slide3.shapes if hasattr(shape, "image")], key=lambda shape: shape.left)
        self.assertEqual(len(pictures), 2)
        self.assertNotEqual(pictures[0].image.blob, source_picture_blobs[0])
        self.assertNotEqual(pictures[1].image.blob, source_picture_blobs[1])
        render_fit = payload["debug"]["render_fit"]
        self.assertEqual(set(render_fit), {"current", "benchmark"})
        self.assertEqual(
            len({item["font_fallback"] for item in render_fit.values()}),
            1,
        )

    @unittest.skipUnless(TEMPLATE.exists(), "New strategic template is unavailable")
    def test_missing_current_preserves_left_side_and_populates_benchmark(self) -> None:
        source = Presentation(str(TEMPLATE))
        source_slide = _slide3(source)
        source_pictures = sorted([shape for shape in source_slide.shapes if hasattr(shape, "image")], key=lambda shape: shape.left)
        payload = build_slide3_search_benchmark(
            {
                "current": [],
                "benchmark": [
                    {
                        "role": "Benchmark",
                        "sourceRow": 4,
                        "searchTerm": "nut butter",
                        "screenshotDataUrl": _image_data_url((20, 20, 100)),
                        "orderedMainResultProducts": [
                            {"position": 1, "title": "Nut Butter", "brand": "Brand A", "reviewCount": 90},
                        ],
                    }
                ],
            },
            client_name="Client Brand",
        )
        generated = Presentation(
            io.BytesIO(
                generate_new_audit_powerpoint_from_template(
                    export_plan={
                        "audit_metadata": {"client_name": "Client Brand", "client_company_name": "Client Brand"},
                        "slide3_search_benchmark": payload,
                    },
                    template_path=str(TEMPLATE),
                    include_slide_9=True,
                )
            )
        )
        pictures = sorted([shape for shape in _slide3(generated).shapes if hasattr(shape, "image")], key=lambda shape: shape.left)
        self.assertEqual(pictures[0].image.blob, source_pictures[0].image.blob)
        self.assertNotEqual(pictures[1].image.blob, source_pictures[1].image.blob)

    @unittest.skipUnless(TEMPLATE.exists(), "New strategic template is unavailable")
    def test_slide3_population_does_not_change_slides_2_4_5_or_6(self) -> None:
        payload = build_slide3_search_benchmark(
            {
                "current": [
                    {
                        "role": "Current",
                        "sourceRow": 3,
                        "searchTerm": "peanut butter spread",
                        "screenshotDataUrl": _image_data_url((100, 20, 20)),
                        "orderedMainResultProducts": [
                            {"position": 1, "title": "Client Brand Peanut Butter", "brand": "Client Brand"},
                        ],
                    }
                ],
                "benchmark": [
                    {
                        "role": "Benchmark",
                        "sourceRow": 4,
                        "searchTerm": "nut butter",
                        "screenshotDataUrl": _image_data_url((20, 20, 100)),
                        "orderedMainResultProducts": [
                            {"position": 1, "title": "Nut Butter", "brand": "Brand A"},
                        ],
                    }
                ],
            },
            client_name="Client Brand",
        )
        base_plan = {"audit_metadata": {"client_name": "Client Brand", "client_company_name": "Client Brand"}}
        without_slide3 = Presentation(
            io.BytesIO(
                generate_new_audit_powerpoint_from_template(
                    export_plan=base_plan,
                    template_path=str(TEMPLATE),
                    include_slide_9=True,
                )
            )
        )
        with_slide3 = Presentation(
            io.BytesIO(
                generate_new_audit_powerpoint_from_template(
                    export_plan={**base_plan, "slide3_search_benchmark": payload},
                    template_path=str(TEMPLATE),
                    include_slide_9=True,
                )
            )
        )
        for slide_number in (2, 4, 5, 6):
            self.assertEqual(
                without_slide3.slides[slide_number - 1].element.xml,
                with_slide3.slides[slide_number - 1].element.xml,
            )


if __name__ == "__main__":
    unittest.main()
