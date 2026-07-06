from __future__ import annotations

import io
import unittest
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

from app.audit_helpers.slide6_visibility import (
    SEGMENT_PACKS,
    VISIBILITY_LABELS,
    _segment,
    _validate_visibility_label,
    _visibility_label,
    build_slide6_visibility,
)
from audit_export import build_audit_export_plan
from audit_powerpoint_new import generate_new_audit_powerpoint_from_template


ROOT = Path(__file__).resolve().parents[1]
NEW_TEMPLATE = ROOT / "templates" / "Audit_Template_New.pptx"


def _record(
    *,
    title: str = "",
    category: str = "",
    product_type: str = "",
    description: str = "",
    key_features: list[str] | None = None,
    ocr_tokens: list[str] | None = None,
) -> dict:
    record = {
        "product_title": title,
        "category": category,
        "product_type": product_type,
        "description": description,
        "key_features": key_features or [],
    }
    if ocr_tokens is not None:
        record["image_analysis"] = {
            "images": [
                {
                    "status": "analyzed",
                    "ocr_tokens": ocr_tokens,
                    "ocr_text": " ".join(ocr_tokens),
                }
            ]
        }
    return record


class Slide6VisibilityTests(unittest.TestCase):
    def test_all_rating_thresholds(self) -> None:
        self.assertEqual(_visibility_label(7.0, 10, 7), "Strong")
        self.assertEqual(_visibility_label(6.999999, 10, 7), "Moderate")
        self.assertEqual(_visibility_label(4.0, 10, 4), "Moderate")
        self.assertEqual(_visibility_label(3.999999, 10, 4), "Partial")
        self.assertEqual(_visibility_label(1.5, 10, 2), "Partial")
        self.assertEqual(_visibility_label(1.499999, 10, 2), "Limited")

    def test_invalid_or_missing_labels_fall_back_to_limited_with_warning(self) -> None:
        warnings: list[str] = []
        self.assertEqual(_validate_visibility_label("Excellent", warnings, "test"), "Limited")
        self.assertEqual(_validate_visibility_label("", warnings, "test"), "Limited")
        self.assertEqual(len(warnings), 2)

    def test_small_sample_does_not_become_strong_from_one_weak_match(self) -> None:
        payload = build_slide6_visibility(
            [_record(title="Organic strawberry jam")],
            [],
            audit_metadata={"client_name": "Test Client"},
        )
        organic = next(item for item in payload["segments"] if item["segment"] in {"organic strawberry", "organic strawberry jam"})
        self.assertEqual(organic["client_visibility"], "Partial")
        self.assertTrue(any("sparse" in warning.lower() for warning in payload["warnings"]))
        self.assertTrue(any("small-sample" in warning.lower() for warning in organic["warnings"]))

    def test_exactly_six_segments_and_jam_pack(self) -> None:
        payload = build_slide6_visibility(
            [_record(title="Strawberry Preserves", category="Jam and Jelly")],
            [],
        )
        self.assertEqual(payload["pack_id"], "jam_fruit_spreads")
        self.assertEqual(len(payload["segments"]), 6)
        self.assertIn("jam", payload["segments"][0]["segment"])
        self.assertTrue(payload["segments"][0]["segment"].islower())

    def test_nut_butter_pack(self) -> None:
        payload = build_slide6_visibility(
            [_record(title="Natural Peanut Butter", product_type="Nut Butter")],
            [],
        )
        self.assertEqual(payload["pack_id"], "nut_butter_spreads")
        self.assertEqual(len(payload["segments"]), 6)

    def test_all_controlled_packs_are_selected(self) -> None:
        cases = (
            ("baby_care", _record(category="Baby Care", product_type="Baby Wipes")),
            ("jam_fruit_spreads", _record(category="Jam and Jelly", title="Fruit Spread")),
            ("nut_butter_spreads", _record(product_type="Peanut Butter")),
            ("skin_care", _record(category="Skin Care", product_type="Face Moisturizer")),
            ("household_cleaning", _record(category="Household Cleaning", product_type="Multi-Surface Cleaner")),
        )
        for expected_pack, record in cases:
            with self.subTest(expected_pack=expected_pack):
                payload = build_slide6_visibility([record], [])
                self.assertEqual(payload["pack_id"], expected_pack)
                self.assertEqual(len(payload["segments"]), 6)

    def test_rows_are_realistic_shopper_queries_not_strategy_labels(self) -> None:
        payload = build_slide6_visibility(
            [
                _record(
                    category="Beauty/Skin Care",
                    product_type="Facial Cleansers",
                    title="Hydrating Face Wash for Sensitive Skin",
                )
            ],
            [],
        )
        rows = [item["segment"] for item in payload["segments"]]
        joined = " ".join(rows).lower()
        self.assertIn("face cleanser", rows)
        self.assertTrue(any("sensitive" in row or "hydrating" in row for row in rows))
        self.assertNotIn("dermatologist recommended", joined)
        self.assertNotIn("daily skin routine", joined)
        self.assertNotIn("shopper education", joined)
        self.assertTrue(all(row == row.lower() for row in rows))

    def test_style_guide_keywords_influence_rows_and_negative_keywords_suppress_leakage(self) -> None:
        payload = build_slide6_visibility(
            [_record(product_type="Nut Butters & Spreads", title="Natural Peanut Butter")],
            [],
        )
        rows = [item["segment"] for item in payload["segments"]]
        joined = " ".join(rows)
        self.assertTrue(any(row in rows for row in ("peanut butter", "natural peanut butter")))
        self.assertTrue(any("hazelnut" in row or "almond" in row or "nut butter" in row for row in rows))
        self.assertNotIn("jam", joined)
        self.assertNotIn("jelly", joined)
        self.assertEqual(payload["debug"]["segment_packs_role"], "fallback_seed_only")
        self.assertTrue(payload["debug"]["row_generation"]["guide_terms_used"]["negative_keywords"])

    def test_final_rows_are_diverse_and_near_duplicates_are_rejected(self) -> None:
        payload = build_slide6_visibility(
            [
                _record(
                    category="Beauty/Skin Care",
                    product_type="Facial Cleansers",
                    title="Gentle Facial Cleanser Face Wash",
                )
            ],
            [],
        )
        rows = [item["segment"] for item in payload["segments"]]
        self.assertEqual(len(rows), 6)
        self.assertEqual(len(set(rows)), 6)
        self.assertFalse({"face cleanser", "face wash"}.issubset(set(rows)))
        self.assertTrue(payload["debug"]["row_selection"]["rejected_similar_rows"])

    def test_jam_and_nut_butter_packs_are_distinct(self) -> None:
        jam = build_slide6_visibility([_record(category="Jam", title="Fruit Spread")], [])
        nut = build_slide6_visibility([_record(product_type="Nut Butter", title="Peanut Butter")], [])
        self.assertNotEqual(
            [item["segment"] for item in jam["segments"]],
            [item["segment"] for item in nut["segments"]],
        )

    def test_duplicate_pack_segments_are_replaced_by_generic_rows(self) -> None:
        duplicate_pack = {
            **SEGMENT_PACKS["baby_care"],
            "segments": (
                _segment("duplicate", "Duplicate", ("baby",)),
                _segment("duplicate", "Duplicate", ("baby",)),
                _segment("unique", "Unique", ("baby",)),
            ),
        }
        with patch.dict(SEGMENT_PACKS, {"baby_care": duplicate_pack}):
            payload = build_slide6_visibility([_record(category="Baby Care")], [])
        names = [item["segment"] for item in payload["segments"]]
        self.assertEqual(len(names), 6)
        self.assertEqual(len({_normalize_for_test(name) for name in names}), 6)
        self.assertTrue(all(name.strip() for name in names))
        self.assertTrue(any("duplicate" in warning.lower() for warning in payload["warnings"]))

    def test_client_and_competitor_are_scored_independently(self) -> None:
        payload = build_slide6_visibility(
            [_record(title="Classic Strawberry Jam", category="Jam")],
            [
                _record(
                    title="Organic Low Sugar Strawberry Fruit Spread for Breakfast Toast",
                    category="Jam",
                    description="Organic fruit spread with reduced sugar for breakfast and toast.",
                ),
                _record(
                    title="Organic Strawberry Preserves",
                    category="Fruit Spread",
                    description="Low sugar breakfast spread for toast.",
                ),
            ],
        )
        rows = {item["segment"]: item for item in payload["segments"]}
        self.assertGreater(
            rows["low sugar jam"]["competitor_supporting_count"],
            rows["low sugar jam"]["client_supporting_count"],
        )
        self.assertNotEqual(
            rows["low sugar jam"]["competitor_visibility"],
            rows["low sugar jam"]["client_visibility"],
        )

    def test_ocr_is_supporting_evidence_only(self) -> None:
        payload = build_slide6_visibility(
            [_record(category="Jam", ocr_tokens=["organic", "fruit", "spread"])],
            [],
        )
        row = next(item for item in payload["segments"] if item["debug"]["ocr_only_support"]["client"])
        self.assertEqual(row["client_supporting_count"], 0)
        self.assertEqual(row["client_visibility"], "Limited")
        self.assertEqual(row["debug"]["ocr_only_support"]["client"], 1)
        self.assertIn("image_analysis.ocr", row["debug"]["matched_fields"]["client"])
        self.assertTrue(any("did not count as support" in warning for warning in row["warnings"]))

    def test_ocr_can_corroborate_but_not_replace_text_evidence(self) -> None:
        payload = build_slide6_visibility(
            [_record(category="Jam", title="Organic option", ocr_tokens=["fruit", "spread"])],
            [],
        )
        organic = next(item for item in payload["segments"] if item["segment"] == "organic fruit spread")
        self.assertEqual(organic["client_supporting_count"], 1)
        self.assertEqual(organic["client_visibility"], "Partial")
        self.assertIn("product_title", organic["debug"]["matched_fields"]["client"])
        self.assertIn("image_analysis.ocr", organic["debug"]["matched_fields"]["client"])

    def test_missing_records_are_not_counted_in_segment_denominator(self) -> None:
        payload = build_slide6_visibility(
            [_record(category="Jam"), {}, {"description": ""}],
            [],
        )
        row = next(item for item in payload["segments"] if item["segment"] in {"jam", "jelly"} or "jam" in item["segment"])
        self.assertEqual(row["client_analyzed_count"], 1)
        self.assertEqual(row["client_fraction"], "1/1")

    def test_traceability_fields_are_complete_and_labels_are_controlled(self) -> None:
        payload = build_slide6_visibility(
            [_record(category="Jam", title="Strawberry Jam")],
            [_record(category="Fruit Spread", description="Breakfast spread")],
        )
        for row in payload["segments"]:
            self.assertIn(row["client_visibility"], VISIBILITY_LABELS)
            self.assertIn(row["competitor_visibility"], VISIBILITY_LABELS)
            self.assertTrue(row["debug"]["segment_id"])
            self.assertEqual(row["debug"]["selected_pack"], payload["pack_id"])
            self.assertIn("client_fraction", row)
            self.assertIn("competitor_fraction", row)
            self.assertIn("client_percentage", row)
            self.assertIn("competitor_percentage", row)
            self.assertIn("matched_fields", row["debug"])
            self.assertIn("matched_terms", row["debug"])
            self.assertIn("row_selection_reason", row["debug"])
            self.assertIn("candidate_source", row["debug"])
            self.assertIn("ranking_factors", row["debug"])
            self.assertIn("score_inputs", row["debug"])
            self.assertIn("label_reason", row["debug"])
            self.assertTrue(row["debug"]["label_reason"]["client"])
            self.assertTrue(row["debug"]["label_reason"]["competitor"])
            self.assertTrue(row["debug"]["reason"])
            self.assertIn("warnings", row)
        self.assertIn("row_generation", payload["debug"])
        self.assertIn("row_selection", payload["debug"])

    def test_generic_fallback_is_restrained_and_non_repetitive(self) -> None:
        payload = build_slide6_visibility(
            [
                _record(
                    title="Premium Widget",
                    category="Specialty Supplies",
                    product_type="Widget Kit",
                    key_features=["durable", "portable", "family"],
                )
            ],
            [],
        )
        names = [item["segment"] for item in payload["segments"]]
        self.assertEqual(payload["pack_id"], "generic")
        self.assertEqual(len(names), 6)
        self.assertEqual(len(set(names)), 6)

    def test_missing_data_warnings_and_client_name(self) -> None:
        payload = build_slide6_visibility(
            [],
            [],
            audit_metadata={"client_company_name": "Example Company"},
        )
        self.assertEqual(payload["client_label"], "Example Company")
        self.assertEqual(len(payload["segments"]), 6)
        self.assertGreaterEqual(len(payload["warnings"]), 3)
        self.assertTrue(
            all(
                item["client_visibility"] == "Limited"
                and item["competitor_visibility"] == "Limited"
                for item in payload["segments"]
            )
        )

    def test_export_plan_includes_slide6_without_restructuring_existing_keys(self) -> None:
        plan = build_audit_export_plan(
            audit_record={
                "client_name": "Short Name",
                "client_company_name": "Full Client Company",
            },
            primary_entries=[],
            competitor_assignments=[],
            competitor_records=[],
        )
        self.assertIn("slide2_summary", plan)
        self.assertIn("slide4_findings", plan)
        self.assertIn("slide6_visibility", plan)
        self.assertEqual(plan["slide6_visibility"]["client_label"], "Full Client Company")
        self.assertEqual(len(plan["slide6_visibility"]["segments"]), 6)

    @unittest.skipUnless(NEW_TEMPLATE.exists(), "New strategic template is unavailable")
    def test_powerpoint_replaces_slide6_in_place(self) -> None:
        client_name = "Extremely Long Client Company Name Incorporated"
        payload = build_slide6_visibility(
            [_record(title="Strawberry Jam", category="Jam")],
            [_record(title="Organic Low Sugar Fruit Spread", category="Jam")],
            audit_metadata={"client_company_name": client_name},
        )
        source = Presentation(str(NEW_TEMPLATE))
        source_slide = source.slides[5]
        source_shape_count = len(source_slide.shapes)
        source_textbox_count = sum(
            1 for shape in source_slide.shapes if getattr(shape, "has_text_frame", False)
        )
        source_table_shape = next(shape for shape in source_slide.shapes if getattr(shape, "has_table", False))
        source_table = source_table_shape.table
        source_geometry = (
            source_table_shape.left,
            source_table_shape.top,
            source_table_shape.width,
            source_table_shape.height,
        )
        source_cell_format = [
            [
                (
                    cell.margin_left,
                    cell.margin_right,
                    cell.margin_top,
                    cell.margin_bottom,
                    cell.vertical_anchor,
                    cell._tc.get_or_add_tcPr().xml,
                )
                for cell in row.cells
            ]
            for row in source_table.rows
        ]
        source_text_format = [
            [
                (
                    cell.text_frame.paragraphs[0].alignment,
                    cell.text_frame.paragraphs[0].runs[0].font.name,
                    cell.text_frame.paragraphs[0].runs[0].font.bold,
                    cell.text_frame.paragraphs[0].runs[0].font.italic,
                    cell.text_frame.paragraphs[0].runs[0].font.color.type,
                )
                for cell in row.cells
            ]
            for row in source_table.rows
        ]

        deck_bytes = generate_new_audit_powerpoint_from_template(
            export_plan={
                "audit_metadata": {
                    "client_name": client_name,
                    "client_company_name": client_name,
                },
                "slide6_visibility": payload,
            },
            template_path=str(NEW_TEMPLATE),
            include_slide_9=True,
        )
        generated = Presentation(io.BytesIO(deck_bytes))
        slide = generated.slides[5]
        self.assertEqual(len(slide.shapes), source_shape_count)
        self.assertEqual(
            sum(1 for shape in slide.shapes if getattr(shape, "has_text_frame", False)),
            source_textbox_count,
        )
        table = next(shape.table for shape in slide.shapes if getattr(shape, "has_table", False))
        generated_table_shape = next(shape for shape in slide.shapes if getattr(shape, "has_table", False))
        self.assertEqual(
            (
                generated_table_shape.left,
                generated_table_shape.top,
                generated_table_shape.width,
                generated_table_shape.height,
            ),
            source_geometry,
        )
        self.assertEqual(table.cell(0, 2).text.strip(), client_name)
        self.assertTrue(table.cell(0, 2).text_frame.word_wrap)
        self.assertEqual(
            table.cell(0, 2).text_frame.auto_size,
            MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
        )
        self.assertTrue(
            all(
                run.font.size is not None
                for paragraph in table.cell(0, 2).text_frame.paragraphs
                for run in paragraph.runs
            )
        )
        for index, segment in enumerate(payload["segments"], start=1):
            self.assertEqual(table.cell(index, 0).text.strip(), segment["segment"])
            self.assertEqual(table.cell(index, 1).text.strip(), segment["competitor_visibility"])
            self.assertEqual(table.cell(index, 2).text.strip(), segment["client_visibility"])
        for row_index, row in enumerate(table.rows):
            for column_index, cell in enumerate(row.cells):
                original = source_cell_format[row_index][column_index]
                self.assertEqual(cell.margin_left, original[0])
                self.assertEqual(cell.margin_right, original[1])
                self.assertEqual(cell.margin_top, original[2])
                self.assertEqual(cell.margin_bottom, original[3])
                self.assertEqual(cell.vertical_anchor, original[4])
                self.assertEqual(cell._tc.get_or_add_tcPr().xml, original[5])
                generated_run = cell.text_frame.paragraphs[0].runs[0]
                text_format = source_text_format[row_index][column_index]
                self.assertEqual(cell.text_frame.paragraphs[0].alignment, text_format[0])
                self.assertEqual(generated_run.font.name, text_format[1])
                self.assertEqual(generated_run.font.bold, text_format[2])
                self.assertEqual(generated_run.font.italic, text_format[3])
                self.assertEqual(generated_run.font.color.type, text_format[4])
        all_text = " ".join(
            shape.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        self.assertIn(payload["intro"], all_text)
        self.assertIn(payload["takeaway"], all_text)
        self.assertNotIn("currently own more Walmart search paths", all_text)
        self.assertNotIn("Baby Care", [table.cell(row, 0).text.strip() for row in range(1, 7)])
        self.assertNotIn("Diapers", [table.cell(row, 0).text.strip() for row in range(1, 7)])

    @unittest.skipUnless(NEW_TEMPLATE.exists(), "New strategic template is unavailable")
    def test_powerpoint_invalid_labels_fall_back_without_new_shapes(self) -> None:
        payload = build_slide6_visibility([_record(category="Jam")], [])
        payload["segments"][0]["client_visibility"] = ""
        payload["segments"][0]["competitor_visibility"] = "Excellent"
        source = Presentation(str(NEW_TEMPLATE))
        shape_count = len(source.slides[5].shapes)
        deck_bytes = generate_new_audit_powerpoint_from_template(
            export_plan={
                "audit_metadata": {"client_name": "Test"},
                "slide6_visibility": payload,
            },
            template_path=str(NEW_TEMPLATE),
            include_slide_9=True,
        )
        generated = Presentation(io.BytesIO(deck_bytes))
        slide = generated.slides[5]
        table = next(shape.table for shape in slide.shapes if getattr(shape, "has_table", False))
        self.assertEqual(table.cell(1, 1).text.strip(), "Limited")
        self.assertEqual(table.cell(1, 2).text.strip(), "Limited")
        self.assertEqual(len(slide.shapes), shape_count)
        self.assertGreaterEqual(len(payload["warnings"]), 2)


def _normalize_for_test(value: str) -> str:
    return " ".join(value.lower().split())


if __name__ == "__main__":
    unittest.main()
