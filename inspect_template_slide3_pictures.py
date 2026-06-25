from pptx import Presentation
from pathlib import Path

prs = Presentation(Path('templates/Audit_Template_New.pptx'))
slide = prs.slides[2]
for idx, shape in enumerate(slide.shapes):
    print(idx, shape.shape_type, type(shape).__name__, getattr(shape, 'name', ''), getattr(shape, 'left', None), getattr(shape, 'top', None), getattr(shape, 'width', None), getattr(shape, 'height', None))
    if hasattr(shape, 'image'):
        print('  has image')
    if getattr(shape, 'has_text_frame', False):
        print('  text=', repr(shape.text[:200]))
    if hasattr(shape, 'shapes'):
        print(' group children', len(shape.shapes))
