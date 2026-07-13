from __future__ import annotations

import base64
import io
import math
import re
import ssl
from collections import Counter
from typing import Any
from urllib.request import Request, urlopen

from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from app.audit_helpers.image_guides import (
    get_image_guide_page,
    load_image_guide,
    resolve_image_guide_category,
)
from app.audit_helpers.bullet_uniqueness import normalize_bullet_text
from app.audit_helpers.slide4_findings import build_slide4_group_findings
from app.audit_helpers.audit_language_resolver import strategic_bullet_text
from app.audit_helpers.strategic_cue_engine import aggregate_strategic_cues
from app.audit_helpers.strategic_identity import resolve_strategic_identity
from audit_powerpoint import _format_cover_date


PLACEHOLDER_RE = re.compile(r"\{\{[^{}]+\}\}")
ONLY_PLACEHOLDERS_RE = re.compile(r"^(?:\s*\{\{[^{}]+\}\}\s*)+$")
SLIDE4_SLOT_LABELS = {
    "silo_front": "Front Silo",
    "graphic_ingredients": "Ingredients",
    "silo_alt_in_pack": "Alt In-Pack",
    "lifestyle_in_use": "Lifestyle",
    "graphic_guarantee": "Guarantee",
    "feature_graphic": "Features",
    "dimensions": "Dimensions",
    "graphic_nutrition": "Nutrition",
}
SLIDE4_BULLET_FAMILIES = (
    "positioning_title",
    "detail_compliance",
    "education_storytelling",
    "trust_visual",
)


def _safe_text(value: Any) -> str:
    return str(value or "").strip()


def _shape_text(shape: Any) -> str:
    """Get text from a shape, handling text frames."""
    if not getattr(shape, "has_text_frame", False):
        return ""
    return _safe_text(getattr(shape, "text", ""))


def _find_shape_contains(slide: Any, token: str) -> Any:
    """Find a shape by searching for token text (case-insensitive)."""
    token_l = token.lower()
    for shape in slide.shapes:
        text = _shape_text(shape).lower()
        if token_l in text:
            return shape
    return None


def _replace_shape_text_preserve_style(shape: Any, text: str) -> None:
    """Replace text in a shape while preserving existing run styling."""
    if not shape or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        shape.text = _safe_text(text)
        return
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = _safe_text(text)
        for run in first_para.runs[1:]:
            run.text = ""
        for para in tf.paragraphs[1:]:
            for run in para.runs:
                run.text = ""
        return
    shape.text = _safe_text(text)


def _replace_paragraph_text_preserve_style(paragraph: Any, text: str) -> None:
    """Replace text in a paragraph while preserving styling."""
    if paragraph.runs:
        paragraph.runs[0].text = _safe_text(text)
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = _safe_text(text)


def _replace_text_in_paragraph(paragraph: Any, replacements: dict[str, str]) -> dict[str, int]:
    """Replace template text within a paragraph without creating new shapes."""
    counts = {target: 0 for target in replacements}
    if not paragraph:
        return counts

    for run in getattr(paragraph, "runs", []) or []:
        run_text = run.text or ""
        if not run_text:
            continue
        new_text = run_text
        for target, replacement in replacements.items():
            found = new_text.count(target)
            if found:
                counts[target] += found
                new_text = new_text.replace(target, replacement)
        if new_text != run_text:
            run.text = new_text

    paragraph_text = paragraph.text or ""
    if not any(target in paragraph_text for target in replacements):
        return counts

    new_text = paragraph_text
    for target, replacement in replacements.items():
        found = new_text.count(target)
        if found:
            counts[target] += found
            new_text = new_text.replace(target, replacement)

    _replace_paragraph_text_preserve_style(paragraph, new_text)
    return counts


def _merge_counts(total: dict[str, int], increment: dict[str, int]) -> None:
    for target, count in increment.items():
        total[target] = total.get(target, 0) + int(count or 0)


def _walk_shapes(shapes: Any) -> list[Any]:
    found = []
    for shape in shapes:
        found.append(shape)
        if hasattr(shape, "shapes"):
            found.extend(_walk_shapes(shape.shapes))
    return found


def _find_slide_by_title(prs: Any, title_text: str) -> Any:
    """Find a slide by exact title text first, then fall back to containment."""
    title_l = title_text.strip().lower()
    for slide in prs.slides:
        for shape in _walk_shapes(slide.shapes):
            if _shape_text(shape).strip().lower() == title_l:
                return slide
    for slide in prs.slides:
        for shape in _walk_shapes(slide.shapes):
            text = _shape_text(shape).lower()
            if title_l in text:
                return slide
    return None


def _remove_slide_by_title(prs: Any, title_text: str) -> bool:
    """
    Remove a slide by its title text.
    Returns True if a slide was removed, False otherwise.
    """
    title_l = title_text.lower()
    slide_idx = None
    for idx, slide in enumerate(prs.slides):
        for shape in _walk_shapes(slide.shapes):
            text = _shape_text(shape).lower()
            if title_l in text:
                slide_idx = idx
                break
        if slide_idx is not None:
            break
    
    if slide_idx is not None:
        # Remove the slide via the slide layout's id
        rId = prs.slides._sldIdLst[slide_idx].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[slide_idx]
        return True
    return False


def replace_existing_template_text(prs: Any, replacements: dict[str, str]) -> dict[str, int]:
    """
    Replace existing template strings in-place across slides, groups, and tables.
    Returns a replacement count by target string for debug reporting.
    """
    counts = {target: 0 for target in replacements}
    for slide in prs.slides:
        for shape in _walk_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    _merge_counts(counts, _replace_text_in_paragraph(paragraph, replacements))
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            _merge_counts(counts, _replace_text_in_paragraph(paragraph, replacements))
    for target in replacements:
        if counts.get(target, 0):
            print(f"[audit_powerpoint_new] Replaced '{target}' {counts[target]} time(s).")
        else:
            print(f"[audit_powerpoint_new] Template text '{target}' was not found.")
    return counts


def _contains_unresolved_placeholder(text: str) -> bool:
    return bool(PLACEHOLDER_RE.search(text or ""))


def _only_unresolved_placeholders(text: str) -> bool:
    return bool(ONLY_PLACEHOLDERS_RE.fullmatch(text or ""))


def _remove_shape(shape: Any) -> bool:
    element = getattr(shape, "element", None)
    parent = element.getparent() if element is not None else None
    if parent is None:
        return False
    parent.remove(element)
    return True


def _clear_unresolved_placeholders_in_paragraph(paragraph: Any) -> int:
    text = paragraph.text or ""
    if not _contains_unresolved_placeholder(text):
        return 0

    if _only_unresolved_placeholders(text):
        _replace_paragraph_text_preserve_style(paragraph, "")
        return len(PLACEHOLDER_RE.findall(text))

    count = 0
    for run in getattr(paragraph, "runs", []) or []:
        run_text = run.text or ""
        if not _contains_unresolved_placeholder(run_text):
            continue
        if _only_unresolved_placeholders(run_text):
            count += len(PLACEHOLDER_RE.findall(run_text))
            run.text = ""
        else:
            cleaned = PLACEHOLDER_RE.sub("", run_text)
            count += len(PLACEHOLDER_RE.findall(run_text))
            run.text = cleaned

    # Fallback for placeholders split across runs.
    remaining = paragraph.text or ""
    if _contains_unresolved_placeholder(remaining):
        count += len(PLACEHOLDER_RE.findall(remaining))
        _replace_paragraph_text_preserve_style(paragraph, PLACEHOLDER_RE.sub("", remaining))
    return count


def clear_unresolved_placeholder_text(prs: Any) -> dict[str, int]:
    """Remove placeholder-only shapes and clear any remaining {{...}} text."""
    counts = {"removed_shapes": 0, "cleared_placeholders": 0}
    for slide in prs.slides:
        for shape in list(_walk_shapes(slide.shapes)):
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            counts["cleared_placeholders"] += _clear_unresolved_placeholders_in_paragraph(paragraph)

            if not getattr(shape, "has_text_frame", False):
                continue

            shape_text = getattr(shape, "text", "") or ""
            if _only_unresolved_placeholders(shape_text):
                if _remove_shape(shape):
                    counts["removed_shapes"] += 1
                continue

            for paragraph in shape.text_frame.paragraphs:
                counts["cleared_placeholders"] += _clear_unresolved_placeholders_in_paragraph(paragraph)

    print(
        "[audit_powerpoint_new] Removed "
        f"{counts['removed_shapes']} unresolved placeholder shape(s); cleared "
        f"{counts['cleared_placeholders']} unresolved placeholder token(s)."
    )
    return counts


def _find_pictures_in_region(slide: Any, region_left: float, region_top: float, region_right: float, region_bottom: float) -> list[Any]:
    """Find all picture shapes within a region."""
    pictures = []
    for shape in slide.shapes:
        if hasattr(shape, "image") or "picture" in str(shape.shape_type).lower():
            if hasattr(shape, "left") and hasattr(shape, "top"):
                shape_left = shape.left.inches if hasattr(shape.left, "inches") else shape.left
                shape_top = shape.top.inches if hasattr(shape.top, "inches") else shape.top
                if (shape_left >= region_left and shape_left <= region_right and
                    shape_top >= region_top and shape_top <= region_bottom):
                    pictures.append(shape)
    return pictures


def _load_image_from_url(url: str) -> bytes | None:
    """Download image from URL and return bytes. Returns None on failure."""
    try:
        request = Request(
            url,
            headers={
                "User-Agent": (
                    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                    "AppleWebKit/537.36 (KHTML, like Gecko) "
                    "Chrome/126.0.0.0 Safari/537.36"
                )
            },
        )
        try:
            with urlopen(request, timeout=10) as response:
                return response.read()
        except Exception as first_exc:
            if not str(url).lower().startswith("https://"):
                raise
            context = ssl._create_unverified_context()
            with urlopen(request, timeout=10, context=context) as response:
                print(
                    "[audit_powerpoint_new] Slide 4 image download used SSL "
                    f"fallback for {url} ({first_exc})"
                )
                return response.read()
    except Exception as exc:
        print(f"[audit_powerpoint_new] Slide 4 image download failed: {url} ({exc})")
        return None


def _add_contained_picture(
    slide: Any,
    *,
    left: int,
    top: int,
    width: int,
    height: int,
    image_bytes: bytes,
) -> Any:
    """Add an uncropped picture contained and centered within an EMU cell."""
    try:
        with Image.open(io.BytesIO(image_bytes)) as image:
            image_width, image_height = image.size
        if image_width <= 0 or image_height <= 0:
            return None
        scale = min(width / image_width, height / image_height)
        rendered_width = max(1, int(image_width * scale))
        rendered_height = max(1, int(image_height * scale))
        rendered_left = left + (width - rendered_width) // 2
        rendered_top = top + (height - rendered_height) // 2
        return slide.shapes.add_picture(
            io.BytesIO(image_bytes),
            left=rendered_left,
            top=rendered_top,
            width=rendered_width,
            height=rendered_height,
        )
    except Exception as exc:
        print(f"[audit_powerpoint_new] Image placement failed: {exc}")
        return None


def _normalize_ordered_images(record: dict[str, Any]) -> list[dict[str, Any]]:
    images = record.get("ordered_images")
    if not isinstance(images, list):
        images = record.get("images", []) or []
    normalized: list[dict[str, Any]] = []
    for position, image in enumerate(images):
        if not isinstance(image, dict):
            continue
        normalized.append(
            {
                "index": image.get("index", position),
                "url": _safe_text(image.get("url", "")),
                "is_hero": bool(image.get("is_hero", position == 0)),
                "width": image.get("width"),
                "height": image.get("height"),
                "dimensions": image.get("dimensions", ""),
                "dimensions_text": image.get("dimensions_text", ""),
            }
        )
    return normalized[:12]


def _build_image_guide_match(category: str, product_type: str) -> dict[str, Any]:
    category_key = resolve_image_guide_category(category)
    page = get_image_guide_page(category_key, product_type)
    if not page:
        return {
            "matched": False,
            "category_key": category_key,
            "product_type": product_type,
        }
    guide = load_image_guide(category_key)
    slot_definitions = guide.get("slot_definitions", {}) or {}
    required_slots = list(page.get("required_slots", []) or [])
    return {
        "matched": True,
        "category_key": category_key,
        "product_type": product_type,
        "page_key": page.get("page_key"),
        "page_display_name": page.get("display_name", ""),
        "required_slots": required_slots,
        "required_slot_labels": [
            SLIDE4_SLOT_LABELS.get(slot)
            or _safe_text((slot_definitions.get(slot, {}) or {}).get("label"))
            or str(slot).replace("_", " ").title()
            for slot in required_slots
        ],
        "additional_recommendations": list(page.get("additional_recommendations", []) or []),
    }


def _parse_image_dimensions(image: dict[str, Any]) -> tuple[int, int] | None:
    try:
        width = int(float(str(image.get("width", "")).replace(",", "").strip()))
        height = int(float(str(image.get("height", "")).replace(",", "").strip()))
        if width > 0 and height > 0:
            return width, height
    except (TypeError, ValueError):
        pass
    raw = _safe_text(image.get("dimensions_text") or image.get("dimensions"))
    match = re.search(r"(\d{2,5})\s*(?:x|×)\s*(\d{2,5})", raw, flags=re.IGNORECASE)
    if match:
        return int(match.group(1)), int(match.group(2))
    return None


def _first_value(record: dict[str, Any], *keys: str, default: Any = "") -> Any:
    for key in keys:
        current: Any = record
        found = True
        for part in key.split("."):
            if not isinstance(current, dict) or part not in current:
                found = False
                break
            current = current[part]
        if found and current not in (None, "", [], {}):
            return current
    return default


def _as_text_list(value: Any) -> list[str]:
    if isinstance(value, list):
        items = value
    elif value in (None, "", {}, []):
        items = []
    else:
        items = [value]
    texts: list[str] = []
    for item in items:
        if isinstance(item, dict):
            text = _safe_text(
                _first_value(item, "text", "title", "description", "value", "name")
            )
        else:
            text = _safe_text(item)
        if text:
            texts.append(text)
    return list(dict.fromkeys(texts))


def _truthy_evidence(value: Any) -> bool:
    if isinstance(value, bool):
        return value
    text = _safe_text(value).lower()
    return text in {"true", "yes", "y", "1", "present", "available", "included"}


def _number_value(value: Any) -> float:
    try:
        return float(str(value).replace(",", "").strip())
    except (TypeError, ValueError):
        return 0.0


def _slide4_record_facts(record: dict[str, Any]) -> dict[str, Any]:
    images = _normalize_ordered_images(record)
    title = _safe_text(_first_value(record, "product_title", "Product Title", "title", "productTitle"))
    brand = _safe_text(_first_value(record, "brand", "Brand", "brandName"))
    category = _safe_text(_first_value(record, "category", "Category", "categoryPathName"))
    product_type = _safe_text(_first_value(record, "product_type", "subcategory", "Product Type", "productType"))
    description = _safe_text(_first_value(record, "description_body", "Description Body", "descriptionBody"))
    description_bullets = _as_text_list(
        _first_value(record, "description_bullets", "Description Bullets", "descriptionBullets", default=[])
    )
    key_features = _as_text_list(
        _first_value(record, "key_features", "Key Features", "keyFeatures", default=[])
    )
    text_blob = " ".join(
        [title, brand, category, product_type, description, *description_bullets, *key_features]
    ).lower()
    image_text = " ".join(
        _safe_text(image.get("url"))
        for image in images
        if isinstance(image, dict)
    ).lower()
    review_count = int(
        _number_value(_first_value(record, "review_count", "Review Count", "reviews_summary.review_count", "reviews_summary.count"))
    )
    average_rating = _number_value(_first_value(record, "average_rating", "Average Rating", "reviews_summary.average_rating", "reviews_summary.rating"))
    return {
        "title": title,
        "brand": brand,
        "category": category,
        "product_type": product_type,
        "description": description,
        "description_bullets": description_bullets,
        "key_features": key_features,
        "text_blob": text_blob,
        "image_text": image_text,
        "images": images,
        "image_count": int(_number_value(_first_value(record, "image_count", "Image Count", default=len(images))) or len(images)),
        "review_count": review_count,
        "average_rating": average_rating,
        "seller_name": _safe_text(_first_value(record, "seller_name", "Seller Name")),
        "sold_by_walmart": _truthy_evidence(_first_value(record, "sold_by_walmart", "Sold by Walmart")),
        "shipped_by_walmart": _truthy_evidence(_first_value(record, "shipped_by_walmart", "Shipped by Walmart")),
        "ebc_present": _truthy_evidence(
            _first_value(
                record,
                "enhanced_brand_content",
                "enhanced_brand_content_present",
                "Enhanced Brand Content",
                "enhancedBrandContentPresent",
            )
        ),
    }


def _has_any(text: str, *terms: str) -> bool:
    return any(term in text for term in terms)


def _append_unique_bullet(
    bullets: list[str],
    debug: list[dict[str, Any]],
    *,
    text: str,
    bullet_type: str,
    dimension: str,
    signals: list[str],
    reason: str,
    used: set[str],
) -> None:
    clean = _safe_text(text)
    key = normalize_bullet_text(clean)
    if not clean or key in used:
        return
    used.add(key)
    bullets.append(clean)
    debug.append(
        {
            "text": clean,
            "type": bullet_type,
            "dimension": dimension,
            "signals": signals,
            "reason": reason,
        }
    )


def _slide4_content_blob(facts: dict[str, Any]) -> str:
    return " ".join(
        [
            _safe_text(facts.get("title")),
            _safe_text(facts.get("brand")),
            _safe_text(facts.get("product_type")),
            _safe_text(facts.get("description")),
            *[_safe_text(value) for value in facts.get("description_bullets", []) or []],
            *[_safe_text(value) for value in facts.get("key_features", []) or []],
        ]
    ).lower()


SLIDE4_SURFACE_TERM_BLOCKLIST = {
    "audience and ingredient needs",
    "audience ingredient needs",
    "ear acupressure seed",
    "ear acupressure seeds",
    "ear vaccaria seed",
    "taxonomy path",
    "vaccaria seed",
    "vaccaria seeds",
}
SLIDE4_SURFACE_BLOCKED_TOKENS = {
    "acupressure",
    "audience",
    "recommended",
    "taxonomy",
    "vaccaria",
}
SLIDE4_METADATA_RE = re.compile(
    r"\b\d{3,5}\s*x\s*\d{3,5}\b|\b[\w-]+\.(?:jpg|jpeg|png|webp)\b|\b(?:jpg|jpeg|png|webp)\b",
    re.I,
)
SLIDE4_INTERNAL_RE = re.compile(
    r"\b(?:framework|resolution path|signal bucket|taxonomy path|inferred category node)\b",
    re.I,
)


def _slide4_sanitize_surface_text(text: Any, fallback: str = "PDP content can work harder") -> str:
    clean = re.sub(r"\s+", " ", _safe_text(text)).strip()
    clean = SLIDE4_METADATA_RE.sub("", clean)
    clean = SLIDE4_INTERNAL_RE.sub("", clean)
    clean = re.sub(r"\s+", " ", clean).strip(" .;:-")
    return clean or fallback


def _slide4_sane_surface_phrase(term: Any, evidence: str) -> str:
    cleaned = re.sub(r"\s+", " ", _safe_text(term).replace("&", " and ")).strip().lower()
    if not cleaned or cleaned in {"category", "product", "product type"}:
        return ""
    if SLIDE4_METADATA_RE.search(cleaned) or SLIDE4_INTERNAL_RE.search(cleaned):
        return ""
    normalized_tokens = [token for token in re.split(r"[^a-z0-9]+", cleaned) if token]
    if cleaned in SLIDE4_SURFACE_TERM_BLOCKLIST:
        return ""
    if any(token in SLIDE4_SURFACE_BLOCKED_TOKENS for token in normalized_tokens):
        return ""
    if cleaned in {"dermatologist recommended", "doctor recommended", "clinically tested"}:
        return ""
    if len(normalized_tokens) > 5:
        return ""
    if " and " in f" {cleaned} " and not any(
        phrase in cleaned
        for phrase in ("jams jellies", "nut butters", "fragrance free", "normal to oily")
    ):
        return ""
    if len(normalized_tokens) >= 3 and any(
        token in {"alignment", "coverage", "discoverability", "framework", "internal", "merchandising", "visibility"}
        for token in normalized_tokens
    ):
        return ""
    evidence_tokens = set(token for token in re.split(r"[^a-z0-9]+", evidence.lower()) if len(token) > 3)
    phrase_tokens = [token for token in normalized_tokens if len(token) > 3]
    if evidence_tokens and phrase_tokens and not any(token in evidence_tokens for token in phrase_tokens):
        return ""
    return cleaned


def _slide4_safe_fallback_phrase(facts: dict[str, Any]) -> str:
    blob = _slide4_content_blob(facts)
    if any(term in blob for term in ("antacid", "heartburn", "acid reducer", "upset stomach", "stomach")):
        if "heartburn" in blob:
            return "heartburn relief"
        if "acid" in blob and "reducer" in blob:
            return "acid reducer"
        if "stomach" in blob:
            return "stomach relief"
        return "antacid"
    if any(term in blob for term in ("face cleanser", "facial cleanser", "face wash", "skin cleanser")):
        return "facial cleanser"
    if any(term in blob for term in ("hazelnut", "cocoa", "nutella")):
        return "hazelnut-and-cocoa spread"
    if any(term in blob for term in ("peanut butter", "jif", "fresh roasted", "fresh-roasted")):
        return "peanut butter"
    if "almond butter" in blob:
        return "almond butter"
    return "product role"


def _slide4_product_phrase(facts: dict[str, Any]) -> str:
    content_blob = _slide4_content_blob(facts)
    if _has_any(content_blob, "hazelnut", "cocoa", "nutella"):
        return "hazelnut-and-cocoa spread"
    if _has_any(content_blob, "peanut butter", "jif", "fresh roasted", "fresh-roasted"):
        return "peanut butter"
    if _has_any(content_blob, "almond butter"):
        return "almond butter"
    product_type = _safe_text(facts.get("product_type")).lower()
    if product_type:
        safe_product_type = _slide4_sane_surface_phrase(product_type.replace("&", "and"), content_blob)
        if safe_product_type:
            return safe_product_type
    category = _safe_text(facts.get("category")).split("/")[-1].lower()
    safe_category = _slide4_sane_surface_phrase(category.replace("&", "and"), content_blob)
    return safe_category or _slide4_safe_fallback_phrase(facts)


def _slide4_content_phrases(facts: dict[str, Any]) -> dict[str, str]:
    product_phrase = _slide4_product_phrase(facts)
    category = _safe_text(facts.get("category")).split("/")[-1].replace("&", "and").strip().lower()
    category = _slide4_sane_surface_phrase(category, _slide4_content_blob(facts)) or product_phrase
    return {
        "product": product_phrase,
        "positioning": "product role",
        "discovery": "search discovery",
        "education": "usage guidance",
        "content": "PDP content",
        "category_context": f"{category} category" if category != product_phrase else f"{product_phrase} space",
    }


def _slide4_evidence_theme(facts: dict[str, Any]) -> str:
    blob = _slide4_content_blob(facts)
    if _has_any(blob, "hazelnut", "cocoa", "nutella"):
        return "hazelnut_cocoa"
    if _has_any(blob, "peanut butter", "peanut", "protein", "jif", "fresh roasted", "fresh-roasted"):
        return "peanut_protein"
    if _has_any(blob, "almond butter", "almond"):
        return "almond"
    return "product"


def _slide4_theme_terms(facts: dict[str, Any]) -> dict[str, str]:
    theme = _slide4_evidence_theme(facts)
    image_count = int(facts.get("image_count", 0) or 0)
    if theme == "hazelnut_cocoa":
        terms = {
            "theme": "hazelnut-cocoa spread",
            "positioning": "Hazelnut-cocoa positioning supports breakfast occasions",
            "benefit": "Cocoa and hazelnut cues sharpen flavor appeal",
            "detail": "Clear hazelnut and cocoa flavor communication",
            "story": "Breakfast and snack storytelling builds usage relevance",
        }
    elif theme == "peanut_protein":
        terms = {
            "theme": "peanut butter",
            "positioning": "Peanut butter positioning reinforces pantry relevance",
            "benefit": "Protein cues strengthen snack and pantry relevance",
            "detail": "Clear peanut and protein detail",
            "story": "Pantry and recipe cues support comparison",
        }
    elif theme == "almond":
        terms = {
            "theme": "almond butter",
            "positioning": "Almond butter positioning clarifies variant choice",
            "benefit": "Ingredient cues strengthen premium shelf relevance",
            "detail": "Clear almond butter pack and nutrition detail",
            "story": "Snack and recipe storytelling supports shopper comparison",
        }
    else:
        product = _slide4_product_phrase(facts)
        terms = {
            "theme": product,
            "positioning": "Title clarifies product role",
            "benefit": "Benefit communication is clearer",
            "detail": "Feature detail supports comparison",
            "story": "Image stack extends usage guidance",
        }
    terms["visual"] = (
        f"{image_count}-image carousel supports visual education"
        if image_count >= 6
        else "Cohesive PDP visual education"
    )
    return terms


def _slide4_theme_bullet(cue_key: str, classification: str, facts: dict[str, Any], identity: dict[str, Any]) -> str:
    if cue_key == "pack_or_spec_detail":
        return "Clear pack and nutrition detail"
    if cue_key == "review_or_trust_signals":
        review_count = int(facts.get("review_count", 0) or 0)
        if review_count >= 25:
            return "Review depth strengthens purchase confidence"
        return "Opportunity to strengthen trust signals"
    return strategic_bullet_text(
        {"cue_key": cue_key, "classification": classification},
        identity,
        slide_key="slide4",
        evidence_terms=_slide4_theme_terms(facts),
    )


def _slide4_product_phrase_for_identity(facts: dict[str, Any], identity: dict[str, Any]) -> str:
    product = _safe_text(identity.get("product_type_display") or "").lower().replace("&", "and")
    content_blob = _slide4_content_blob(facts)
    if product and product not in {"category", "product"}:
        if product.endswith("s") and not product.endswith("ss"):
            product = product[:-1]
        if product == "facial cleanser":
            return "facial cleanser"
        safe_product = _slide4_sane_surface_phrase(product, content_blob)
        if safe_product:
            return safe_product
    return _slide4_product_phrase(facts)


def _slide4_category_context(facts: dict[str, Any], identity: dict[str, Any]) -> str:
    blob = " ".join(
        [
            _safe_text(identity.get("category_key")),
            _safe_text(identity.get("category_display")),
            _safe_text(facts.get("category")),
            _safe_text(facts.get("product_type")),
            _safe_text(facts.get("title")),
        ]
    ).lower()
    if any(term in blob for term in ("beauty", "skin", "facial", "cleanser", "cosmetic")):
        return "beauty"
    if any(term in blob for term in ("food", "pantry", "snack", "beverage", "spread", "butter", "jam", "jell")):
        return "food"
    if any(term in blob for term in ("electronics", "device", "charger", "battery", "headphone", "speaker")):
        return "electronics"
    return "general"


def _slide4_language_allowed(text: str, category_context: str) -> tuple[bool, str]:
    normalized = _safe_text(text).lower()
    invalid_by_context = {
        "beauty": ("nutrition", "protein", "breakfast", "snack", "recipe", "serving", "pantry"),
        "food": ("regimen", "skin-care", "skin care", "dermatologist", "clinical routine", "device setup"),
        "electronics": ("ingredient", "ingredients", "formula", "nutrition", "regimen", "skin"),
    }
    for term in invalid_by_context.get(category_context, ()):
        if term in normalized:
            return False, f"blocked_wrong_category_term:{term}"
    return True, "allowed"


def _slide4_family_for_cue(cue_key: str) -> str:
    if cue_key in {"product_positioning", "benefit_communication"}:
        return "positioning_title"
    if cue_key in {"ingredient_or_formula_communication", "pack_or_spec_detail", "conversion_guidance"}:
        return "detail_compliance"
    if cue_key in {"shopper_education", "usage_storytelling"}:
        return "education_storytelling"
    if cue_key in {"visual_identity", "review_or_trust_signals"}:
        return "trust_visual"
    return "education_storytelling"


def _slide4_add_candidate(
    candidates: list[dict[str, Any]],
    rejected: list[dict[str, str]],
    *,
    text: str,
    family: str,
    evidence_source: str,
    score: float,
    bullet_type: str,
    dimension: str,
    signals: list[str],
    reason: str,
    product_context: str,
    category_context: str,
    guide_context: dict[str, Any],
) -> None:
    clean = _slide4_sanitize_surface_text(text)
    if not clean:
        return
    allowed, guard_reason = _slide4_language_allowed(clean, category_context)
    if not allowed:
        rejected.append({"text": clean, "reason": guard_reason})
        return
    candidates.append(
        {
            "text": clean,
            "family": family if family in SLIDE4_BULLET_FAMILIES else "education_storytelling",
            "evidence_source": evidence_source,
            "score": score,
            "type": bullet_type,
            "dimension": dimension,
            "signals": [signal for signal in signals if signal],
            "reason": reason,
            "product_context": product_context,
            "category_context": category_context,
            "guide_context": guide_context,
        }
    )


def _slide4_image_analysis_facts(record: dict[str, Any]) -> dict[str, Any]:
    analysis = record.get("image_analysis", {}) or {}
    images = [
        image
        for image in (analysis.get("images", []) or [])
        if isinstance(image, dict) and image.get("status") == "analyzed"
    ]
    formats = Counter(_safe_text(image.get("probable_format")) for image in images if _safe_text(image.get("probable_format")))
    detected: set[str] = set()
    tokens: set[str] = set()
    for image in images:
        detected.update(_safe_text(signal).lower() for signal in (image.get("detected_signals") or []) if _safe_text(signal))
        tokens.update(_safe_text(token).lower() for token in (image.get("ocr_tokens") or []) if _safe_text(token))
    return {
        "analyzed_image_count": len(images),
        "formats": formats,
        "detected_signals": detected,
        "ocr_tokens": tokens,
        "has_lifestyle": formats["lifestyle_or_scene"] > 0,
        "has_detail": bool({"nutrition_or_ingredients", "dimensions_or_instructions"} & set(formats))
        or bool({"nutrition_or_ingredients", "size_or_count", "dimensions_or_scale"} & detected),
        "has_usage": bool({"usage_or_instructions", "recipe_or_serving"} & detected)
        or bool({"recipe", "serving", "toast", "routine", "apply", "use"} & tokens),
        "has_trust": bool({"organic_or_certification", "clinical_or_dermatologist", "guarantee"} & detected),
        "has_benefit": "feature_or_benefit_claim" in detected,
    }


def _slide4_title_has_product(facts: dict[str, Any], identity: dict[str, Any], product_phrase: str) -> bool:
    title = _safe_text(facts.get("title")).lower()
    product_terms = [
        product_phrase,
        _safe_text(identity.get("product_type_display")).lower(),
        *_as_text_list((identity.get("style_product_type") or {}).get("aliases") if isinstance(identity.get("style_product_type"), dict) else []),
        *_as_text_list((identity.get("style_product_type") or {}).get("title_keywords") if isinstance(identity.get("style_product_type"), dict) else []),
    ]
    return any(_safe_text(term).lower() and _safe_text(term).lower() in title for term in product_terms)


def _slide4_build_candidate_pool(
    record: dict[str, Any],
    *,
    side: str,
    existing_texts: set[str],
) -> tuple[list[dict[str, Any]], list[dict[str, str]], dict[str, Any]]:
    facts = _slide4_record_facts(record)
    identity = resolve_strategic_identity(
        [record],
        fallback_category=_safe_text(facts.get("category")),
        fallback_product_type=_safe_text(facts.get("product_type")),
    )
    guide_match = _build_image_guide_match(_safe_text(facts.get("category")), _safe_text(facts.get("product_type")))
    image_facts = _slide4_image_analysis_facts(record)
    blob = _slide4_content_blob(facts)
    product_phrase = _slide4_product_phrase_for_identity(facts, identity)
    category_context = _slide4_category_context(facts, identity)
    brand = facts["brand"] or side.replace("_", " ").title()
    is_client = side == "client"
    style_product = identity.get("style_product_type") if isinstance(identity.get("style_product_type"), dict) else {}
    guide_context = {
        "style_guide_path": identity.get("style_guide_path", ""),
        "image_guide_path": identity.get("image_guide_path", ""),
        "style_product_type_score": identity.get("style_product_type_score"),
        "title_priorities": list(identity.get("recommended_title_priorities") or [])[:4],
        "visual_priorities": list(identity.get("recommended_visual_priorities") or [])[:4],
        "image_required_slots": list(guide_match.get("required_slot_labels") or [])[:4],
    }
    candidates: list[dict[str, Any]] = []
    rejected: list[dict[str, str]] = []

    def add(**kwargs: Any) -> None:
        _slide4_add_candidate(
            candidates,
            rejected,
            product_context=product_phrase,
            category_context=category_context,
            guide_context=guide_context,
            **kwargs,
        )

    if _has_any(blob, "hazelnut", "cocoa", "nutella"):
        add(
            text="Hazelnut-cocoa positioning supports breakfast occasions",
            family="positioning_title",
            evidence_source="title/product_identification",
            score=93,
            bullet_type="strength",
            dimension="product_positioning",
            signals=["hazelnut", "cocoa", "title"],
            reason="Title or PDP copy identifies hazelnut/cocoa spread positioning.",
        )
    elif _has_any(blob, "peanut butter", "peanut", "protein", "jif", "fresh roasted", "fresh-roasted"):
        add(
            text="Peanut butter positioning reinforces pantry relevance",
            family="positioning_title",
            evidence_source="title/product_identification",
            score=93,
            bullet_type="strength",
            dimension="product_positioning",
            signals=["peanut butter", "protein", "title"],
            reason="Title, product type, or PDP copy identifies peanut butter/protein positioning.",
        )
    elif _slide4_title_has_product(facts, identity, product_phrase):
        add(
            text=f"Title clarity makes the {product_phrase} role easier to understand",
            family="positioning_title",
            evidence_source="title_formula",
            score=90,
            bullet_type="strength",
            dimension="title_product_type_identification",
            signals=["title", product_phrase],
            reason="Product type or guide-supported aliases appear in the title.",
        )
    else:
        add(
            text="Title wording can make the product role more direct",
            family="positioning_title",
            evidence_source="title_formula",
            score=82,
            bullet_type="opportunity",
            dimension="title_product_type_identification",
            signals=["title_gap", product_phrase],
            reason="Resolved product type was not strongly represented in the title.",
        )

    key_features = facts.get("key_features", []) or []
    description = _safe_text(facts.get("description"))
    if _has_any(blob, "protein", "nutrition", "calories", "ingredients") and category_context == "food":
        add(
            text="Protein cues strengthen snack and pantry relevance",
            family="detail_compliance",
            evidence_source="pdp_detail_compliance",
            score=91,
            bullet_type="strength",
            dimension="detail_compliance",
            signals=["protein", "nutrition", "key_features"],
            reason="Food PDP evidence includes nutrition, ingredient, or protein detail.",
        )
        add(
            text="Clear pack and nutrition detail",
            family="detail_compliance",
            evidence_source="style_title_and_pdp_detail",
            score=89,
            bullet_type="strength",
            dimension="pack_or_spec_detail",
            signals=["nutrition", "ingredients", "style_guide"],
            reason="PDP and guide context support food-specific pack/nutrition detail.",
        )
    elif category_context == "beauty":
        detail_text = (
            "Formula and skin-benefit details make comparison easier"
            if _has_any(blob, "hydrating", "sensitive", "fragrance", "ingredient", "formula")
            else "Benefit detail can clarify skin-fit and routine relevance"
        )
        add(
            text=detail_text,
            family="detail_compliance",
            evidence_source="pdp_detail_compliance",
            score=88 if _has_any(blob, "hydrating", "sensitive", "ingredient", "formula") else 80,
            bullet_type="strength" if _has_any(blob, "hydrating", "sensitive", "ingredient", "formula") else "opportunity",
            dimension="formula_detail",
            signals=["formula", "skin_benefit", "description"],
            reason="Beauty PDP evidence is evaluated with formula/benefit language instead of food detail.",
        )
    elif category_context == "electronics":
        add(
            text="Spec detail helps shoppers compare options",
            family="detail_compliance",
            evidence_source="pdp_detail_compliance",
            score=86,
            bullet_type="strength" if description or key_features else "opportunity",
            dimension="spec_detail",
            signals=["spec", "description", "key_features"],
            reason="Electronics PDP detail is framed around specs and comparison support.",
        )
    elif description or key_features:
        add(
            text="Feature detail supports easier comparison",
            family="detail_compliance",
            evidence_source="pdp_detail_compliance",
            score=84,
            bullet_type="strength",
            dimension="feature_detail",
            signals=["description", "key_features"],
            reason="PDP description or key features provide comparison detail.",
        )
    else:
        add(
            text="Description detail can make comparison easier",
            family="detail_compliance",
            evidence_source="pdp_detail_compliance",
            score=76,
            bullet_type="opportunity",
            dimension="description_quality",
            signals=["description_gap"],
            reason="Description and key-feature evidence was sparse.",
        )

    if category_context == "food" and _has_any(blob, "breakfast", "snack", "recipe", "toast", "pantry"):
        add(
            text=f"{product_phrase.title()} cues connect PDP content to breakfast",
            family="education_storytelling",
            evidence_source="pdp_storytelling",
            score=88,
            bullet_type="strength",
            dimension="usage_storytelling",
            signals=["breakfast", "snack", "recipe", "description"],
            reason="Usage occasion language was present in PDP title, description, or features.",
        )
    elif category_context == "beauty" and _has_any(blob, "daily", "routine", "use", "sensitive", "hydrating"):
        add(
            text="Usage cues help connect the product to routine fit",
            family="education_storytelling",
            evidence_source="pdp_storytelling",
            score=88,
            bullet_type="strength",
            dimension="usage_storytelling",
            signals=["routine", "sensitive", "hydrating"],
            reason="Beauty PDP evidence includes usage, routine, or skin-need language.",
        )
    elif image_facts["has_usage"] or image_facts["has_lifestyle"]:
        add(
            text="Image variety adds usage and routine context",
            family="education_storytelling",
            evidence_source="image_support",
            score=83,
            bullet_type="strength",
            dimension="visual_storytelling",
            signals=["image_usage", "lifestyle"],
            reason="Image analysis detected usage, recipe, routine, or lifestyle support.",
        )
    else:
        add(
            text="Usage storytelling can make product fit clearer",
            family="education_storytelling",
            evidence_source="pdp_storytelling",
            score=77,
            bullet_type="opportunity",
            dimension="usage_storytelling",
            signals=["limited_usage_language"],
            reason="PDP and image evidence did not show strong usage education.",
        )

    if facts["review_count"] >= 100:
        add(
            text=(
                f"{brand} review volume raises comparison confidence"
                if not is_client
                else "Review depth helps reinforce shopper confidence"
            ),
            family="trust_visual",
            evidence_source="review_trust",
            score=87,
            bullet_type="strength",
            dimension="review_authority",
            signals=[f"{facts['review_count']}_reviews"],
            reason="Review count evidence was high enough to support shopper confidence.",
        )
    if facts["image_count"] >= 6 or image_facts["analyzed_image_count"] >= 3:
        add(
            text=(
                f"{facts['image_count']}-image carousel supports visual education"
                if is_client
                else f"{brand} carousel supports visual education"
            ),
            family="trust_visual",
            evidence_source="image_guide_support",
            score=84,
            bullet_type="strength",
            dimension="image_sequence",
            signals=[f"{facts['image_count']}_images", *guide_context["image_required_slots"][:2]],
            reason="Carousel depth and image-guide context support visual comparison.",
        )
    if facts["sold_by_walmart"] or facts["shipped_by_walmart"] or facts["ebc_present"] or image_facts["has_trust"]:
        add(
            text="Trust cues help reduce hesitation at purchase",
            family="trust_visual",
            evidence_source="review_trust",
            score=82,
            bullet_type="strength",
            dimension="trust_support",
            signals=["walmart_fulfillment" if facts["sold_by_walmart"] or facts["shipped_by_walmart"] else "", "enhanced_content" if facts["ebc_present"] else "", "image_trust" if image_facts["has_trust"] else ""],
            reason="Fulfillment, enhanced content, or image trust signals were present.",
        )

    cue_bullets, cue_debug, cue_context_debug = _translate_slide4_strategic_cues(
        record,
        facts,
        side=side,
        existing_texts=set(existing_texts),
    )
    for text, cue in zip(cue_bullets, cue_debug):
        dimension = _safe_text(cue.get("dimension", "cue"))
        add(
            text=text,
            family=_slide4_family_for_cue(dimension),
            evidence_source="strategic_cue_engine",
            score=72,
            bullet_type=_safe_text(cue.get("type", "context")) or "context",
            dimension=dimension,
            signals=[*_as_text_list(cue.get("signals"))[:3]],
            reason=_safe_text(cue.get("reason")) or "Selected from strategic cue engine output.",
        )

    return candidates, rejected, {
        "identity": {
            "category_key": identity.get("category_key"),
            "product_type_display": identity.get("product_type_display"),
            "style_guide_path": identity.get("style_guide_path", ""),
            "image_guide_path": identity.get("image_guide_path", ""),
        },
        "product_context": product_phrase,
        "category_context": category_context,
        "style_terms_used": {
            "aliases": _as_text_list(style_product.get("aliases"))[:5],
            "title_keywords": _as_text_list(style_product.get("title_keywords"))[:5],
            "context_keywords": _as_text_list(style_product.get("context_keywords"))[:5],
            "attributes": _as_text_list(style_product.get("attributes"))[:5],
            "negative_keywords": _as_text_list(style_product.get("negative_keywords"))[:5],
        },
        "image_guide": guide_context,
        "image_analysis": {
            "analyzed_image_count": image_facts["analyzed_image_count"],
            "detected_signals": sorted(image_facts["detected_signals"])[:8],
            "formats": dict(image_facts["formats"]),
        },
        "cue_context": cue_context_debug,
        "rejected_language": rejected,
    }


def _select_slide4_balanced_candidates(
    candidates: list[dict[str, Any]],
    *,
    existing_texts: set[str],
) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    ranked = sorted(candidates, key=lambda item: float(item.get("score", 0)), reverse=True)
    selected: list[dict[str, Any]] = []
    family_counts: Counter[str] = Counter()
    used = set(existing_texts)
    rejected: list[dict[str, str]] = []

    def try_add(candidate: dict[str, Any], *, allow_second_family: bool) -> bool:
        text = _safe_text(candidate.get("text"))
        key = normalize_bullet_text(text)
        family = _safe_text(candidate.get("family"))
        if not text or key in used:
            rejected.append({"text": text, "reason": "duplicate_or_cross_column_repeat"})
            return False
        if family_counts[family] >= 2:
            rejected.append({"text": text, "reason": "family_cap_reached"})
            return False
        if family_counts[family] >= 1 and not allow_second_family:
            return False
        used.add(key)
        family_counts[family] += 1
        selected.append(candidate)
        return True

    for family in SLIDE4_BULLET_FAMILIES:
        family_candidates = [candidate for candidate in ranked if candidate.get("family") == family]
        for candidate in family_candidates:
            if try_add(candidate, allow_second_family=False):
                break
        if len(selected) == 6:
            break

    if len(selected) < 6:
        for candidate in ranked:
            if len(selected) == 6:
                break
            if candidate in selected:
                continue
            try_add(candidate, allow_second_family=True)

    return selected[:6], {
        "family_counts": dict(family_counts),
        "ranked_candidates": [
            {
                "text": candidate.get("text"),
                "family": candidate.get("family"),
                "evidence_source": candidate.get("evidence_source"),
                "score": candidate.get("score"),
            }
            for candidate in ranked[:12]
        ],
        "rejected_candidates": rejected[:12],
    }


def _slide4_candidate_sort_key(candidate: dict[str, Any]) -> tuple[int, int, float, float]:
    class_rank = {"strength": 0, "opportunity": 1, "context": 2, "pressure": 3}
    cue_rank = {
        "product_positioning": 0,
        "benefit_communication": 1,
        "ingredient_or_formula_communication": 2,
        "pack_or_spec_detail": 3,
        "shopper_education": 4,
        "usage_storytelling": 5,
        "visual_identity": 6,
        "review_or_trust_signals": 7,
        "conversion_guidance": 8,
    }
    return (
        class_rank.get(candidate.get("classification"), 9),
        cue_rank.get(candidate.get("cue_key"), 99),
        -float(candidate.get("coverage_ratio", 0) or 0),
        -float(candidate.get("strength_ratio", 0) or 0),
    )


def _translate_slide4_strategic_cues(
    record: dict[str, Any],
    facts: dict[str, Any],
    *,
    side: str,
    existing_texts: set[str],
) -> tuple[list[str], list[dict[str, Any]], dict[str, Any]]:
    context = aggregate_strategic_cues([record])
    candidates = [
        candidate
        for candidate in context.get("candidate_cues", [])
        if "slide4_pdp_benchmark" in (candidate.get("slide_objective_tags") or [])
    ]
    bullets: list[str] = []
    debug: list[dict[str, Any]] = []
    used = existing_texts
    for candidate in sorted(candidates, key=_slide4_candidate_sort_key):
        text = _slide4_theme_bullet(
            candidate.get("cue_key", ""),
            candidate.get("classification", ""),
            facts,
            context.get("identity", {}),
        )
        key = normalize_bullet_text(text)
        if not text or key in used:
            continue
        used.add(key)
        bullets.append(text)
        debug.append(
            {
                "text": text,
                "type": candidate.get("classification", "context"),
                "dimension": candidate.get("cue_key", "cue"),
                "signals": [
                    candidate.get("classification", ""),
                    candidate.get("cue_key", ""),
                    *candidate.get("evidence_sources", [])[:2],
                ],
                "reason": candidate.get("debug_reason", "Selected from strategic cue engine output."),
                "cue_debug": candidate,
                "strategic_cue_context": context.get("debug", {}),
            }
        )
        if len(bullets) >= 6:
            break
    fallback_candidates = [
        "Benefit-forward product positioning",
        "Clear pack and nutrition detail",
        "Structured shopper education",
        "Opportunity to deepen serving guidance",
        "Review and trust cues support conversion confidence",
        "Visual sequencing can clarify shopper priorities",
    ]
    for text in fallback_candidates:
        if len(bullets) >= 6:
            break
        key = normalize_bullet_text(text)
        if key in used:
            continue
        used.add(key)
        bullets.append(text)
        debug.append(
            {
                "text": text,
                "type": "opportunity",
                "dimension": "fallback",
                "signals": ["fallback"],
                "reason": "Controlled fallback used to preserve Slide 4 bullet count.",
                "strategic_cue_context": context.get("debug", {}),
            }
        )
    return bullets[:6], debug[:6], context.get("debug", {})


def _build_slide4_evidence_bullets(
    record: dict[str, Any],
    *,
    side: str,
    existing_texts: set[str],
) -> tuple[list[str], list[dict[str, Any]], list[str]]:
    facts = _slide4_record_facts(record)
    candidates, rejected_language, evidence_context = _slide4_build_candidate_pool(
        record,
        side=side,
        existing_texts=existing_texts,
    )
    selected_candidates, selection_debug = _select_slide4_balanced_candidates(
        candidates,
        existing_texts=existing_texts,
    )
    if len(selected_candidates) == 6:
        for candidate in selected_candidates:
            existing_texts.add(normalize_bullet_text(candidate["text"]))
        return (
            [candidate["text"] for candidate in selected_candidates],
            [
                {
                    "text": candidate["text"],
                    "type": candidate.get("type", "context"),
                    "dimension": candidate.get("dimension", candidate.get("family", "")),
                    "signals": candidate.get("signals", []),
                    "reason": candidate.get("reason", ""),
                    "bullet_family": candidate.get("family", ""),
                    "evidence_family_source": candidate.get("evidence_source", ""),
                    "selected_because": (
                        "Chosen by Slide 4 balanced PDP selection using supported evidence, "
                        "family diversity, product-type fit, and cross-column uniqueness."
                    ),
                    "product_type_context": candidate.get("product_context", ""),
                    "category_context": candidate.get("category_context", ""),
                    "guide_context": candidate.get("guide_context", {}),
                    "selection_debug": selection_debug,
                    "evidence_context": evidence_context,
                    "rejected_language": rejected_language,
                }
                for candidate in selected_candidates
            ],
            [],
        )
    blob = _slide4_content_blob(facts)
    image_blob = facts["image_text"]
    bullets: list[str] = []
    debug: list[dict[str, Any]] = []
    warnings: list[str] = []
    used = existing_texts
    product_phrase = _slide4_product_phrase(facts)
    phrases = _slide4_content_phrases(facts)
    brand = facts["brand"] or side.replace("_", " ").title()
    is_client = side == "client"
    cue_bullets, cue_debug, _cue_context_debug = _translate_slide4_strategic_cues(
        record,
        facts,
        side=side,
        existing_texts=used,
    )
    if len(cue_bullets) == 6:
        return cue_bullets, cue_debug[:6], warnings

    if _has_any(blob, "hazelnut", "cocoa", "nutella"):
        _append_unique_bullet(
            bullets,
            debug,
            text=(
                "Client hazelnut-cocoa positioning supports breakfast relevance"
                if is_client
                else f"{brand} sets a hazelnut-cocoa comparison benchmark"
            ),
            bullet_type="strength",
            dimension="ingredient_positioning",
            signals=["hazelnut", "cocoa"],
            reason="Product title or PDP text references hazelnut/cocoa spread positioning.",
            used=used,
        )
    if _has_any(blob, "peanut", "protein", "fresh roasted", "fresh-roasted", "jif"):
        _append_unique_bullet(
            bullets,
            debug,
            text=(
                "Client peanut-butter PDP content reinforces pantry cues"
                if is_client
                else f"{brand} frames peanut-butter protein cues for comparison"
            ),
            bullet_type="strength",
            dimension="ingredient_positioning",
            signals=["peanut", "protein"],
            reason="Product title, description, or features reference peanut butter/protein cues.",
            used=used,
        )
    if facts["image_count"] >= 6:
        _append_unique_bullet(
            bullets,
            debug,
            text=(
                "Image variety adds usage and routine context"
                if is_client
                else f"{brand} uses the carousel to add comparison context"
            ),
            bullet_type="opportunity",
            dimension="image_sequence",
            signals=[f"{facts['image_count']}_images"],
            reason="Image count indicates enough carousel depth to sequence shopper education.",
            used=used,
        )
    if _has_any(blob, "nutrition", "calories", "protein", "ingredients") or _has_any(image_blob, "nutrition", "ingredient"):
        _append_unique_bullet(
            bullets,
            debug,
            text=(
                "Pack and nutrition cues make product value clearer"
                if is_client
                else f"{brand} uses pack and nutrition as proof points"
            ),
            bullet_type="strength",
            dimension="nutrition_detail",
            signals=["nutrition", "ingredients"],
            reason="Nutrition, ingredient, or pack-detail terms were detected in PDP evidence.",
            used=used,
        )
    if _has_any(blob, "breakfast", "snack", "recipe", "toast", "pantry"):
        _append_unique_bullet(
            bullets,
            debug,
            text=(
                f"{product_phrase.title()} cues connect PDP content to breakfast"
                if is_client
                else f"{brand} ties discovery to pantry occasions"
            ),
            bullet_type="strength",
            dimension="usage_occasions",
            signals=["breakfast/snack/pantry"],
            reason="Usage occasion language was present in title, description, or features.",
            used=used,
        )
    else:
        _append_unique_bullet(
            bullets,
            debug,
            text=(
                "PDP content can clarify usage occasions"
                if is_client
                else f"{brand} can make usage occasions clearer"
            ),
            bullet_type="opportunity",
            dimension="usage_occasions",
            signals=["limited_usage_language"],
            reason="No clear breakfast, snack, recipe, or pantry use-case language was detected.",
            used=used,
        )
    if facts["review_count"] >= 100:
        _append_unique_bullet(
            bullets,
            debug,
            text=(
                "Review depth helps reinforce shopper confidence"
                if is_client
                else f"{brand} review volume raises comparison expectations"
            ),
            bullet_type="strength",
            dimension="review_authority",
            signals=[f"{facts['review_count']}_reviews"],
            reason="Review count evidence was available and materially high.",
            used=used,
        )
    if facts["ebc_present"]:
        _append_unique_bullet(
            bullets,
            debug,
            text=(
                "Enhanced content can deepen usage guidance"
                if is_client
                else f"{brand} enhanced content strengthens shopper education"
            ),
            bullet_type="strength",
            dimension="enhanced_content",
            signals=["enhanced_brand_content"],
            reason="Enhanced Brand Content is marked present in the PDP evidence.",
            used=used,
        )
    if facts["sold_by_walmart"] or facts["shipped_by_walmart"]:
        _append_unique_bullet(
            bullets,
            debug,
            text=(
                "Walmart fulfillment cues support purchase confidence"
                if is_client
                else f"{brand} fulfillment cues support purchase confidence"
            ),
            bullet_type="strength",
            dimension="retail_trust",
            signals=["sold_by_walmart" if facts["sold_by_walmart"] else "", "shipped_by_walmart" if facts["shipped_by_walmart"] else ""],
            reason="Sold/shipped by Walmart signals were present in PDP evidence.",
            used=used,
        )
    if not _has_any(image_blob, "recipe", "serving", "lifestyle", "breakfast", "snack"):
        _append_unique_bullet(
            bullets,
            debug,
            text=(
                "Usage imagery can make product fit easier to picture"
                if is_client
                else f"{brand} can add usage context to discovery"
            ),
            bullet_type="opportunity",
            dimension="visual_storytelling",
            signals=["limited_recipe_lifestyle_image_hints"],
            reason="Image metadata did not show recipe, serving, or lifestyle cues.",
            used=used,
        )
    if not facts["ebc_present"]:
        _append_unique_bullet(
            bullets,
            debug,
            text=(
                "Enhanced storytelling can add clearer usage guidance"
                if is_client
                else f"{brand} can add richer shopper education"
            ),
            bullet_type="opportunity",
            dimension="enhanced_content",
            signals=["ebc_not_present"],
            reason="Enhanced Brand Content was not marked present.",
            used=used,
        )

    fallback_candidates = [
        (
            f"{brand} keeps the product role visible",
            "strength",
            "brand_specificity",
            ["brand", product_phrase],
            "Fallback used brand/product-type evidence to avoid generic duplicate wording.",
        ),
        (
            "Opening image sequence can clarify shopper priorities",
            "opportunity",
            "opening_sequence",
            ["opening_sequence"],
            "Controlled fallback opportunity used when richer evidence was limited.",
        ),
        (
            "Carousel education can strengthen PDP content",
            "opportunity",
            "shopper_education",
            ["carousel_education"],
            "Controlled fallback opportunity used to preserve template bullet count.",
        ),
        (
            f"{facts['brand'] or side.title()} can broaden variant storytelling",
            "opportunity",
            "variant_storytelling",
            ["variant_occasion_storytelling"],
            "Brand-specific fallback used after duplicate bullets were removed.",
        ),
        (
            f"{facts['brand'] or side.title()} can make shopper education more distinctive",
            "opportunity",
            "distinctive_education",
            ["distinctive_shopper_education"],
            "Brand-specific fallback used after duplicate bullets were removed.",
        ),
    ]
    for text, bullet_type, dimension, signals, reason in fallback_candidates:
        if len(bullets) >= 6:
            break
        _append_unique_bullet(
            bullets,
            debug,
            text=text,
            bullet_type=bullet_type,
            dimension=dimension,
            signals=signals,
            reason=reason,
            used=used,
        )
    if len(bullets) < 6:
        for index in range(len(bullets), 6):
            _append_unique_bullet(
                bullets,
                debug,
                text=f"{facts['brand'] or side.title()} can add clearer shopper-facing detail",
                bullet_type="opportunity",
                dimension=f"controlled_fallback_{index + 1}",
                signals=["controlled_fallback"],
                reason="Last-resort side-specific fallback used to fill the existing template bullet structure without duplicates.",
                used=used,
            )
    if len(bullets) < 6:
        warnings.append("Slide 4 used restrained fallback bullets because PDP evidence was sparse.")
    return bullets[:6], debug[:6], warnings


def _build_slide4_bullets(images: list[dict[str, Any]], guide_match: dict[str, Any]) -> list[str]:
    bullets = [f"Carousel: {len(images)} ordered images"]
    dimensions = [dims for image in images if (dims := _parse_image_dimensions(image))]
    unique_dimensions = list(dict.fromkeys(dimensions))
    if len(unique_dimensions) == 1:
        suffix = " throughout" if len(dimensions) == len(images) else " detected"
        bullets.append(f"Asset sizing is consistent{suffix}")
    elif len(unique_dimensions) > 1:
        bullets.append("Asset sizing varies across the image stack")

    if guide_match.get("matched"):
        labels = list(guide_match.get("required_slot_labels", []) or [])
        if labels:
            bullets.append("Recommended opening: " + " / ".join(labels[:3]))
        if len(labels) > 3:
            bullets.append("Recommended support: " + " / ".join(labels[3:6]))
        additional = list(guide_match.get("additional_recommendations", []) or [])
        if additional:
            opportunity = _safe_text(additional[0]).split(":", 1)[0]
            if opportunity:
                bullets.append(f"Guide opportunity: {opportunity}")
    return bullets[:5]


def _build_slide4_column(
    record: dict[str, Any],
    fallback_label: str,
    findings: dict[str, Any] | None = None,
    *,
    side: str = "client",
    existing_bullet_texts: set[str] | None = None,
) -> dict[str, Any]:
    existing_bullet_texts = existing_bullet_texts if existing_bullet_texts is not None else set()
    if not record:
        return {
            "label": fallback_label,
            "brand": "",
            "product_title": "",
            "product_id": "",
            "category": "",
            "product_type": "",
            "ordered_images": [],
            "image_count": 0,
            "image_guide_match": {"matched": False},
            "bullets": [],
            "bullet_debug": [],
            "warnings": [],
            "findings": findings or {},
            "active": False,
        }
    images = _normalize_ordered_images(record)
    product_type = _safe_text(_first_value(record, "product_type", "subcategory", "Product Type", "productType"))
    category = _safe_text(_first_value(record, "category", "Category", "categoryPathName"))
    guide_match = _build_image_guide_match(category, product_type)
    findings = findings or build_slide4_group_findings([record], fallback_label)
    evidence_bullets, bullet_debug, warnings = _build_slide4_evidence_bullets(
        record,
        side=side,
        existing_texts=existing_bullet_texts,
    )
    if not evidence_bullets:
        finding_bullets = [
            _safe_text(bullet)
            for bullet in (findings.get("slide4_bullets", []) if isinstance(findings, dict) else [])
            if _safe_text(bullet) and _safe_text(bullet).lower() not in existing_bullet_texts
        ]
        for bullet in finding_bullets:
            existing_bullet_texts.add(bullet.lower())
        evidence_bullets = finding_bullets[:6] or _build_slide4_bullets(images, guide_match)
    return {
        "label": _safe_text(_first_value(record, "brand", "Brand", "brandName")) or fallback_label,
        "brand": _safe_text(_first_value(record, "brand", "Brand", "brandName")),
        "product_title": _safe_text(_first_value(record, "product_title", "Product Title", "title", "productTitle")),
        "product_id": _safe_text(_first_value(record, "item_id", "product_id", "Product ID", "productId", "record_id")),
        "category": category,
        "product_type": product_type,
        "ordered_images": images,
        "image_count": int(_number_value(_first_value(record, "image_count", "Image Count", default=len(images))) or len(images)),
        "review_count": int(_number_value(_first_value(record, "review_count", "Review Count", "reviews_summary.review_count", "reviews_summary.count"))),
        "average_rating": _number_value(_first_value(record, "average_rating", "Average Rating", "reviews_summary.average_rating", "reviews_summary.rating")),
        "ebc_present": _truthy_evidence(_first_value(record, "enhanced_brand_content", "enhanced_brand_content_present", "Enhanced Brand Content", "enhancedBrandContentPresent")),
        "sold_by_walmart": _truthy_evidence(_first_value(record, "sold_by_walmart", "Sold by Walmart")),
        "shipped_by_walmart": _truthy_evidence(_first_value(record, "shipped_by_walmart", "Shipped by Walmart")),
        "image_guide_match": guide_match,
        "bullets": evidence_bullets[:6],
        "bullet_debug": bullet_debug,
        "warnings": warnings,
        "findings": findings,
        "active": True,
    }


def build_slide4_pdp_benchmark_payload(
    export_plan: dict, competitor_records: list[dict[str, Any]] | None = None
) -> dict[str, Any]:
    """Build deterministic Slide 4 content from one client and two competitor PDPs."""
    metadata = export_plan.get("audit_metadata", {}) or {}
    product_pairs = export_plan.get("product_slide_pairs", []) or []
    client_company_name = _safe_text(metadata.get("client_company_name", ""))
    client_name = _safe_text(metadata.get("client_name", ""))
    client_record = (
        dict((product_pairs[0].get("pdp_slide", {}) or {}))
        if product_pairs
        else {}
    )
    client_fallback = (
        client_company_name
        or client_name
        or _safe_text(client_record.get("brand"))
        or "Client"
    )
    slide4_findings = export_plan.get("slide4_findings", {}) or {}
    used_bullet_texts: set[str] = set()
    client_column = _build_slide4_column(
        client_record,
        client_fallback,
        slide4_findings.get("client") if isinstance(slide4_findings, dict) else None,
        side="client",
        existing_bullet_texts=used_bullet_texts,
    )
    client_column["label"] = client_company_name or client_name or client_column["label"]

    competitor_records = list(competitor_records or [])
    competitor_columns = [
        _build_slide4_column(
            competitor_records[index] if index < len(competitor_records) else {},
            f"Competitor {index + 1}",
            (
                slide4_findings.get(f"competitor_{index + 1}")
                if isinstance(slide4_findings, dict)
                else None
            ),
            side=f"competitor_{index + 1}",
            existing_bullet_texts=used_bullet_texts,
        )
        for index in range(2)
    ]
    active_competitors = [column for column in competitor_columns if column.get("active")]
    layout_mode = "2-column" if len(active_competitors) == 1 else "3-column"
    active_columns_debug = [
        {"index": index, "label": column.get("label", ""), "active": bool(column.get("active"))}
        for index, column in enumerate([client_column, *competitor_columns])
        if column.get("active")
    ]
    hidden_columns = [
        {
            "index": index + 1,
            "label": column["label"],
            "reason": "No PDP evidence was available.",
        }
        for index, column in enumerate(competitor_columns)
        if not column.get("active")
    ]
    columns = [client_column, *competitor_columns]
    all_bullets: set[str] = set()
    removed_duplicates: list[dict[str, str]] = []
    for column in columns:
        final_bullets = []
        for bullet in column.get("bullets", []) or []:
            normalized = normalize_bullet_text(bullet)
            if normalized in all_bullets:
                removed_duplicates.append({"column": column.get("label", ""), "bullet": bullet})
                continue
            all_bullets.add(normalized)
            final_bullets.append(bullet)
        column["bullets"] = final_bullets[:6]
    return {
        "columns": columns,
        "slide4_findings": slide4_findings,
        "hidden_columns": hidden_columns,
        "layout_mode": layout_mode,
        "warnings": [
            f"{item['label']} was cleared because no PDP evidence was available."
            for item in hidden_columns
        ],
        "debug": {
            "layout_mode": layout_mode,
            "layout_debug": {
                "layout_mode": layout_mode,
                "active_columns": active_columns_debug,
                "hidden_columns": hidden_columns,
                "active_column_count": len(active_columns_debug),
                "render_slots": [0, 2] if layout_mode == "2-column" and len(active_columns_debug) == 2 else [0, 1, 2],
                "reused_regions": (
                    "outer_template_regions_rebalanced_as_left_right_columns"
                    if layout_mode == "2-column" and len(active_columns_debug) == 2
                    else "native_three_column_template_regions"
                ),
                "divider_cleanup": (
                    "suppress_middle_column_label_rule_and_vertical_3_column_dividers"
                    if layout_mode == "2-column" and len(active_columns_debug) == 2
                    else "preserve_three_column_chrome"
                ),
            },
            "hidden_columns": hidden_columns,
            "resolved_labels": [column.get("label", "") for column in columns],
            "resolved_category": client_column.get("category", ""),
            "resolved_product_type": client_column.get("product_type", ""),
            "evidence_used_per_bullet": {
                column.get("label", f"column_{index + 1}"): column.get("bullet_debug", [])
                for index, column in enumerate(columns)
            },
            "bullets_removed_for_duplication": removed_duplicates,
            "ocr_usage": False,
            "final_bullets": {
                column.get("label", f"column_{index + 1}"): column.get("bullets", [])
                for index, column in enumerate(columns)
            },
            "render_targets": {
                "target_bullet_count": 6,
                "final_bullet_counts": {
                    column.get("label", f"column_{index + 1}"): len(column.get("bullets", []) or [])
                    for index, column in enumerate(columns)
                    if column.get("active")
                },
            },
        },
    }


def _slide4_label_shapes(slide: Any) -> list[Any]:
    candidates = []
    for shape in _walk_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        text = _shape_text(shape)
        if not text or len(text) > 80:
            continue
        top = int(getattr(shape, "top", 0) or 0)
        if Inches(2.1) <= top <= Inches(2.8):
            candidates.append(shape)
    return sorted(candidates, key=lambda shape: int(getattr(shape, "left", 0) or 0))[:3]


def _slide4_bullet_shapes(slide: Any) -> list[Any]:
    candidates = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and int(getattr(shape, "top", 0) or 0) >= Inches(5.4)
        and int(getattr(shape, "height", 0) or 0) >= Inches(0.8)
    ]
    return sorted(candidates, key=lambda shape: int(getattr(shape, "left", 0) or 0))[:3]


def _copy_basic_paragraph_style(source: Any, target: Any) -> None:
    target.level = getattr(source, "level", 0)
    target.alignment = getattr(source, "alignment", None)
    target.space_before = getattr(source, "space_before", None)
    target.space_after = getattr(source, "space_after", None)
    target.line_spacing = getattr(source, "line_spacing", None)
    source_runs = list(getattr(source, "runs", []) or [])
    if not source_runs:
        return
    if not target.runs:
        target.add_run()
    source_font = source_runs[0].font
    target_font = target.runs[0].font
    target_font.name = source_font.name
    target_font.size = source_font.size
    target_font.bold = source_font.bold
    target_font.italic = source_font.italic


def _replace_bullet_shape_text(
    shape: Any,
    bullets: list[str],
    *,
    ensure_paragraph_count: bool = False,
) -> dict[str, int]:
    text_frame = shape.text_frame
    existing_count = len(text_frame.paragraphs)
    created_count = 0
    if ensure_paragraph_count and bullets:
        while len(text_frame.paragraphs) < len(bullets):
            template_paragraph = text_frame.paragraphs[-1]
            new_paragraph = text_frame.add_paragraph()
            _copy_basic_paragraph_style(template_paragraph, new_paragraph)
            created_count += 1
    paragraphs = list(text_frame.paragraphs)
    for index, paragraph in enumerate(paragraphs):
        _replace_paragraph_text_preserve_style(
            paragraph, bullets[index] if index < len(bullets) else ""
        )
    return {
        "existing_paragraph_count": existing_count,
        "created_paragraph_count": created_count,
        "available_paragraph_count": len(text_frame.paragraphs),
    }


def _apply_bullet_spacing_and_font(
    shape: Any,
    font_size: int | None,
    *,
    font_name: str | None = None,
    line_spacing: float = 1.0,
    space_after: int = 0,
    margin_top: int | None = None,
    margin_bottom: int | None = None,
    margin_left: int | None = None,
    margin_right: int | None = None,
) -> None:
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    if margin_top is not None:
        text_frame.margin_top = Pt(margin_top)
    if margin_bottom is not None:
        text_frame.margin_bottom = Pt(margin_bottom)
    if margin_left is not None:
        text_frame.margin_left = Pt(margin_left)
    if margin_right is not None:
        text_frame.margin_right = Pt(margin_right)
    if font_size is not None:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for paragraph in text_frame.paragraphs:
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(space_after)
        paragraph.line_spacing = line_spacing
        if font_size is None:
            if font_name is None:
                continue
        for run in paragraph.runs:
            if font_name is not None:
                run.font.name = font_name
            if font_size is not None:
                run.font.size = Pt(font_size)


def _fit_bullet_shape_group(
    items: list[tuple[str, Any, list[str]]],
    *,
    max_lines: int,
    drop_threshold: int,
    font_threshold: int,
    base_font_size: int | None = None,
    fallback_font_size: int = 10,
    allow_drop: bool = False,
    ensure_paragraph_count: bool = False,
    font_name: str | None = None,
    line_spacing: float = 1.0,
    space_after: int = 0,
    margin_top: int | None = None,
    margin_bottom: int | None = None,
    margin_left: int | None = None,
    margin_right: int | None = None,
) -> dict[str, dict[str, Any]]:
    prepared: dict[str, dict[str, Any]] = {}
    global_font = base_font_size
    for key, shape, bullets in items:
        rendered = [
            _safe_text(bullet)
            for bullet in bullets
            if _safe_text(bullet)
        ][:max_lines]
        dropped: list[str] = []
        while allow_drop and len(rendered) > 1 and sum(len(bullet) for bullet in rendered) > drop_threshold:
            dropped.insert(0, rendered.pop())
        if sum(len(bullet) for bullet in rendered) > font_threshold:
            global_font = fallback_font_size
        prepared[key] = {
            "shape": shape,
            "rendered": rendered,
            "dropped": dropped,
            "render_target_count": min(len([bullet for bullet in bullets if _safe_text(bullet)]), max_lines),
            "target_bullet_count": max_lines,
            "source_bullet_count": len([bullet for bullet in bullets if _safe_text(bullet)]),
            "dropped_bullets_reason": (
                "length_fit_threshold"
                if dropped
                else "not_dropped"
            ),
        }

    for key, data in prepared.items():
        shape = data["shape"]
        paragraph_debug = _replace_bullet_shape_text(
            shape,
            data["rendered"],
            ensure_paragraph_count=ensure_paragraph_count,
        )
        _apply_bullet_spacing_and_font(
            shape,
            global_font,
            font_name=font_name,
            line_spacing=line_spacing,
            space_after=space_after,
            margin_top=margin_top,
            margin_bottom=margin_bottom,
            margin_left=margin_left,
            margin_right=margin_right,
        )
        data["font_fallback"] = global_font
        data["font_size_selected"] = global_font
        data["font_name_selected"] = font_name
        data["line_spacing_selected"] = line_spacing
        data["space_after_selected"] = space_after
        data["margin_top_selected"] = margin_top
        data["margin_bottom_selected"] = margin_bottom
        data["margin_left_selected"] = margin_left
        data["margin_right_selected"] = margin_right
        data["shared_fallback_font_size_used"] = (
            base_font_size is not None
            and global_font != base_font_size
        )
        data.update(paragraph_debug)
        data["rendered_bullet_count"] = len(data["rendered"])
        data["visible_count_expectation_met"] = (
            len(data["rendered"]) == data["render_target_count"]
            and data["available_paragraph_count"] >= len(data["rendered"])
        )
        del data["shape"]
    return prepared


def _fit_bullet_shape_text(shape: Any, bullets: list[str], *, max_lines: int = 4) -> dict[str, Any]:
    final_bullets = [bullet for bullet in bullets if _safe_text(bullet)][:max_lines]
    font_fallback = None
    if sum(len(bullet) for bullet in final_bullets) > 360:
        dropped = final_bullets[-1:]
        final_bullets = final_bullets[:-1]
    else:
        dropped = []
    if sum(len(bullet) for bullet in final_bullets) > 300:
        font_fallback = 10
        shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_fallback)
    _replace_bullet_shape_text(shape, final_bullets)
    return {
        "rendered": final_bullets,
        "dropped": dropped,
        "font_fallback": font_fallback,
    }


def _apply_slide5_bullet_rhythm(shape: Any) -> None:
    _apply_bullet_spacing_and_font(
        shape,
        11,
        line_spacing=0.86,
        space_after=0,
    )


def _slide3_side_shapes(prs: Any, slide: Any) -> dict[str, dict[str, Any]]:
    midpoint = int(prs.slide_width) // 2
    text_shapes = [
        shape
        for shape in _walk_shapes(slide.shapes)
        if getattr(shape, "has_text_frame", False)
    ]
    labels = [
        shape
        for shape in text_shapes
        if int(getattr(shape, "top", 0) or 0) >= Inches(1.8)
        and int(getattr(shape, "top", 0) or 0) <= Inches(2.6)
        and ("“" in _shape_text(shape) or '"' in _shape_text(shape))
    ]
    bullets = [
        shape
        for shape in text_shapes
        if int(getattr(shape, "top", 0) or 0) >= Inches(2.6)
        and int(getattr(shape, "height", 0) or 0) >= Inches(2.0)
        and len(getattr(shape.text_frame, "paragraphs", []) or []) >= 5
    ]
    pictures = [
        shape
        for shape in slide.shapes
        if hasattr(shape, "image") or "picture" in str(shape.shape_type).lower()
    ]
    pictures = [
        shape
        for shape in pictures
        if int(getattr(shape, "top", 0) or 0) >= Inches(2.3)
        and int(getattr(shape, "top", 0) or 0) <= Inches(5.8)
    ]
    pictures = sorted(pictures, key=lambda shape: int(getattr(shape, "left", 0) or 0))
    return {
        "current": {
            "label": next((shape for shape in labels if int(shape.left) < midpoint), None),
            "bullets": next((shape for shape in bullets if int(shape.left) < midpoint), None),
            "picture": pictures[0] if len(pictures) >= 1 else None,
        },
        "benchmark": {
            "label": next((shape for shape in labels if int(shape.left) >= midpoint), None),
            "bullets": next((shape for shape in bullets if int(shape.left) >= midpoint), None),
            "picture": pictures[1] if len(pictures) >= 2 else None,
        },
    }


def _apply_slide3_search_benchmark(prs: Any, slide: Any, payload: dict[str, Any]) -> None:
    if not payload:
        return
    current = payload.get("current")
    benchmark = payload.get("benchmark")
    if not isinstance(current, dict) and not isinstance(benchmark, dict):
        print(
            "[audit_powerpoint_new] Slide 3 search evidence was unavailable; "
            "the slide was left unchanged."
        )
        return

    intro_shape = next(
        (
            shape
            for shape in _walk_shapes(slide.shapes)
            if getattr(shape, "has_text_frame", False)
            and "walmart search results within the" in _shape_text(shape).lower()
        ),
        None,
    )
    if intro_shape is not None:
        text = _shape_text(intro_shape)
        current_phrase = _safe_text((current or {}).get("category_phrase")) or "category search"
        benchmark_phrase = _safe_text((benchmark or {}).get("category_phrase")) or "category search"
        updated = text.replace("baby bath", current_phrase).replace(
            "clean baby care",
            benchmark_phrase,
        )
        _replace_shape_text_preserve_style(intro_shape, updated)
    else:
        print("[audit_powerpoint_new] Slide 3 intro shape was not found.")

    shapes = _slide3_side_shapes(prs, slide)
    bullet_render_items: list[tuple[str, Any, list[str]]] = []
    for side in ("current", "benchmark"):
        side_payload = payload.get(side)
        if not isinstance(side_payload, dict) or not side_payload.get("source_row"):
            print(
                f"[audit_powerpoint_new] Slide 3 {side} evidence was unavailable; "
                "that side was left unchanged."
            )
            continue
        side_shapes = shapes.get(side, {})
        label = side_shapes.get("label")
        if label is not None:
            term = _safe_text(side_payload.get("search_term")) or "category search"
            _replace_shape_text_preserve_style(label, f"“{term}”")
        else:
            print(f"[audit_powerpoint_new] Slide 3 {side} search label was not found.")

        bullets = [
            _safe_text(value)
            for value in (side_payload.get("bullets", []) or [])
            if _safe_text(value)
        ]
        bullet_shape = side_shapes.get("bullets")
        if bullet_shape is not None and len(bullets) == 4:
            bullet_render_items.append((side, bullet_shape, bullets))
        elif bullet_shape is None:
            print(f"[audit_powerpoint_new] Slide 3 {side} bullet box was not found.")
        else:
            print(
                f"[audit_powerpoint_new] Slide 3 {side} did not contain exactly four bullets; "
                "the template bullet box was preserved."
            )

        picture = side_shapes.get("picture")
        image_bytes = _decode_data_image(side_payload.get("screenshot"))
        if picture is None:
            print(f"[audit_powerpoint_new] Slide 3 {side} screenshot placeholder was not found.")
        elif image_bytes is None:
            print(
                f"[audit_powerpoint_new] Slide 3 {side} screenshot was invalid; "
                "the template picture was preserved."
            )
        else:
            bounds = (
                int(picture.left),
                int(picture.top),
                int(picture.width),
                int(picture.height),
            )
            inserted = _add_contained_picture(
                slide,
                left=bounds[0],
                top=bounds[1],
                width=bounds[2],
                height=bounds[3],
                image_bytes=image_bytes,
            )
            if inserted is None:
                print(
                    f"[audit_powerpoint_new] Slide 3 {side} screenshot could not be inserted; "
                    "the template picture was preserved."
                )
            else:
                _remove_shape(picture)
    if bullet_render_items:
        payload.setdefault("debug", {})["render_fit"] = _fit_bullet_shape_group(
            bullet_render_items,
            max_lines=4,
            drop_threshold=260,
            font_threshold=235,
            base_font_size=12,
            fallback_font_size=11,
        )


def _fit_slide4_label(shape: Any) -> None:
    text_frame = shape.text_frame
    text_frame.word_wrap = False
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    label_text = _shape_text(shape)
    if len(label_text) > 12:
        font_size = Pt(10 if len(label_text) > 20 else 12)
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = font_size


def _slide4_picture_containers(prs: Any, slide: Any) -> list[Any]:
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    pictures = [
        shape
        for shape in slide.shapes
        if hasattr(shape, "image") or "picture" in str(shape.shape_type).lower()
    ]
    on_slide = [
        shape
        for shape in pictures
        if int(shape.left) < slide_width
        and int(shape.top) < slide_height
        and int(shape.left + shape.width) > 0
        and int(shape.top + shape.height) > 0
    ]
    return sorted(on_slide, key=lambda shape: int(shape.left))[:3]


def _populate_slide4_carousel(
    slide: Any,
    *,
    container: tuple[int, int, int, int],
    images: list[dict[str, Any]],
) -> None:
    usable_images = [image for image in images[:12] if _safe_text(image.get("url"))]
    if not usable_images:
        return
    left, top, width, height = container
    columns = 3
    rows = max(1, math.ceil(len(usable_images) / columns))
    gutter = Inches(0.06)
    cell_width = int((width - gutter * (columns - 1)) / columns)
    cell_height = int((height - gutter * (rows - 1)) / rows)
    for position, image in enumerate(usable_images):
        row, column = divmod(position, columns)
        cell_left = int(left + column * (cell_width + gutter))
        cell_top = int(top + row * (cell_height + gutter))
        image_bytes = _load_image_from_url(_safe_text(image.get("url")))
        if image_bytes is None:
            continue
        _add_contained_picture(
            slide,
            left=cell_left,
            top=cell_top,
            width=cell_width,
            height=cell_height,
            image_bytes=image_bytes,
        )


def _slide4_line_shapes(slide: Any) -> list[Any]:
    return [
        shape
        for shape in _walk_shapes(slide.shapes)
        if "LINE" in str(getattr(shape, "shape_type", "")).upper()
    ]


def _cleanup_slide4_two_column_chrome(
    slide: Any,
    *,
    target_indices: list[int],
    inactive_indices: list[int],
    two_column_bounds: dict[int, dict[str, tuple[int, int, int, int]]],
) -> dict[str, Any]:
    removed: list[dict[str, Any]] = []
    relocated: list[dict[str, Any]] = []
    lines = _slide4_line_shapes(slide)
    vertical_lines = [
        shape
        for shape in lines
        if int(getattr(shape, "width", 0) or 0) <= int(Inches(0.03))
        and int(getattr(shape, "height", 0) or 0) >= int(Inches(1.0))
    ]
    for shape in vertical_lines:
        removed.append({"name": getattr(shape, "name", ""), "reason": "removed_3_column_vertical_divider"})
        _remove_shape(shape)

    horizontal_lines = sorted(
        [
            shape
            for shape in lines
            if shape not in vertical_lines
            and int(getattr(shape, "width", 0) or 0) >= int(Inches(1.0))
            and int(getattr(shape, "height", 0) or 0) <= int(Inches(0.03))
        ],
        key=lambda shape: int(getattr(shape, "left", 0) or 0),
    )
    for index in inactive_indices:
        if index < len(horizontal_lines):
            shape = horizontal_lines[index]
            removed.append({"name": getattr(shape, "name", ""), "reason": "removed_inactive_column_label_rule"})
            _remove_shape(shape)

    for index in target_indices:
        if index >= len(horizontal_lines) or index not in two_column_bounds:
            continue
        shape = horizontal_lines[index]
        if getattr(shape, "element", None) is None or shape.element.getparent() is None:
            continue
        label_left, label_top, label_width, label_height = two_column_bounds[index]["label"]
        bullet_left, _bullet_top, bullet_width, _bullet_height = two_column_bounds[index]["bullet"]
        shape.left = int(min(label_left, bullet_left))
        shape.top = int(label_top + label_height)
        shape.width = int(max(label_width, bullet_width))
        shape.height = 0
        relocated.append(
            {
                "name": getattr(shape, "name", ""),
                "target_index": index,
                "bounds": [int(shape.left), int(shape.top), int(shape.width), int(shape.height)],
            }
        )

    return {
        "removed": removed,
        "relocated_label_rules": relocated,
        "suppressed_vertical_dividers": len(vertical_lines),
        "inactive_indices": inactive_indices,
    }


def _apply_slide4_content(prs: Any, slide: Any, payload: dict[str, Any]) -> None:
    columns = list(payload.get("columns", []) or [])[:3]
    while len(columns) < 3:
        columns.append({"label": f"Competitor {len(columns)}", "ordered_images": [], "bullets": []})
    active_indices = [index for index, column in enumerate(columns) if column.get("active", True)]
    layout_mode = payload.get("layout_mode") or ("2-column" if len(active_indices) == 2 else "3-column")
    layout_debug = payload.setdefault("debug", {}).setdefault("layout_debug", {})

    for shape in _walk_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        text = _shape_text(shape)
        if "Best-in-class Walmart PDPs" in text:
            _replace_shape_text_preserve_style(
                shape,
                text.replace("Best-in-class Walmart PDPs", "Strong Walmart PDP patterns"),
            )
            break

    labels = _slide4_label_shapes(slide)
    bullets = _slide4_bullet_shapes(slide)
    containers = _slide4_picture_containers(prs, slide)
    if len(labels) < 3 or len(bullets) < 3 or len(containers) < 3:
        print(
            "[audit_powerpoint_new] Slide 4 template shapes incomplete; "
            "leaving the slide without floating replacement content."
        )
        return

    container_bounds = [
        (int(shape.left), int(shape.top), int(shape.width), int(shape.height))
        for shape in containers
    ]
    label_bounds = [
        (int(shape.left), int(shape.top), int(shape.width), int(shape.height))
        for shape in labels
    ]
    bullet_bounds = [
        (int(shape.left), int(shape.top), int(shape.width), int(shape.height))
        for shape in bullets
    ]
    for shape in containers:
        _remove_shape(shape)
    for shape in list(slide.shapes):
        if not (hasattr(shape, "image") or "picture" in str(shape.shape_type).lower()):
            continue
        if (
            int(shape.left) < 0
            or int(shape.top) < 0
            or int(shape.left + shape.width) > int(prs.slide_width)
            or int(shape.top + shape.height) > int(prs.slide_height)
        ):
            _remove_shape(shape)

    if layout_mode == "2-column" and len(active_indices) == 2:
        target_indices = [0, 2]
        inactive_indices = [index for index in range(3) if index not in target_indices]
        for inactive_index in inactive_indices:
            _remove_shape(labels[inactive_index])
            _remove_shape(bullets[inactive_index])
        bullet_width = max(
            int(bullet_bounds[0][2]),
            int(bullet_bounds[2][2]),
            int(Inches(4.05)),
        )
        gutter = int(Inches(0.55))
        total_width = (bullet_width * 2) + gutter
        start_left = int((int(prs.slide_width) - total_width) / 2)
        image_width = max(int(container_bounds[0][2]), int(container_bounds[2][2]))
        label_width = max(int(label_bounds[0][2]), int(label_bounds[2][2]))
        two_column_bounds: dict[int, dict[str, tuple[int, int, int, int]]] = {}
        for slot, target_index in enumerate(target_indices):
            column_left = start_left + slot * (bullet_width + gutter)
            image_left = int(column_left + (bullet_width - image_width) / 2)
            label_left = int(column_left + (bullet_width - label_width) / 2)
            two_column_bounds[target_index] = {
                "label": (label_left, label_bounds[target_index][1], label_width, label_bounds[target_index][3]),
                "bullet": (column_left, bullet_bounds[target_index][1], bullet_width, bullet_bounds[target_index][3]),
                "container": (image_left, container_bounds[target_index][1], image_width, container_bounds[target_index][3]),
            }
        cleanup_debug = _cleanup_slide4_two_column_chrome(
            slide,
            target_indices=target_indices,
            inactive_indices=inactive_indices,
            two_column_bounds=two_column_bounds,
        )
        layout_debug["render_cleanup"] = cleanup_debug
        layout_debug["final_active_geometry"] = {
            str(index): {key: list(value) for key, value in bounds.items()}
            for index, bounds in two_column_bounds.items()
        }
        render_pairs = list(zip(target_indices, [columns[index] for index in active_indices]))
    else:
        two_column_bounds = {}
        layout_debug["render_cleanup"] = {
            "removed": [],
            "relocated_label_rules": [],
            "suppressed_vertical_dividers": 0,
            "inactive_indices": [],
        }
        layout_debug["final_active_geometry"] = {
            str(index): {
                "label": list(label_bounds[index]),
                "bullet": list(bullet_bounds[index]),
                "container": list(container_bounds[index]),
            }
            for index in range(min(3, len(columns)))
        }
        render_pairs = [(index, column) for index, column in enumerate(columns)]

    shared_bullet_top = min(
        int((two_column_bounds.get(index, {}) or {}).get("bullet", bullet_bounds[index])[1])
        if two_column_bounds
        else int(bullet_bounds[index][1])
        for index, _column in render_pairs
    )
    layout_debug["shared_bullet_block_top"] = shared_bullet_top
    slide4_bullet_render_items: list[tuple[str, Any, list[str]]] = []
    for target_index, column in render_pairs:
        if not column.get("active", True):
            _replace_shape_text_preserve_style(labels[target_index], "")
            _replace_bullet_shape_text(bullets[target_index], [])
            continue
        bounds = two_column_bounds.get(target_index)
        labels[target_index].left, labels[target_index].top, labels[target_index].width, labels[target_index].height = (
            bounds["label"] if bounds else label_bounds[target_index]
        )
        bullets[target_index].left, bullets[target_index].top, bullets[target_index].width, bullets[target_index].height = (
            bounds["bullet"] if bounds else bullet_bounds[target_index]
        )
        bullets[target_index].top = shared_bullet_top
        _replace_shape_text_preserve_style(labels[target_index], column.get("label", ""))
        _fit_slide4_label(labels[target_index])
        slide4_bullet_render_items.append(
            (
                str(target_index),
                bullets[target_index],
                list(column.get("bullets", []) or []),
            )
        )
        left, top, width, height = bounds["container"] if bounds else container_bounds[target_index]
        bullet_top = int(bullets[target_index].top)
        height = max(1, min(height, bullet_top - Inches(0.05) - top))
        _populate_slide4_carousel(
            slide,
            container=(left, top, width, height),
            images=list(column.get("ordered_images", []) or []),
        )
    if slide4_bullet_render_items:
        layout_debug["bullet_render_fit"] = _fit_bullet_shape_group(
            slide4_bullet_render_items,
            max_lines=6,
            drop_threshold=520,
            font_threshold=470,
            base_font_size=11,
            fallback_font_size=10,
            ensure_paragraph_count=True,
            line_spacing=0.9,
            space_after=2,
            margin_top=0,
            margin_bottom=0,
            margin_left=0,
            margin_right=0,
        )


def _decode_data_image(value: Any) -> bytes | None:
    data_url = _safe_text(value)
    if not data_url.lower().startswith("data:image/") or "," not in data_url:
        return None
    header, encoded = data_url.split(",", 1)
    if ";base64" not in header.lower():
        return None
    try:
        image_bytes = base64.b64decode(encoded, validate=True)
        with Image.open(io.BytesIO(image_bytes)) as image:
            image.verify()
        return image_bytes
    except Exception as exc:
        print(f"[audit_powerpoint_new] Slide 5 screenshot decode failed: {exc}")
        return None


def _slide5_side_shapes(prs: Any, slide: Any) -> dict[str, dict[str, Any]]:
    midpoint = int(prs.slide_width) // 2
    pictures = sorted(
        [
            shape
            for shape in slide.shapes
            if (hasattr(shape, "image") or "picture" in str(shape.shape_type).lower())
            and int(getattr(shape, "top", 0) or 0) >= Inches(2.5)
        ],
        key=lambda shape: int(shape.left),
    )
    bullet_shapes = sorted(
        [
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and int(getattr(shape, "top", 0) or 0) >= Inches(2.5)
            and int(getattr(shape, "height", 0) or 0) >= Inches(2.0)
            and len(shape.text_frame.paragraphs) >= 6
        ],
        key=lambda shape: int(shape.left),
    )
    left_header = next(
        (
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and _shape_text(shape) == "Current Structure"
        ),
        None,
    )
    right_header = next(
        (
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and _shape_text(shape) == "Competitive Benchmark"
        ),
        None,
    )
    lines = sorted(
        [
            shape
            for shape in slide.shapes
            if "LINE" in str(getattr(shape, "shape_type", "")).upper()
            and int(getattr(shape, "top", 0) or 0) >= Inches(2.0)
            and int(getattr(shape, "top", 0) or 0) <= Inches(3.0)
        ],
        key=lambda shape: int(shape.left),
    )
    return {
        "client": {
            "picture": next(
                (shape for shape in pictures if int(shape.left) < midpoint),
                None,
            ),
            "bullets": next(
                (shape for shape in bullet_shapes if int(shape.left) < midpoint),
                None,
            ),
            "header": left_header,
            "divider": lines[0] if lines else None,
        },
        "competitor": {
            "picture": next(
                (shape for shape in pictures if int(shape.left) >= midpoint),
                None,
            ),
            "bullets": next(
                (shape for shape in bullet_shapes if int(shape.left) >= midpoint),
                None,
            ),
            "header": right_header,
            "divider": lines[-1] if lines else None,
        },
    }


def _apply_slide5_no_brand_shop(
    prs: Any,
    slide: Any,
    payload: dict[str, Any],
    shapes: dict[str, dict[str, Any]],
) -> None:
    competitor = payload.get("competitor")
    if not isinstance(competitor, dict):
        print(
            "[audit_powerpoint_new] Slide 5 No Brand Shop mode lacked valid "
            "Competitor evidence; the slide was left unchanged."
        )
        return
    competitor_picture = shapes["competitor"].get("picture")
    competitor_bullets = shapes["competitor"].get("bullets")
    competitor_header = shapes["competitor"].get("header")
    competitor_divider = shapes["competitor"].get("divider")
    if any(
        shape is None
        for shape in (
            competitor_picture,
            competitor_bullets,
            competitor_header,
            competitor_divider,
        )
    ):
        print(
            "[audit_powerpoint_new] Slide 5 No Brand Shop template shapes were "
            "not resolved; the slide was left unchanged."
        )
        return
    image_bytes = _decode_data_image(competitor.get("screenshot"))
    bullets = [
        _safe_text(value)
        for value in (competitor.get("bullets", []) or [])
        if _safe_text(value)
    ]
    if image_bytes is None or len(bullets) < 1:
        print(
            "[audit_powerpoint_new] Slide 5 No Brand Shop payload was incomplete; "
            "the slide was left unchanged."
        )
        return

    inserted = _add_contained_picture(
        slide,
        left=Inches(1.4),
        top=Inches(2.88),
        width=Inches(7.4),
        height=Inches(3.77),
        image_bytes=image_bytes,
    )
    if inserted is None:
        print(
            "[audit_powerpoint_new] Slide 5 No Brand Shop screenshot could not be "
            "inserted; the slide was left unchanged."
        )
        return

    for key in ("header", "divider", "picture", "bullets"):
        shape = shapes["client"].get(key)
        if shape is not None:
            _remove_shape(shape)

    competitor_header.left = Inches(3.45)
    competitor_header.width = Inches(3.2)
    competitor_divider.left = Inches(1.4)
    competitor_divider.width = Inches(11.2)

    _remove_shape(competitor_picture)
    payload.setdefault("debug", {})["render_fit"] = _fit_bullet_shape_group(
        [("competitor", competitor_bullets, bullets)],
        max_lines=7,
        drop_threshold=640,
        font_threshold=590,
        base_font_size=14,
        fallback_font_size=14,
        ensure_paragraph_count=True,
        font_name="Raleway",
        line_spacing=0.9,
        space_after=1,
    )


def _apply_slide5_brand_shop(prs: Any, slide: Any, payload: dict[str, Any]) -> None:
    if not payload:
        return
    shapes = _slide5_side_shapes(prs, slide)
    if payload.get("mode") == "no_brand_shop":
        _apply_slide5_no_brand_shop(prs, slide, payload, shapes)
        return
    bullet_render_items: list[tuple[str, Any, list[str]]] = []
    for side in ("client", "competitor"):
        side_payload = payload.get(side)
        if not isinstance(side_payload, dict):
            print(
                f"[audit_powerpoint_new] Slide 5 {side} evidence was unavailable; "
                "the template side was left unchanged."
            )
            continue
        picture = shapes[side].get("picture")
        bullet_shape = shapes[side].get("bullets")
        if picture is None or bullet_shape is None:
            print(
                f"[audit_powerpoint_new] Slide 5 {side} template shapes were not resolved; "
                "the template side was left unchanged."
            )
            continue
        image_bytes = _decode_data_image(side_payload.get("screenshot"))
        bullets = [
            _safe_text(value)
            for value in (side_payload.get("bullets", []) or [])
            if _safe_text(value)
        ]
        if image_bytes is None:
            print(
                f"[audit_powerpoint_new] Slide 5 {side} screenshot was invalid; "
                "the template picture was preserved."
            )
        else:
            bounds = (
                int(picture.left),
                int(picture.top),
                int(picture.width),
                int(picture.height),
            )
            inserted = _add_contained_picture(
                slide,
                left=bounds[0],
                top=bounds[1],
                width=bounds[2],
                height=bounds[3],
                image_bytes=image_bytes,
            )
            if inserted is None:
                print(
                    f"[audit_powerpoint_new] Slide 5 {side} screenshot could not be "
                    "inserted; the template picture was preserved."
                )
            else:
                _remove_shape(picture)
        if bullets:
            bullet_render_items.append((side, bullet_shape, bullets))
        else:
            print(
                f"[audit_powerpoint_new] Slide 5 {side} did not contain bullets; "
                "the template bullet box was preserved."
            )
    if bullet_render_items:
        payload.setdefault("debug", {})["render_fit"] = _fit_bullet_shape_group(
            bullet_render_items,
            max_lines=7,
            drop_threshold=640,
            font_threshold=590,
            base_font_size=14,
            fallback_font_size=14,
            ensure_paragraph_count=True,
            font_name="Raleway",
            line_spacing=0.9,
            space_after=1,
        )


def _slide2_text_shapes(slide: Any) -> list[Any]:
    return [
        shape
        for shape in _walk_shapes(slide.shapes)
        if getattr(shape, "has_text_frame", False) and _shape_text(shape)
    ]


def _shape_contains_any(shape: Any, tokens: tuple[str, ...]) -> bool:
    text = _shape_text(shape).lower()
    return any(token.lower() in text for token in tokens)


def _apply_slide2_summary(slide: Any, payload: dict[str, Any]) -> None:
    if not payload:
        return
    shapes = _slide2_text_shapes(slide)
    phrases = payload.get("phrases", {}) or {}
    category_phrase = _safe_text(phrases.get("category_phrase"))
    left_replaced = False
    if category_phrase:
        intro_shape = next(
            (
                shape
                for shape in shapes
                if int(getattr(shape, "left", 0) or 0) < Inches(5.0)
                and _shape_contains_any(
                    shape,
                    (
                        "baby care and clean lifestyle categories",
                    ),
                )
            ),
            None,
        )
        if intro_shape is not None:
            text = _shape_text(intro_shape)
            updated = text.replace(
                "baby care and clean lifestyle categories",
                category_phrase,
            )
            if updated != text:
                _replace_shape_text_preserve_style(intro_shape, updated)
                left_replaced = True
    payload.setdefault("debug", {})["left_category_phrase_replacement_succeeded"] = left_replaced
    if category_phrase and not left_replaced:
        print(
            "[audit_powerpoint_new] Slide 2 category phrase anchor was not found; "
            "left paragraph was left unchanged."
        )

    sections = payload.get("sections", {}) or {}
    rating_shapes = sorted(
        [
            shape
            for shape in shapes
            if _shape_text(shape) in {"Strong", "Significant", "Evolving", "Emerging", "Limited", "Meaningful", "Selective", "Competitive"}
            and int(getattr(shape, "left", 0) or 0) >= Inches(9.0)
        ],
        key=lambda shape: int(getattr(shape, "top", 0) or 0),
    )
    section_order = ("consumer_demand", "walmart_opportunity", "competitive_benchmark")
    for index, section_key in enumerate(section_order):
        if index < len(rating_shapes):
            rating = _safe_text((sections.get(section_key, {}) or {}).get("rating"))
            if rating:
                _replace_shape_text_preserve_style(rating_shapes[index], rating)

    bullet_shapes = sorted(
        [
            shape
            for shape in shapes
            if int(getattr(shape, "left", 0) or 0) >= Inches(5.0)
            and int(getattr(shape, "top", 0) or 0) >= Inches(1.45)
            and int(getattr(shape, "height", 0) or 0) >= Inches(0.7)
            and not _shape_contains_any(
                shape,
                (
                    "consumer demand",
                    "walmart opportunity",
                    "competitive benchmark",
                    "walmart ecommerce opportunity",
                ),
            )
            and _shape_text(shape) not in {"Strong", "Significant", "Evolving", "Emerging", "Limited", "Meaningful", "Selective", "Competitive"}
        ],
        key=lambda shape: int(getattr(shape, "top", 0) or 0),
    )[:3]
    bullet_render_items: list[tuple[str, Any, list[str]]] = []
    for index, section_key in enumerate(section_order):
        if index >= len(bullet_shapes):
            continue
        bullets = [
            _safe_text(bullet)
            for bullet in ((sections.get(section_key, {}) or {}).get("bullets", []) or [])
            if _safe_text(bullet)
        ]
        bullet_render_items.append((section_key, bullet_shapes[index], bullets[:4]))
    if bullet_render_items:
        payload.setdefault("debug", {})["render_fit"] = _fit_bullet_shape_group(
            bullet_render_items,
            max_lines=4,
            drop_threshold=250,
            font_threshold=220,
            base_font_size=12,
            fallback_font_size=11,
            ensure_paragraph_count=True,
        )


def _replace_cell_text_preserve_style(cell: Any, text: str) -> None:
    text_frame = getattr(cell, "text_frame", None)
    if text_frame is None or not text_frame.paragraphs:
        cell.text = _safe_text(text)
        return
    _replace_paragraph_text_preserve_style(text_frame.paragraphs[0], text)
    for paragraph in text_frame.paragraphs[1:]:
        _replace_paragraph_text_preserve_style(paragraph, "")


def _slide6_display_label(payload: dict[str, Any], value: Any, context: str) -> str:
    label = _safe_text(value)
    if label in {"Strong", "Moderate", "Partial", "Limited"}:
        return label
    warning = (
        f"{context} produced invalid visibility '{label or '<empty>'}'; "
        "replaced with Limited."
    )
    payload.setdefault("warnings", []).append(warning)
    print(f"[audit_powerpoint_new] Slide 6 warning: {warning}")
    return "Limited"


def _fit_slide6_client_header(cell: Any, client_label: str) -> None:
    if len(client_label) <= 18:
        return
    text_frame = cell.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    font_size = 12 if len(client_label) <= 28 else 10 if len(client_label) <= 40 else 8
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)


def _apply_slide6_visibility(slide: Any, payload: dict[str, Any]) -> None:
    if not payload:
        return
    segments = list(payload.get("segments", []) or [])
    if len(segments) != 6:
        print(
            "[audit_powerpoint_new] Slide 6 visibility payload did not contain "
            "exactly six segments; table replacement was skipped."
        )
        return

    table_shape = next(
        (
            shape
            for shape in _walk_shapes(slide.shapes)
            if getattr(shape, "has_table", False)
            and len(shape.table.rows) >= 7
            and len(shape.table.columns) >= 3
            and "search segment" in _safe_text(shape.table.cell(0, 0).text).lower()
        ),
        None,
    )
    if table_shape is None:
        print(
            "[audit_powerpoint_new] Slide 6 search-segment table was not found; "
            "visibility rows were skipped."
        )
    else:
        table = table_shape.table
        client_label = _safe_text(payload.get("client_label")) or "Client"
        _replace_cell_text_preserve_style(
            table.cell(0, 2),
            client_label,
        )
        _fit_slide6_client_header(table.cell(0, 2), client_label)
        for row_index, segment in enumerate(segments, start=1):
            _replace_cell_text_preserve_style(table.cell(row_index, 0), segment.get("segment", ""))
            _replace_cell_text_preserve_style(
                table.cell(row_index, 1),
                _slide6_display_label(
                    payload,
                    segment.get("competitor_visibility"),
                    f"row {row_index} competitor",
                ),
            )
            _replace_cell_text_preserve_style(
                table.cell(row_index, 2),
                _slide6_display_label(
                    payload,
                    segment.get("client_visibility"),
                    f"row {row_index} client",
                ),
            )

    text_shapes = [
        shape
        for shape in _walk_shapes(slide.shapes)
        if getattr(shape, "has_text_frame", False)
        and _shape_text(shape).strip().lower() != "digital shelf ownership"
    ]
    intro_shape = next(
        (
            shape
            for shape in text_shapes
            if "competitors currently own more walmart search paths" in _shape_text(shape).lower()
        ),
        None,
    )
    if intro_shape is None:
        intro_shape = min(
            (shape for shape in text_shapes if int(getattr(shape, "top", 0) or 0) < Inches(2.5)),
            key=lambda shape: int(getattr(shape, "top", 0) or 0),
            default=None,
        )
    if intro_shape is not None:
        _replace_shape_text_preserve_style(intro_shape, payload.get("intro", ""))
    else:
        print("[audit_powerpoint_new] Slide 6 intro shape was not found.")

    takeaway_shape = next(
        (
            shape
            for shape in text_shapes
            if "walmart shoppers are highly search-driven" in _shape_text(shape).lower()
        ),
        None,
    )
    if takeaway_shape is None:
        takeaway_shape = max(
            text_shapes,
            key=lambda shape: int(getattr(shape, "top", 0) or 0),
            default=None,
        )
    if takeaway_shape is not None:
        _replace_shape_text_preserve_style(takeaway_shape, payload.get("takeaway", ""))
    else:
        print("[audit_powerpoint_new] Slide 6 takeaway shape was not found.")


def generate_new_audit_powerpoint_from_template(
    *, export_plan: dict, template_path: str, include_slide_9: bool = False, competitor_records: list[dict[str, Any]] | None = None
) -> bytes:
    prs = Presentation(template_path)
    metadata = export_plan.get("audit_metadata", {}) or {}
    client_name = _safe_text(metadata.get("client_name", ""))
    client_company_name = _safe_text(metadata.get("client_company_name", "")) or client_name
    audit_date = _format_cover_date(_safe_text(metadata.get("audit_date", "")))
    
    # Handle Slide 9 removal if not included
    if not include_slide_9:
        _remove_slide_by_title(prs, "Walmart Cash Program Visibility")
    _remove_slide_by_title(prs, "If they only have bandwidth for 5 things:")
    
    # Existing template strings are the replacement targets. Do not require manual placeholders.
    template_replacements = {
        "The Honest Company": client_company_name,
        "May 27, 2026": audit_date,
        "Honest": client_name,
    }
    replace_existing_template_text(prs, template_replacements)

    slide2 = _find_slide_by_title(prs, "Walmart eCommerce Opportunity")
    if slide2 is None:
        print(
            "[audit_powerpoint_new] Slide 2 title was not found; "
            "Walmart eCommerce Opportunity summary was skipped."
        )
    else:
        _apply_slide2_summary(slide2, export_plan.get("slide2_summary", {}) or {})

    slide3 = _find_slide_by_title(prs, "Search & Discoverability Benchmarking")
    if slide3 is None:
        print(
            "[audit_powerpoint_new] Slide 3 title was not found; "
            "Search & Discoverability Benchmarking was skipped."
        )
    else:
        _apply_slide3_search_benchmark(
            prs,
            slide3,
            export_plan.get("slide3_search_benchmark", {}) or {},
        )

    slide4 = _find_slide_by_title(prs, "PDP Content Benchmarking")
    if slide4 is None:
        print(
            "[audit_powerpoint_new] Slide 4 title was not found; "
            "PDP benchmarking was skipped."
        )
    else:
        slide4_payload = build_slide4_pdp_benchmark_payload(
            export_plan,
            competitor_records,
        )
        _apply_slide4_content(prs, slide4, slide4_payload)

    slide5 = _find_slide_by_title(prs, "Brand Shop Content Benchmarking")
    if slide5 is None:
        print(
            "[audit_powerpoint_new] Slide 5 title was not found; "
            "Brand Shop benchmarking was skipped."
        )
    else:
        _apply_slide5_brand_shop(
            prs,
            slide5,
            export_plan.get("slide5_brand_shop", {}) or {},
        )

    slide6 = _find_slide_by_title(prs, "Digital Shelf Ownership")
    if slide6 is None:
        print(
            "[audit_powerpoint_new] Slide 6 title was not found; "
            "digital shelf alignment was skipped."
        )
    else:
        _apply_slide6_visibility(
            slide6,
            export_plan.get("slide6_visibility", {}) or {},
        )

    clear_unresolved_placeholder_text(prs)
    
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue()
