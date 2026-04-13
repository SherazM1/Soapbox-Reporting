from __future__ import annotations

import copy
import io
import os
from datetime import date as dt_date, datetime
from urllib.request import Request, urlopen
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


def resolve_audit_template_path() -> str:
    candidates = [
        os.path.join("templates", "audit_template.pptx"),
        os.path.join("templates", "Audit Template.pptx"),
    ]
    for path in candidates:
        if os.path.exists(path):
            return path

    templates_dir = "templates"
    if os.path.isdir(templates_dir):
        for name in os.listdir(templates_dir):
            if name.lower().endswith(".pptx"):
                return os.path.join(templates_dir, name)
    raise FileNotFoundError("No .pptx template found in templates/.")


def _safe_text(value: Any) -> str:
    return str(value or "").strip()


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return _safe_text(getattr(shape, "text", ""))


def _find_slide_by_text(prs: Presentation, token: str) -> Any:
    token_l = token.lower()
    for slide in prs.slides:
        for shape in slide.shapes:
            text = _shape_text(shape).lower()
            if token_l in text:
                return slide
    return None


def _find_shape_contains(slide: Any, token: str) -> Any:
    token_l = token.lower()
    for shape in slide.shapes:
        text = _shape_text(shape).lower()
        if token_l in text:
            return shape
    return None


def _find_pdp_title_shape(slide: Any) -> Any:
    # First-pass template token.
    token_shape = _find_shape_contains(slide, "(Product Title)")
    if token_shape is not None:
        return token_shape

    # Duplicated PDP slides may no longer contain template tokens after first population.
    # In that case, resolve the same title block via the stable "Item ID:" marker.
    item_id_shape = _find_shape_contains(slide, "Item ID:")
    if item_id_shape is not None:
        return item_id_shape

    return None


def _set_shape_text(shape: Any, text: str, font_size: int | None = None) -> None:
    if not shape or not getattr(shape, "has_text_frame", False):
        return
    shape.text_frame.clear()
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = _safe_text(text)
    if font_size is not None:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)


def _replace_shape_text_preserve_style(shape: Any, text: str) -> None:
    """Replace text while preserving existing shape/run styling as much as possible."""
    if not shape or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        shape.text = _safe_text(text)
        return
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = _safe_text(text)
        for run in first_para.runs[1:]:
            run.text = ""
        for para in tf.paragraphs[1:]:
            for run in para.runs:
                run.text = ""
        return
    shape.text = _safe_text(text)


def _set_bullet_block(shape: Any, heading: str, bullets: list[str]) -> None:
    if not shape or not getattr(shape, "has_text_frame", False):
        return
    lines = [heading]
    lines.extend([f"- {_safe_text(b)}" for b in bullets if _safe_text(b)])
    _set_shape_text(shape, "\n".join(lines))


def _replace_paragraph_text_preserve_style(paragraph: Any, text: str) -> None:
    if paragraph.runs:
        paragraph.runs[0].text = _safe_text(text)
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = _safe_text(text)


def _force_paragraph_bullet(paragraph: Any) -> None:
    p = paragraph._p  # pylint: disable=protected-access
    pPr = p.get_or_add_pPr()
    # Remove numbered bullet settings if present.
    for child in list(pPr):
        if child.tag in {qn("a:buNone"), qn("a:buAutoNum")}:
            pPr.remove(child)
    if pPr.find(qn("a:buChar")) is None:
        bu_char = OxmlElement("a:buChar")
        bu_char.set("char", "\u2022")
        pPr.insert(0, bu_char)


def _force_paragraph_no_list(paragraph: Any) -> None:
    p = paragraph._p  # pylint: disable=protected-access
    pPr = p.get_or_add_pPr()
    for child in list(pPr):
        if child.tag in {qn("a:buChar"), qn("a:buAutoNum"), qn("a:buNone")}:
            pPr.remove(child)
    bu_none = OxmlElement("a:buNone")
    pPr.insert(0, bu_none)


def _force_paragraph_numbered(paragraph: Any, *, start_at: int = 1) -> None:
    p = paragraph._p  # pylint: disable=protected-access
    pPr = p.get_or_add_pPr()
    for child in list(pPr):
        if child.tag in {qn("a:buChar"), qn("a:buAutoNum"), qn("a:buNone")}:
            pPr.remove(child)
    bu_auto = OxmlElement("a:buAutoNum")
    bu_auto.set("type", "arabicPeriod")
    bu_auto.set("startAt", str(max(1, int(start_at))))
    pPr.insert(0, bu_auto)


def _set_pdp_title_and_item_id(shape: Any, title: str, item_id: str) -> None:
    if not shape or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        shape.text = _safe_text(title)
        return

    # Keep template title styling on first paragraph.
    _replace_paragraph_text_preserve_style(tf.paragraphs[0], title)

    item_text = f"Item ID: {item_id}" if _safe_text(item_id) else ""
    if len(tf.paragraphs) >= 2:
        _replace_paragraph_text_preserve_style(tf.paragraphs[1], item_text)
        for para in tf.paragraphs[2:]:
            _replace_paragraph_text_preserve_style(para, "")
    else:
        p = tf.add_paragraph()
        _replace_paragraph_text_preserve_style(p, item_text)
        p.level = 0


def _set_pdp_image_recommendations(shape: Any, heading: str, bullets: list[str]) -> None:
    if not shape or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame

    seen: set[str] = set()
    clean_bullets: list[str] = []
    for bullet in bullets or []:
        t = " ".join(_safe_text(bullet).split()).lstrip("-").strip()
        if not t:
            continue
        key = t.lower()
        if key in seen:
            continue
        seen.add(key)
        clean_bullets.append(t)
        if len(clean_bullets) >= 8:
            break

    if not tf.paragraphs:
        shape.text = _safe_text(heading)

    paragraphs = tf.paragraphs
    _replace_paragraph_text_preserve_style(paragraphs[0], heading)
    paragraphs[0].level = 0
    _force_paragraph_no_list(paragraphs[0])

    required_paragraphs = 1 + len(clean_bullets)
    while len(tf.paragraphs) < required_paragraphs:
        tf.add_paragraph()

    for idx, bullet in enumerate(clean_bullets, start=1):
        p = tf.paragraphs[idx]
        _replace_paragraph_text_preserve_style(p, bullet)
        p.level = 0
        _force_paragraph_numbered(p, start_at=idx)

    for p in tf.paragraphs[required_paragraphs:]:
        _replace_paragraph_text_preserve_style(p, "")
        _force_paragraph_no_list(p)


def _title_case_text(text: str) -> str:
    clean = _safe_text(text)
    if not clean:
        return ""

    def cap_token(tok: str) -> str:
        if not tok:
            return tok
        parts = tok.split("-")
        out_parts: list[str] = []
        for part in parts:
            if not part:
                out_parts.append(part)
                continue
            if "'" in part:
                sub = part.split("'")
                sub = [s[:1].upper() + s[1:].lower() if s else s for s in sub]
                out_parts.append("'".join(sub))
            else:
                out_parts.append(part[:1].upper() + part[1:].lower())
        return "-".join(out_parts)

    return " ".join(cap_token(t) for t in clean.split())


def _clean_bullets_for_slide(items: list[str], max_items: int = 8) -> list[str]:
    seen: set[str] = set()
    out: list[str] = []
    for item in items or []:
        txt = _safe_text(item).strip().lstrip("-").strip()
        txt = " ".join(txt.split())
        if not txt:
            continue
        key = txt.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(txt)
        if len(out) >= max_items:
            break
    return out


def _set_heading_and_real_bullets(shape: Any, heading: str, bullets: list[str]) -> None:
    """Populate existing placeholder with heading + real PPT bullet paragraphs."""
    if not shape or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    clean_bullets = _clean_bullets_for_slide(bullets, max_items=8)
    if not tf.paragraphs:
        shape.text = _safe_text(heading)

    _replace_paragraph_text_preserve_style(tf.paragraphs[0], heading)
    tf.paragraphs[0].level = 0
    _force_paragraph_no_list(tf.paragraphs[0])

    required_paragraphs = 1 + len(clean_bullets)
    while len(tf.paragraphs) < required_paragraphs:
        tf.add_paragraph()

    for idx, bullet in enumerate(clean_bullets, start=1):
        p = tf.paragraphs[idx]
        _replace_paragraph_text_preserve_style(p, bullet)
        p.level = 0
        _force_paragraph_bullet(p)
    _normalize_bullet_paragraphs(tf.paragraphs[1:required_paragraphs])

    for p in tf.paragraphs[required_paragraphs:]:
        _replace_paragraph_text_preserve_style(p, "")
        _force_paragraph_no_list(p)


def _normalize_bullet_paragraphs(paragraphs: list[Any]) -> None:
    if not paragraphs:
        return
    ref_para = paragraphs[1] if len(paragraphs) > 1 else paragraphs[0]
    ref_ppr = copy.deepcopy(ref_para._p.get_or_add_pPr())  # pylint: disable=protected-access
    ref_runs = list(ref_para.runs)
    ref_run = ref_runs[0] if ref_runs else None
    ref_run_rpr = None
    if ref_run is not None:
        try:
            ref_run_rpr = copy.deepcopy(ref_run._r.get_or_add_rPr())  # pylint: disable=protected-access
        except Exception:
            ref_run_rpr = None
    for para in paragraphs:
        # Keep all bullet paragraphs aligned/list-styled the same (including bullet 1).
        p = para._p  # pylint: disable=protected-access
        pPr = p.get_or_add_pPr()
        for child in list(pPr):
            pPr.remove(child)
        for child in list(ref_ppr):
            pPr.append(copy.deepcopy(child))
        for run in para.runs:
            if ref_run_rpr is not None:
                try:
                    r = run._r  # pylint: disable=protected-access
                    cur_rpr = r.rPr
                    if cur_rpr is not None:
                        r.remove(cur_rpr)
                    r.insert(0, copy.deepcopy(ref_run_rpr))
                    continue
                except Exception:
                    pass
            if ref_run is None:
                continue
            for attr in ("name", "size", "bold", "italic", "underline"):
                try:
                    setattr(run.font, attr, getattr(ref_run.font, attr))
                except Exception:
                    continue


def _emphasize_heading_paragraph(paragraph: Any, min_size_pt: int = 18) -> None:
    paragraph.level = 0
    for run in paragraph.runs:
        run.font.bold = True
        if run.font.size is not None:
            run.font.size = Pt(max(min_size_pt, int(run.font.size.pt)))


def _set_title_recommendation_block(shape: Any, heading: str, recommended_title: str) -> None:
    if not shape or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if not tf.paragraphs:
        shape.text = _safe_text(heading)

    _replace_paragraph_text_preserve_style(tf.paragraphs[0], heading)
    _force_paragraph_no_list(tf.paragraphs[0])
    _emphasize_heading_paragraph(tf.paragraphs[0], min_size_pt=18)

    bullet_text = _safe_text(recommended_title)
    required_paragraphs = 2
    while len(tf.paragraphs) < required_paragraphs:
        tf.add_paragraph()

    bullet_para = tf.paragraphs[1]
    _replace_paragraph_text_preserve_style(bullet_para, bullet_text)
    bullet_para.level = 0
    _force_paragraph_no_list(bullet_para)

    for p in tf.paragraphs[2:]:
        _replace_paragraph_text_preserve_style(p, "")
        _force_paragraph_no_list(p)


def _download_image_bytes(url: str) -> bytes | None:
    clean = _safe_text(url)
    if not clean:
        return None
    if clean.startswith("data:image"):
        try:
            header, payload = clean.split(",", 1)
            if ";base64" in header:
                import base64

                return base64.b64decode(payload)
        except Exception:
            return None

    try:
        req = Request(clean, headers={"User-Agent": "Mozilla/5.0"})
        with urlopen(req, timeout=12) as resp:
            return resp.read()
    except Exception:
        return None


def _insert_image_over_shape(slide: Any, target_shape: Any, image_url: str) -> bool:
    image_bytes = _download_image_bytes(image_url)
    if not image_bytes:
        return False
    try:
        left = target_shape.left
        top = target_shape.top
        width = target_shape.width
        height = target_shape.height
        slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width=width, height=height)
        return True
    except Exception:
        return False


def _insert_image_fit_within_shape(
    slide: Any, target_shape: Any, image_url: str, inset_ratio: float = 0.0
) -> bool:
    image_bytes = _download_image_bytes(image_url)
    if not image_bytes:
        return False
    try:
        left = int(target_shape.left)
        top = int(target_shape.top)
        width = int(target_shape.width)
        height = int(target_shape.height)
        inset = max(0.0, min(0.25, float(inset_ratio)))
        inset_x = int(width * inset)
        inset_y = int(height * inset)
        left += inset_x
        top += inset_y
        width = max(1, width - 2 * inset_x)
        height = max(1, height - 2 * inset_y)

        pic = slide.shapes.add_picture(io.BytesIO(image_bytes), 0, 0)
        src_w = max(1, int(pic.width))
        src_h = max(1, int(pic.height))
        scale = min(width / src_w, height / src_h)
        new_w = max(1, int(src_w * scale))
        new_h = max(1, int(src_h * scale))

        pic.width = new_w
        pic.height = new_h
        pic.left = left + int((width - new_w) / 2)
        pic.top = top + int((height - new_h) / 2)
        return True
    except Exception:
        return False


def _remove_shape_box_treatment(shape: Any) -> None:
    if shape is None:
        return
    try:
        if getattr(shape, "has_line_format", False):
            shape.line.fill.background()
            shape.line.width = 0
    except Exception:
        pass
    # Some template shapes keep an outline at the XML layer; remove it explicitly.
    try:
        sp_pr = shape.element.spPr
        if sp_pr is not None:
            ln = sp_pr.find(qn("a:ln"))
            if ln is not None:
                sp_pr.remove(ln)
            # Remove visual effects that can look like a frame/shadow around the shape.
            for tag in ("a:effectLst", "a:effectDag", "a:scene3d", "a:sp3d"):
                node = sp_pr.find(qn(tag))
                if node is not None:
                    sp_pr.remove(node)
        # Some templates inherit a visible line from p:style/a:lnRef; clear it too.
        style = shape.element.find(qn("p:style"))
        if style is not None:
            ln_ref = style.find(qn("a:lnRef"))
            if ln_ref is not None:
                style.remove(ln_ref)
    except Exception:
        pass
    try:
        if getattr(shape, "fill", None) is not None:
            shape.fill.background()
    except Exception:
        pass


def _suppress_pdp_image_placeholders(slide: Any) -> None:
    tokens = ("(insert pdp image here)", "insert pdp image here")
    for shape in list(slide.shapes):
        text = _shape_text(shape).lower()
        if not any(token in text for token in tokens):
            continue
        # Preserve original shape stack but remove placeholder text/frame visibility.
        _set_shape_text(shape, "")
        _remove_shape_box_treatment(shape)


def _largest_autoshape(slide: Any) -> Any:
    autoshapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    if not autoshapes:
        return None
    return max(autoshapes, key=lambda s: int(s.width) * int(s.height))


def _duplicate_slide(prs: Presentation, source_slide: Any) -> Any:
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)

    for shape in source_slide.shapes:
        el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(el, "p:extLst")

    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        try:
            new_slide.part.rels.add_relationship(rel.reltype, rel._target, rel.rId)
        except Exception:
            continue

    return new_slide


def _slide_index(prs: Presentation, slide: Any) -> int:
    for idx, s in enumerate(prs.slides):
        if s.slide_id == slide.slide_id:
            return idx
    return -1


def _slide_id_element(prs: Presentation, slide: Any) -> Any | None:
    sld_id_lst = prs.slides._sldIdLst  # pylint: disable=protected-access
    for el in sld_id_lst:
        try:
            if int(el.id) == int(slide.slide_id):
                return el
        except Exception:
            continue
    return None


def _move_slide_before(prs: Presentation, slide_to_move: Any, anchor_slide: Any) -> None:
    move_el = _slide_id_element(prs, slide_to_move)
    anchor_el = _slide_id_element(prs, anchor_slide)
    if move_el is None or anchor_el is None:
        return
    sld_id_lst = prs.slides._sldIdLst  # pylint: disable=protected-access
    try:
        sld_id_lst.remove(move_el)
        anchor_pos = list(sld_id_lst).index(anchor_el)
        sld_id_lst.insert(anchor_pos, move_el)
    except Exception:
        return


def _move_slide_after(prs: Presentation, slide_to_move: Any, anchor_slide: Any) -> None:
    move_el = _slide_id_element(prs, slide_to_move)
    anchor_el = _slide_id_element(prs, anchor_slide)
    if move_el is None or anchor_el is None:
        return
    sld_id_lst = prs.slides._sldIdLst  # pylint: disable=protected-access
    try:
        sld_id_lst.remove(move_el)
        anchor_pos = list(sld_id_lst).index(anchor_el)
        sld_id_lst.insert(anchor_pos + 1, move_el)
    except Exception:
        return


def _find_first_shared_anchor_slide(prs: Presentation) -> Any | None:
    anchors = []
    for token in ("Competitor Graphics", "Retail Media Optimizations", "Competitor Ad Graphics", "Let’s connect", "Let's connect"):
        slide = _find_slide_by_text(prs, token)
        if slide is not None:
            idx = _slide_index(prs, slide)
            if idx >= 0:
                anchors.append((idx, slide))
    if not anchors:
        return None
    anchors.sort(key=lambda x: x[0])
    return anchors[0][1]


def _populate_pdp_slide(slide: Any, pair_payload: dict[str, Any]) -> None:
    pdp = pair_payload.get("pdp_slide", {}) or {}
    content = pair_payload.get("content_optimization_slide", {}) or {}

    title = _safe_text(pdp.get("product_title"))
    item_id = _safe_text(pdp.get("item_id"))
    image_recs = list(content.get("image_recommendations", []) or [])
    selected_img = (pdp.get("selected_primary_image", {}) or {}).get("url", "")

    title_shape = _find_pdp_title_shape(slide)
    if title_shape:
        _set_pdp_title_and_item_id(title_shape, title, item_id)

    rec_shape = _find_shape_contains(slide, "Image Recommendations")
    if rec_shape:
        _set_pdp_image_recommendations(rec_shape, "Image Recommendations:", image_recs)

    image_box = _largest_autoshape(slide)
    if image_box and _safe_text(selected_img):
        _remove_shape_box_treatment(image_box)
        _insert_image_fit_within_shape(slide, image_box, selected_img, inset_ratio=0.03)
    _suppress_pdp_image_placeholders(slide)


def _populate_content_slide(slide: Any, pair_payload: dict[str, Any]) -> None:
    content = pair_payload.get("content_optimization_slide", {}) or {}
    title = _safe_text(content.get("product_title"))

    title_shape = _find_shape_contains(slide, "Title Recommendations")
    if title_shape:
        rec_title = _safe_text(content.get("recommended_title"))
        value = _title_case_text(rec_title or title)
        _set_title_recommendation_block(
            title_shape,
            "Title Recommendations:",
            value,
        )

    desc_shape = _find_shape_contains(slide, "Description Recommendations")
    if desc_shape:
        _set_heading_and_real_bullets(
            desc_shape,
            "Description Recommendations:",
            list(content.get("description_recommendations", []) or []),
        )

    feat_shape = _find_shape_contains(slide, "Key Features Recommendations")
    if feat_shape:
        bullets = list(content.get("key_features_recommendations", []) or [])
        bullets.extend(list(content.get("top_priority_fixes", []) or []))
        _set_heading_and_real_bullets(
            feat_shape,
            "Key Features Recommendations:",
            bullets,
        )


def _sorted_competitor_slots(slide: Any) -> list[Any]:
    slots = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        shape_text = _shape_text(shape).lower()
        if "(image placeholder)" in shape_text:
            continue
        slots.append(shape)
    slots.sort(key=lambda s: (int(s.top), int(s.left)))
    return slots[:10]


def _remove_shape(slide: Any, shape: Any) -> None:
    try:
        slide.shapes._spTree.remove(shape._element)  # pylint: disable=protected-access
    except Exception:
        return


def _competitor_row_layout_for_count(count: int) -> list[int]:
    # Stable layout families (max 10), with odd counts derived from nearest even family
    # while avoiding empty/dead slots.
    if count <= 1:
        return [1]
    if count == 2:
        return [2]
    if count == 3:
        return [1, 2]
    if count == 4:
        return [2, 2]
    if count == 5:
        return [2, 3]
    if count == 6:
        return [3, 3]
    if count == 7:
        return [3, 4]
    if count == 8:
        return [4, 4]
    if count == 9:
        return [4, 5]
    return [5, 5]


def _layout_competitor_slots(slide: Any, slots: list[Any], active_count: int) -> list[Any]:
    if not slots or active_count <= 0:
        for shape in slots:
            _remove_shape(slide, shape)
        return []

    active_count = min(active_count, len(slots))
    row_layout = _competitor_row_layout_for_count(active_count)
    rows = len(row_layout)
    cols = max(row_layout) if row_layout else 1

    left = min(int(s.left) for s in slots)
    top = min(int(s.top) for s in slots)
    right = max(int(s.left + s.width) for s in slots)
    bottom = max(int(s.top + s.height) for s in slots)
    total_w = max(1, right - left)
    total_h = max(1, bottom - top)

    gap_x = int(total_w * 0.02)
    gap_y = int(total_h * 0.03)
    if cols == 1:
        gap_x = 0
    if rows == 1:
        gap_y = 0

    cell_w = max(1, int((total_w - (cols - 1) * gap_x) / cols))
    cell_h = max(1, int((total_h - (rows - 1) * gap_y) / rows))

    active_slots = slots[:active_count]
    idx = 0
    for row_idx, row_cols in enumerate(row_layout):
        if idx >= len(active_slots):
            break
        row_w = row_cols * cell_w + max(0, row_cols - 1) * gap_x
        row_left = left + int((total_w - row_w) / 2)
        row_top = top + row_idx * (cell_h + gap_y)
        for col_idx in range(row_cols):
            if idx >= len(active_slots):
                break
            shape = active_slots[idx]
            shape.left = row_left + col_idx * (cell_w + gap_x)
            shape.top = row_top
            shape.width = cell_w
            shape.height = cell_h
            idx += 1

    for shape in slots[active_count:]:
        _remove_shape(slide, shape)
    return active_slots


def _populate_competitor_graphics(slide: Any, ordered_assignments: list[dict[str, Any]], notes: str) -> None:
    slots = _sorted_competitor_slots(slide)
    valid_assignments = [
        a
        for a in (ordered_assignments or [])
        if 1 <= int(a.get("display_order", 0) or 0) <= 10 and _safe_text(a.get("url"))
    ]
    valid_assignments.sort(key=lambda a: int(a.get("display_order", 0)))
    active_assignments = valid_assignments[:10]
    active_slots = _layout_competitor_slots(slide, slots, len(active_assignments))

    for idx, assignment in enumerate(active_assignments):
        if idx >= len(active_slots):
            break
        _insert_image_over_shape(slide, active_slots[idx], _safe_text(assignment.get("url")))

    notes_shape = _find_shape_contains(slide, "(image placeholder)")
    if notes_shape:
        _set_shape_text(notes_shape, _safe_text(notes) or "")


def _populate_competitor_graphics_slides(
    prs: Presentation,
    competitor_template_slide: Any,
    competitor_payload: dict[str, Any],
    notes: str,
) -> None:
    if competitor_template_slide is None:
        return

    slide_specs = list(competitor_payload.get("slides", []) or [])
    if not slide_specs:
        slide_specs = [
            {
                "ordered_assignments": list(competitor_payload.get("ordered_assignments", []) or []),
            }
        ]

    if not slide_specs:
        slide_specs = [{"ordered_assignments": []}]

    duplicate_slides: list[Any] = []
    previous_slide = competitor_template_slide
    for _ in slide_specs[1:]:
        dup = _duplicate_slide(prs, competitor_template_slide)
        duplicate_slides.append(dup)
        _move_slide_after(prs, dup, previous_slide)
        previous_slide = dup

    target_slides = [competitor_template_slide, *duplicate_slides]
    for idx, spec in enumerate(slide_specs):
        if idx >= len(target_slides):
            break
        _populate_competitor_graphics(
            target_slides[idx],
            list(spec.get("ordered_assignments", []) or []),
            notes,
        )


def _populate_shared_note_slide(slide: Any, heading: str, body_text: str) -> None:
    body = _safe_text(body_text)
    target = None
    for shape in slide.shapes:
        text = _shape_text(shape)
        if text and heading.lower() not in text.lower() and "(retailer)" not in text.lower():
            target = shape
            break

    if target and getattr(target, "has_text_frame", False):
        _set_shape_text(target, body)
        return

    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.8), Inches(5.5))
    _set_shape_text(textbox, body)


def _populate_retailer_tokens(prs: Presentation, retailer: str) -> None:
    value = _safe_text(retailer)
    if not value:
        return
    for slide in prs.slides:
        for shape in slide.shapes:
            text = _shape_text(shape)
            if "(Retailer)" in text:
                _set_shape_text(shape, text.replace("(Retailer)", value))


def _format_cover_date(audit_date: str) -> str:
    raw = _safe_text(audit_date)
    if not raw:
        return ""
    parsed: dt_date | None = None
    try:
        parsed = dt_date.fromisoformat(raw)
    except Exception:
        for fmt in ("%m/%d/%Y", "%Y/%m/%d", "%B %d, %Y", "%b %d, %Y"):
            try:
                parsed = datetime.strptime(raw, fmt).date()
                break
            except Exception:
                continue
    if parsed is None:
        return raw
    return f"{parsed.strftime('%B')} {parsed.day}, {parsed.year}"


def _populate_cover_title(prs: Presentation, client_name: str) -> None:
    client = _safe_text(client_name)
    if not client or len(prs.slides) < 1:
        return
    cover = prs.slides[0]
    target_text = f"{client} Audit"
    for shape in cover.shapes:
        text = _shape_text(shape)
        if not text:
            continue
        text_l = text.lower()
        if "hyper audit" in text_l or text_l.endswith(" audit") or text_l == "audit":
            _replace_shape_text_preserve_style(shape, target_text)
            break


def _populate_cover_date(prs: Presentation, audit_date: str) -> None:
    if len(prs.slides) < 1:
        return
    formatted = _format_cover_date(audit_date)
    if not formatted:
        return
    cover = prs.slides[0]
    for shape in cover.shapes:
        text = _shape_text(shape)
        if text and any(ch.isdigit() for ch in text) and len(text) <= 40:
            _replace_shape_text_preserve_style(shape, formatted)
            break


def generate_audit_powerpoint_from_template(*, export_plan: dict[str, Any], template_path: str) -> bytes:
    prs = Presentation(template_path)

    pair_payloads = list(export_plan.get("product_slide_pairs", []) or [])
    if not pair_payloads:
        raise ValueError("No included primary product entries were found for export.")

    pdp_template = _find_slide_by_text(prs, "Image Recommendations")
    content_template = _find_slide_by_text(prs, "Content Optimizations")
    competitor_slide = _find_slide_by_text(prs, "Competitor Graphics")
    retail_media_slide = _find_slide_by_text(prs, "Retail Media Optimizations")
    competitor_ad_slide = _find_slide_by_text(prs, "Competitor Ad Graphics")

    if pdp_template is None or content_template is None:
        raise ValueError("Could not find required primary product template slides.")

    _populate_pdp_slide(pdp_template, pair_payloads[0])
    _populate_content_slide(content_template, pair_payloads[0])

    for pair in pair_payloads[1:]:
        pdp_slide = _duplicate_slide(prs, pdp_template)
        content_slide = _duplicate_slide(prs, content_template)
        _populate_pdp_slide(pdp_slide, pair)
        _populate_content_slide(content_slide, pair)
        shared_anchor = _find_first_shared_anchor_slide(prs)
        if shared_anchor is not None:
            _move_slide_before(prs, pdp_slide, shared_anchor)
            _move_slide_before(prs, content_slide, shared_anchor)

    competitor_payload = export_plan.get("competitor_graphics_payload", {}) or {}
    shared_sections = export_plan.get("shared_sections_payload", {}) or {}
    if competitor_slide is not None:
        _populate_competitor_graphics_slides(
            prs,
            competitor_slide,
            competitor_payload,
            _safe_text(shared_sections.get("competitor_graphics_notes", "")),
        )

    if retail_media_slide is not None:
        _populate_shared_note_slide(
            retail_media_slide,
            "Retail Media Optimizations",
            _safe_text(shared_sections.get("retail_media_optimizations", "")),
        )

    if competitor_ad_slide is not None:
        _populate_shared_note_slide(
            competitor_ad_slide,
            "Competitor Ad Graphics",
            _safe_text(shared_sections.get("competitor_ad_graphics_notes", "")),
        )

    metadata = export_plan.get("audit_metadata", {}) or {}
    _populate_cover_title(prs, _safe_text(metadata.get("client_name", "")))
    _populate_retailer_tokens(prs, _safe_text(metadata.get("retailer", "")))
    _populate_cover_date(prs, _safe_text(metadata.get("audit_date", "")))

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue()
